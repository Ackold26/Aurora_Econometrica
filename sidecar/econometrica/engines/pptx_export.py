"""
PPTX report generator for MMM pipeline results.
Uses python-pptx to create branded presentation with charts and speaker notes.
"""
import json
import logging
from datetime import datetime
from pathlib import Path
from typing import Any

logger = logging.getLogger('econometrica')

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.dml.color import RGBColor
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    logger.warning("python-pptx not installed — PPTX export disabled")


# Aurora AI brand colors
AURORA_BLUE   = RGBColor(0x3B, 0x82, 0xF6) if HAS_PPTX else None
AURORA_GREEN  = RGBColor(0x22, 0xC5, 0x5E) if HAS_PPTX else None
AURORA_RED    = RGBColor(0xEF, 0x44, 0x44) if HAS_PPTX else None
AURORA_AMBER  = RGBColor(0xF5, 0x9E, 0x0B) if HAS_PPTX else None
AURORA_DARK   = RGBColor(0x1E, 0x21, 0x2C) if HAS_PPTX else None
AURORA_TEXT    = RGBColor(0xE2, 0xE8, 0xF0) if HAS_PPTX else None
AURORA_MUTED  = RGBColor(0x94, 0xA3, 0xB8) if HAS_PPTX else None
WHITE         = RGBColor(0xFF, 0xFF, 0xFF) if HAS_PPTX else None
BLACK         = RGBColor(0x00, 0x00, 0x00) if HAS_PPTX else None

CHANNEL_COLORS = [
    RGBColor(0x3B, 0x82, 0xF6),  # blue
    RGBColor(0x22, 0xC5, 0x5E),  # green
    RGBColor(0xF5, 0x9E, 0x0B),  # amber
    RGBColor(0xEF, 0x44, 0x44),  # red
    RGBColor(0x8B, 0x5C, 0xF6),  # violet
    RGBColor(0x06, 0xB6, 0xD4),  # cyan
    RGBColor(0xF9, 0x73, 0x16),  # orange
    RGBColor(0x84, 0xCC, 0x16),  # lime
] if HAS_PPTX else []


def _set_slide_bg(slide, color=None):
    """Set slide background to dark."""
    if not color:
        color = AURORA_DARK
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_title_text(slide, text, left=0.5, top=0.3, width=9, height=0.6, size=28, color=None, bold=True):
    """Add a title textbox."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color or WHITE
    p.font.bold = bold
    return txBox


def _add_body_text(slide, text, left=0.5, top=1.2, width=9, height=5, size=14, color=None):
    """Add body text."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color or AURORA_TEXT
        p.space_after = Pt(6)
    return txBox


def _add_notes(slide, text):
    """Add speaker notes."""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


def _style_chart_text(chart, color=None):
    """Делает весь текст графика (оси, легенда, подписи) читаемым на тёмном фоне.
    python-pptx по умолчанию ставит чёрный шрифт — невидим на AURORA_DARK.
    """
    if not HAS_PPTX:
        return
    text_color = color or AURORA_TEXT
    # Category axis (X)
    try:
        ax = chart.category_axis
        ax.tick_labels.font.color.rgb = text_color
        ax.tick_labels.font.size = Pt(10)
    except Exception:
        pass
    # Value axis (Y)
    try:
        ay = chart.value_axis
        ay.tick_labels.font.color.rgb = text_color
        ay.tick_labels.font.size = Pt(10)
    except Exception:
        pass
    # Legend
    try:
        if chart.has_legend:
            chart.legend.font.color.rgb = text_color
            chart.legend.font.size = Pt(10)
    except Exception:
        pass
    # Data labels on each series
    try:
        for series in chart.series:
            try:
                dl = series.data_labels
                dl.font.color.rgb = text_color
                dl.font.size = Pt(10)
            except Exception:
                continue
    except Exception:
        pass


def build_pptx(model_data: dict, decompose_data: dict, optimize_data: dict, output_path: str) -> dict[str, Any]:
    """Build a branded PPTX presentation from MMM pipeline data."""
    if not HAS_PPTX:
        return {'status': 'error', 'message': 'python-pptx не установлен. pip install python-pptx'}

    # Defensive: backend sometimes hands us None for a section
    model_data = model_data or {}
    decompose_data = decompose_data or {}
    optimize_data = optimize_data or {}

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9

    diag = model_data.get('diagnostics', {}) or {}
    mqs = diag.get('mqs', {}) or {}
    channels = decompose_data.get('channels', []) or []
    opt_channels = optimize_data.get('channels', []) or []
    waterfall = decompose_data.get('waterfall', []) or []

    logger.info(
        f"build_pptx start: channels={len(channels)} opt_channels={len(opt_channels)} "
        f"waterfall={len(waterfall)} mqs={mqs.get('score')} r2={diag.get('r_squared')}"
    )

    failed_phases: list[str] = []

    blank_layout = prs.slide_layouts[6]  # blank

    # ── Slide 1: Title ─────────────────────────────────────
    try:
        slide = prs.slides.add_slide(blank_layout)
        _set_slide_bg(slide)
        _add_title_text(slide, "Marketing Mix Model", top=1.5, size=36)
        _add_body_text(slide, f"Аналитический отчёт\n{datetime.now().strftime('%d.%m.%Y')}", top=2.5, size=18, color=AURORA_MUTED)
        _add_body_text(slide, "Aurora AI Econometrica", top=4.2, size=12, color=AURORA_BLUE)
        logger.info("PPTX phase OK: title")
    except Exception:
        logger.exception("PPTX phase FAILED: title")
        failed_phases.append('title')

    # ── Slide 2: Executive Summary ────────────────────────
    mqs_score = mqs.get('score', 0) or 0
    mqs_label = mqs.get('tier_label', 'N/A') or 'N/A'
    # Backend nests fit metrics under `diagnostics.metrics` (mape under `mape_pct`).
    metrics = diag.get('metrics', {}) or {}
    r_sq = metrics.get('r_squared', diag.get('r_squared', 0)) or 0
    mape_val = metrics.get('mape_pct', diag.get('mape', 0)) or 0
    lift = optimize_data.get('expected_lift_pct', 0) or 0
    budget = optimize_data.get('total_budget', 0) or 0

    try:
        slide = prs.slides.add_slide(blank_layout)
        _set_slide_bg(slide)
        _add_title_text(slide, "Executive Summary", size=24)

        # Крупный MQS-бейдж слева (как в UI)
        mqs_color = AURORA_GREEN if mqs_score >= 80 else (AURORA_BLUE if mqs_score >= 60 else AURORA_AMBER)
        mqs_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(2.4), Inches(2.0))
        mqs_tf = mqs_box.text_frame
        mqs_tf.word_wrap = True
        mqs_label_p = mqs_tf.paragraphs[0]
        mqs_label_p.text = "MQS"
        mqs_label_p.font.size = Pt(11)
        mqs_label_p.font.color.rgb = AURORA_MUTED
        mqs_label_p.font.bold = True
        p1 = mqs_tf.add_paragraph()
        p1.text = f"{mqs_score:.0f}"
        p1.font.size = Pt(64)
        p1.font.color.rgb = mqs_color
        p1.font.bold = True
        p2 = mqs_tf.add_paragraph()
        p2.text = mqs_label.upper()
        p2.font.size = Pt(12)
        p2.font.color.rgb = AURORA_MUTED
        p2.font.bold = True

        # Метрики справа — сетка 2x2
        metric_defs = [
            ("R\u00B2", f"{r_sq:.3f}", f"{r_sq*100:.0f}% объяснённой вариации", AURORA_GREEN if r_sq >= 0.7 else AURORA_AMBER),
            ("MAPE", f"{mape_val:.1f}%", "средняя ошибка прогноза", AURORA_GREEN if mape_val < 10 else AURORA_AMBER),
            ("Прирост", f"{lift:+.1f}%", "при перераспределении", AURORA_GREEN if lift > 5 else AURORA_TEXT),
            ("Бюджет", f"{budget:,.0f}".replace(",", " ") + " ₽", "общий", AURORA_TEXT),
        ]
        for idx, (lbl, val, sub, col) in enumerate(metric_defs):
            col_idx = idx % 2
            row_idx = idx // 2
            x = 3.2 + col_idx * 3.3
            y = 1.25 + row_idx * 1.05
            box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(3.2), Inches(0.95))
            tf = box.text_frame
            tf.word_wrap = True
            lp = tf.paragraphs[0]
            lp.text = lbl
            lp.font.size = Pt(10)
            lp.font.color.rgb = AURORA_MUTED
            lp.font.bold = True
            vp = tf.add_paragraph()
            vp.text = val
            vp.font.size = Pt(24)
            vp.font.color.rgb = col
            vp.font.bold = True
            sp = tf.add_paragraph()
            sp.text = sub
            sp.font.size = Pt(9)
            sp.font.color.rgb = AURORA_MUTED

        # Нижняя строка — контекст
        _add_body_text(
            slide,
            f"Bayesian MMM с {len(channels)} канал{'ами' if len(channels) > 4 else 'ами' if len(channels) > 1 else 'ом'} медиа. Adstock + Hill saturation. MCMC-сэмплер, {len(channels)} параметр{'ов' if len(channels) > 4 else 'а' if len(channels) > 1 else ''}.",
            top=3.7, size=12, color=AURORA_MUTED,
        )

        _add_notes(slide, "MQS > 80 = отлично, 60-80 = хорошо, < 60 = требует доработки. R\u00b2 показывает долю объяснённой вариации.")
        logger.info("PPTX phase OK: summary")
    except Exception:
        logger.exception("PPTX phase FAILED: summary")
        failed_phases.append('summary')

    # ── Slide 2.5: Спецификация модели ────────────────────
    try:
        spec = diag.get('model_spec') or {}
        if not spec:
            from utils.model_spec import bayesian_mmm_spec
            spec = bayesian_mmm_spec()

        slide = prs.slides.add_slide(blank_layout)
        _set_slide_bg(slide)
        _add_title_text(slide, spec.get('title', 'Спецификация модели'), size=24)

        # Subtitle + formula on a single text block
        spec_lines = [
            spec.get('subtitle', ''),
            '',
            'Формула:',
            f"  {spec.get('formula', '')}",
            '',
            'Трансформации:',
        ]
        for tr in spec.get('transformations', []):
            spec_lines.append(f"  • {tr.get('name')}: {tr.get('formula')}")
        spec_lines.append('')
        spec_lines.append('Priors (априорные распределения):')
        for p in spec.get('priors', []):
            spec_lines.append(f"  {p.get('symbol'):>4}  {p.get('distribution'):<22}  — {p.get('name')}")

        _add_body_text(slide, '\n'.join(spec_lines), top=0.95, size=11)

        inf = spec.get('inference', {}) or {}
        notes_parts = [
            spec.get('subtitle', ''),
            f"Инференс: {inf.get('method', 'N/A')}",
            inf.get('note', ''),
            spec.get('normalization', ''),
        ]
        _add_notes(slide, '\n'.join(p for p in notes_parts if p))
        logger.info("PPTX phase OK: spec")
    except Exception:
        logger.exception("PPTX phase FAILED: spec")
        failed_phases.append('spec')

    # ── Slide 3: Декомпозиция ─────────────────────────────
    if waterfall:
        try:
            slide = prs.slides.add_slide(blank_layout)
            _set_slide_bg(slide)
            _add_title_text(slide, "Декомпозиция продаж", size=24)

            from pptx.chart.data import CategoryChartData
            chart_data_obj = CategoryChartData()
            # waterfall формат из backend: {'labels': [...], 'values': [...], 'types': [...]}
            # Legacy-формат (list of {category, value}) — fallback для старых pickle-кэшей.
            if isinstance(waterfall, dict):
                labels = waterfall.get('labels', [])
                values = waterfall.get('values', [])
            else:
                labels = [str(w.get('category', '')) for w in waterfall]
                values = [float(w.get('value', 0) or 0) for w in waterfall]
            chart_data_obj.categories = [str(x) for x in labels]
            chart_data_obj.add_series('Вклад', [float(v or 0) for v in values])

            chart_frame = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(1.2), Inches(9), Inches(3.8), chart_data_obj
            )
            chart = chart_frame.chart
            chart.has_legend = False
            series = chart.series[0]
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = AURORA_BLUE
            _style_chart_text(chart)

            base_pct = decompose_data.get('baseline_pct', decompose_data.get('base_pct', 0)) or 0
            _add_notes(slide, f"Base sales = {base_pct:.0f}%. Это органические продажи без рекламного воздействия. Значение > 60% означает сильный бренд.")
            logger.info("PPTX phase OK: decomposition")
        except Exception:
            logger.exception("PPTX phase FAILED: decomposition")
            failed_phases.append('decomposition')

    # ── Slide 4: ROI по каналам ───────────────────────────
    if channels:
        try:
            slide = prs.slides.add_slide(blank_layout)
            _set_slide_bg(slide)
            _add_title_text(slide, "ROI по каналам", size=24)

            from pptx.chart.data import CategoryChartData
            chart_data_obj = CategoryChartData()
            sorted_chs = sorted(channels, key=lambda c: c.get('roi', 0) or 0, reverse=True)
            chart_data_obj.categories = [str(c.get('name', '')) for c in sorted_chs]
            chart_data_obj.add_series('ROI', [float(c.get('roi', 0) or 0) for c in sorted_chs])

            chart_frame = slide.shapes.add_chart(
                XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.5), Inches(1.2), Inches(9), Inches(3.8), chart_data_obj
            )
            chart = chart_frame.chart
            chart.has_legend = False
            series = chart.series[0]
            for idx in range(len(sorted_chs)):
                pt = series.points[idx]
                pt.format.fill.solid()
                color_idx = idx % len(CHANNEL_COLORS)
                pt.format.fill.fore_color.rgb = CHANNEL_COLORS[color_idx]
            _style_chart_text(chart)

            _add_notes(slide, "ROI > 2.0x считается хорошим. ROI < 1.0x означает, что расходы на канал превышают его вклад в продажи.")
            logger.info("PPTX phase OK: roi")
        except Exception:
            logger.exception("PPTX phase FAILED: roi")
            failed_phases.append('roi')

    # ── Slide 5: Share of Spend vs Effect ─────────────────
    if channels:
        try:
            slide = prs.slides.add_slide(blank_layout)
            _set_slide_bg(slide)
            _add_title_text(slide, "Share of Spend vs Share of Effect", size=22)

            from pptx.chart.data import CategoryChartData
            chart_data_obj = CategoryChartData()
            total_spend = sum(float(c.get('spend', 0) or 0) for c in channels)
            total_contrib = sum(float(c.get('contribution', 0) or 0) for c in channels)

            chart_data_obj.categories = [str(c.get('name', '')) for c in channels]
            chart_data_obj.add_series('% бюджета', [(float(c.get('spend', 0) or 0) / total_spend * 100) if total_spend else 0 for c in channels])
            chart_data_obj.add_series('% эффекта', [(float(c.get('contribution', 0) or 0) / total_contrib * 100) if total_contrib else 0 for c in channels])

            chart_frame = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(1.2), Inches(9), Inches(3.8), chart_data_obj
            )
            chart = chart_frame.chart
            chart.series[0].format.fill.solid()
            chart.series[0].format.fill.fore_color.rgb = AURORA_MUTED
            chart.series[1].format.fill.solid()
            chart.series[1].format.fill.fore_color.rgb = AURORA_GREEN
            chart.has_legend = True
            _style_chart_text(chart)

            _add_notes(slide, "Если % эффекта > % бюджета — канал недоинвестирован (efficiency > 1). Если наоборот — перенасыщен.")
            logger.info("PPTX phase OK: share")
        except Exception:
            logger.exception("PPTX phase FAILED: share")
            failed_phases.append('share')

    # ── Slide 6: Оптимизация ──────────────────────────────
    if opt_channels:
        try:
            slide = prs.slides.add_slide(blank_layout)
            _set_slide_bg(slide)
            _add_title_text(slide, f"Оптимизация бюджета (lift: {lift:+.1f}%)", size=22)

            from pptx.chart.data import CategoryChartData
            chart_data_obj = CategoryChartData()
            chart_data_obj.categories = [str(c.get('name', '')) for c in opt_channels]
            chart_data_obj.add_series('Текущий', [float(c.get('current_spend', 0) or 0) for c in opt_channels])
            chart_data_obj.add_series('Оптимальный', [float(c.get('optimal_spend', 0) or 0) for c in opt_channels])

            chart_frame = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(1.2), Inches(9), Inches(3.8), chart_data_obj
            )
            chart = chart_frame.chart
            chart.series[0].format.fill.solid()
            chart.series[0].format.fill.fore_color.rgb = AURORA_MUTED
            chart.series[1].format.fill.solid()
            chart.series[1].format.fill.fore_color.rgb = AURORA_BLUE
            chart.has_legend = True
            _style_chart_text(chart)

            _add_notes(slide, f"Ожидаемый прирост {lift:+.1f}% при перераспределении бюджета. Рекомендуется пилотный период 4-6 недель.")
            logger.info("PPTX phase OK: optimize")
        except Exception:
            logger.exception("PPTX phase FAILED: optimize")
            failed_phases.append('optimize')

    # ── Slide 7: Рекомендации ─────────────────────────────
    try:
        slide = prs.slides.add_slide(blank_layout)
        _set_slide_bg(slide)
        _add_title_text(slide, "Рекомендации", size=24)

        # (severity, text, color)
        recs: list[tuple[str, str, Any]] = []
        if lift > 5:
            recs.append(("ВЫСОКАЯ", f"Перераспределить бюджет — ожидаемый прирост {lift:+.1f}%", AURORA_GREEN))
        if channels:
            top_ch = max(channels, key=lambda c: c.get('roi', 0) or 0)
            recs.append(("ВЫСОКАЯ", f"Приоритет: {top_ch.get('name', '')} (ROI {top_ch.get('roi', 0) or 0:.1f}×)", AURORA_GREEN))
            bot_ch = min(channels, key=lambda c: c.get('roi', 0) or 0)
            if (bot_ch.get('roi', 0) or 0) < 1:
                recs.append(("СРЕДНЯЯ", f"Сократить {bot_ch.get('name', '')} (ROI {bot_ch.get('roi', 0) or 0:.1f}×)", AURORA_AMBER))
        if r_sq < 0.7:
            recs.append(("СРЕДНЯЯ", "R\u00b2 ниже 0.7 — добавить контрольные переменные (сезонность, промо)", AURORA_AMBER))
        if mqs_score >= 80:
            recs.append(("ВЫСОКАЯ", "Высокий MQS — результаты надёжны для принятия решений", AURORA_GREEN))
        metrics_ratio = diag.get('metrics', {}).get('ratio') if isinstance(diag.get('metrics'), dict) else None
        if metrics_ratio is not None and metrics_ratio < 4:
            recs.append(("ВАЖНАЯ", f"Данных мало (Ratio {metrics_ratio:.1f}:1) — пилот 4-6 недель перед переходом", AURORA_RED))

        if not recs:
            _add_body_text(slide, "Нет специальных рекомендаций.", top=1.2, size=14)
        else:
            for i, (sev, text, color) in enumerate(recs):
                y = 1.2 + i * 0.6
                # Severity badge
                badge = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(1.4), Inches(0.45))
                btf = badge.text_frame
                bp = btf.paragraphs[0]
                bp.text = sev
                bp.font.size = Pt(10)
                bp.font.color.rgb = color
                bp.font.bold = True
                # Recommendation text
                body = slide.shapes.add_textbox(Inches(2.0), Inches(y), Inches(7.6), Inches(0.55))
                btf2 = body.text_frame
                btf2.word_wrap = True
                bp2 = btf2.paragraphs[0]
                bp2.text = text
                bp2.font.size = Pt(14)
                bp2.font.color.rgb = AURORA_TEXT
        logger.info("PPTX phase OK: recommendations")
    except Exception:
        logger.exception("PPTX phase FAILED: recommendations")
        failed_phases.append('recommendations')

    # ── Slide 8: Методология ──────────────────────────────
    try:
        slide = prs.slides.add_slide(blank_layout)
        _set_slide_bg(slide)
        _add_title_text(slide, "Методология", size=24)

        sections = [
            ("Тип модели", "Байесовская Marketing Mix Model (Bayesian MMM) с priors — устойчива к мультиколлинеарности."),
            ("Adstock", "Геометрический/Weibull per-канал — моделирует отложенный эффект рекламы (carry-over)."),
            ("Saturation", "Hill function с параметрами α (крутизна) и γ (точка полу-насыщения) — убывающая отдача."),
            ("Инференс", "MCMC (Markov Chain Monte Carlo), NumPyro NUTS-сэмплер, 4 цепи × 2000 итераций."),
            ("Оптимизация", "scipy SLSQP с бюджетными ограничениями (глобальные Мин/Макс % и per-channel)."),
            ("Данные", f"{len(channels)} медиа-канал{'ов' if len(channels) > 4 else 'а' if len(channels) > 1 else ''}, {len(decompose_data.get('time_series', {}).get('dates', []) or [])} наблюдений."),
        ]
        for i, (label, desc) in enumerate(sections):
            y = 1.1 + i * 0.65
            # Label (left, accent)
            lbl_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(2.3), Inches(0.5))
            lbl_tf = lbl_box.text_frame
            lp = lbl_tf.paragraphs[0]
            lp.text = label
            lp.font.size = Pt(12)
            lp.font.color.rgb = AURORA_BLUE
            lp.font.bold = True
            # Description
            desc_box = slide.shapes.add_textbox(Inches(2.9), Inches(y), Inches(6.8), Inches(0.6))
            desc_tf = desc_box.text_frame
            desc_tf.word_wrap = True
            dp = desc_tf.paragraphs[0]
            dp.text = desc
            dp.font.size = Pt(11)
            dp.font.color.rgb = AURORA_TEXT

        _add_notes(slide, "Байесовский подход позволяет учитывать априорные знания и оценивать неопределённость. Hill function моделирует убывающую отдачу.")
        logger.info("PPTX phase OK: methodology")
    except Exception:
        logger.exception("PPTX phase FAILED: methodology")
        failed_phases.append('methodology')

    # ── Save ──────────────────────────────────────────────
    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    logger.info(f"PPTX saved: {out} ({len(prs.slides)} slides, failed_phases={failed_phases})")

    return {
        'status': 'ok' if not failed_phases else 'partial',
        'path': str(out),
        'slides': len(prs.slides),
        'failed_phases': failed_phases,
    }
