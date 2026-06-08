"""Тест механики встраивания шрифтов в .pptx (OOXML) + OFL-лицензионный guard.

font_embed портирован из Smart Analytica. В Econometrica НЕ вплетён в save()
(PPTX использует MS-проприетарные Georgia/Arial — встраивать нельзя). Тест
проверяет: (1) OOXML-механика встраивания работает; (2) OFL-guard отвергает
MS-проприетарные шрифты.
"""
import os
import zipfile

import pytest

from pptx import Presentation
from pptx.util import Inches, Pt

from aurora_pptx import font_embed
from aurora_pptx import text_metrics as TM


def _tiny_deck(path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = "Шрифт-тест"
    r.font.size = Pt(18)
    prs.save(path)
    return path


def test_embed_fonts_injects_ooxml_parts(tmp_path):
    """Механика OOXML: встроенный шрифт даёт fntdata-часть + embeddedFontLst + флаг.

    Носитель — системный Arial (тест-артефакт, не распространяется). Проверяем РОВНО
    байты пакета, не лицензию (лицензию проверяет OFL-guard, см. ниже)."""
    arial = TM.font_path("Arial", bold=False)
    if not arial:
        pytest.skip("нет системного Arial TTF (не-Windows рантайм) — механику не проверить")
    deck = _tiny_deck(str(tmp_path / "deck.pptx"))
    ok = font_embed.embed_fonts(deck, [{"typeface": "TestFace", "regular": arial}])
    assert ok is True
    with zipfile.ZipFile(deck) as z:
        names = z.namelist()
        assert any(n.startswith("ppt/fonts/font") and n.endswith(".fntdata") for n in names)
        pres = z.read("ppt/presentation.xml").decode("utf-8")
        assert 'embedTrueTypeFonts="1"' in pres
        assert "embeddedFontLst" in pres
        assert "TestFace" in pres


def test_embed_fonts_noop_when_nothing_resolved(tmp_path):
    deck = _tiny_deck(str(tmp_path / "deck.pptx"))
    # путь не существует → ничего не встроено
    assert font_embed.embed_fonts(deck, [{"typeface": "X", "regular": "/no/such.ttf"}]) is False


def test_ofl_guard_rejects_ms_proprietary(tmp_path):
    """OFL-guard: Arial/Georgia (MS-проприетарные) НЕ встраиваются (skipped)."""
    deck = _tiny_deck(str(tmp_path / "deck.pptx"))
    res = font_embed.embed_brand_fonts(deck, ["Arial", "Georgia"])
    assert res["embedded"] == []
    assert set(res["skipped"]) == {"Arial", "Georgia"}


def test_ofl_guard_strict_raises_on_ms(tmp_path):
    deck = _tiny_deck(str(tmp_path / "deck.pptx"))
    with pytest.raises(ValueError, match="OFL"):
        font_embed.embed_brand_fonts(deck, ["Georgia"], strict=True)


def test_ofl_allowlist_contains_brand_candidates():
    # Inter/Lora — OFL-шрифты HTML-отчёта Econometrica; кандидаты для PPTX при
    # переходе с MS-шрифтов. Montserrat — шрифт Smart Analytica.
    assert {"Inter", "Lora", "Montserrat"} <= font_embed.OFL_EMBEDDABLE
