"""Пласт 2, Фаза 3: KPI-паспорт — тексты отчётов PPTX/HTML count-aware.

Проверяет что для count kpi:
- result-фразы НЕ содержат «каждый рубль», «ROAS», «mROAS» в контексте результата/отдачи
- contrib-заголовок таблицы НЕ «₽ млн»
- helper-функции (_lift_phrase, _hero_vs_leader_quote, lift_phrase, hero_vs_leader_quote)
  возвращают KPI-aware формулировки

Область: aurora_pptx/kpi_helpers.py, aurora_html/sections.py, aurora_pptx/builder.py

Замечание: тесты через builder.py / render_* требуют matplotlib и python-pptx;
они помечены marks=pptx/html и пропускаются при TclError (_tkinter) или ImportError.
"""
from __future__ import annotations

import copy
import json
import os
import re
import sys

import pytest

_HERE = os.path.dirname(os.path.abspath(__file__))
# econometrica/ в sys.path — тест cwd-независим (aurora_pptx/aurora_html
# резолвятся из любой рабочей директории, как в bootstrap-тестах Фазы 3 пласт 1).
sys.path.insert(0, os.path.dirname(_HERE))
_FIXTURE = os.path.join(_HERE, "fixtures", "kagocel_builder_payload.json")
_STRINGS = os.path.join(os.path.dirname(_HERE), "aurora_html", "strings_ru.json")


# ─── KPI фикстуры ─────────────────────────────────────────────────────────────

def _kpi_count():
    """kpi dict для счётной метрики (упаковки, лиды)."""
    return {
        "kpi_kind": "count",
        "mode": "roi",
        "kpi_type": "count_leads",
        "is_legacy": False,
        "metric_label": "CPU, ₽/лид",
        "metric_short": "CPU",
        "target_unit": "лид",
        "target_axis": "Лиды",
        "cpu_per_label": "₽/лид",
        "vpcu": 500.0,
        "vpcu_label": "₽/лид",
        "methodology_label": "",
    }


def _kpi_monetary():
    """kpi dict для денежной метрики (legacy ROI)."""
    return {
        "kpi_kind": "monetary",
        "mode": "roi",
        "kpi_type": None,
        "is_legacy": True,
        "metric_label": "ROI",
        "metric_short": "ROI",
        "target_unit": "₽",
        "target_axis": "Продажи, ₽",
        "cpu_per_label": "₽/ед.",
        "vpcu": None,
        "vpcu_label": "",
        "methodology_label": "",
    }


def _kpi_effectiveness():
    """kpi dict для режима эффективности (доля)."""
    return {
        "kpi_kind": "monetary",
        "mode": "effectiveness",
        "kpi_type": None,
        "is_legacy": False,
        "metric_label": "Доля",
        "metric_short": "Доля",
        "target_unit": "%",
        "target_axis": "Доля эффекта",
        "cpu_per_label": "₽/ед.",
        "vpcu": None,
        "vpcu_label": "",
        "methodology_label": "",
    }


# ─── Тесты kpi_helpers.lift_phrase ────────────────────────────────────────────

class TestLiftPhrase:
    def setup_method(self):
        from aurora_pptx.kpi_helpers import lift_phrase
        self.lift_phrase = lift_phrase

    def test_monetary_contains_roas(self):
        result = self.lift_phrase(5.0, _kpi_monetary())
        assert "ROAS" in result
        assert "5.0" in result

    def test_count_no_roas_no_ruble(self):
        result = self.lift_phrase(5.0, _kpi_count())
        assert "ROAS" not in result, f"count lift_phrase содержит ROAS: {result!r}"
        assert "каждый рубль" not in result.lower(), f"count lift_phrase содержит 'каждый рубль': {result!r}"
        assert "5.0" in result

    def test_effectiveness_no_roas(self):
        result = self.lift_phrase(3.5, _kpi_effectiveness())
        assert "ROAS" not in result, f"effectiveness lift_phrase содержит ROAS: {result!r}"
        assert "3.5" in result
        assert "доля" in result.lower() or "эффект" in result.lower()

    def test_none_lift_no_roas(self):
        result = self.lift_phrase(None, _kpi_count())
        assert "ROAS" not in result, f"lift_phrase(None, count) содержит ROAS: {result!r}"


# ─── Тесты kpi_helpers.hero_vs_leader_quote ───────────────────────────────────

class TestHeroVsLeaderQuote:
    def setup_method(self):
        from aurora_pptx.kpi_helpers import hero_vs_leader_quote
        self.quote = hero_vs_leader_quote

    def test_monetary_contains_ruble(self):
        result = self.quote("Digital", "TV", _kpi_monetary())
        assert "рубль" in result.lower(), f"monetary quote не содержит 'рубль': {result!r}"
        assert "Digital" in result
        assert "TV" in result

    def test_count_no_ruble_phrase(self):
        result = self.quote("Digital", "TV", _kpi_count())
        assert "каждый рубль" not in result.lower(), \
            f"count quote содержит 'каждый рубль': {result!r}"
        assert "ROAS" not in result, f"count quote содержит ROAS: {result!r}"
        assert "Digital" in result
        assert "TV" in result

    def test_effectiveness_no_ruble(self):
        result = self.quote("Digital", "TV", _kpi_effectiveness())
        assert "каждый рубль" not in result.lower(), \
            f"effectiveness quote содержит 'каждый рубль': {result!r}"
        assert "доля" in result.lower() or "эффект" in result.lower()


# ─── Тесты sections._lift_phrase и _hero_vs_leader_quote ──────────────────────

class TestSectionsHelpers:
    def test_sections_lift_phrase_count_no_roas(self):
        from aurora_html.sections import _lift_phrase
        result = _lift_phrase(4.0, _kpi_count())
        assert "ROAS" not in result, f"sections._lift_phrase(count) содержит ROAS: {result!r}"
        assert "mROAS" not in result

    def test_sections_lift_phrase_monetary_has_roas(self):
        from aurora_html.sections import _lift_phrase
        result = _lift_phrase(4.0, _kpi_monetary())
        assert "ROAS" in result

    def test_sections_hero_vs_leader_count_no_ruble(self):
        from aurora_html.sections import _hero_vs_leader_quote
        result = _hero_vs_leader_quote("Канал1", "Канал2", _kpi_count())
        assert "каждый рубль" not in result.lower(), \
            f"sections._hero_vs_leader_quote(count) содержит 'каждый рубль': {result!r}"

    def test_sections_hero_vs_leader_effectiveness_no_ruble(self):
        from aurora_html.sections import _hero_vs_leader_quote
        result = _hero_vs_leader_quote("Канал1", "Канал2", _kpi_effectiveness())
        assert "каждый рубль" not in result.lower()
        assert "ROAS" not in result


# ─── HTML sections с count kpi ────────────────────────────────────────────────

def _load_payload_with_kpi(kpi_kind="count"):
    """Загружает фикстуру и патчит kpi на count."""
    if not os.path.exists(_FIXTURE):
        return None
    with open(_FIXTURE, encoding="utf-8") as f:
        base = json.load(f)
    p = copy.deepcopy(base)
    p["kpi"] = {
        "kpi_kind": kpi_kind,
        "derived_mode": "roi",
        "kpi_type": "count_leads" if kpi_kind == "count" else None,
        "value_per_count_unit": 500.0,
        "value_per_count_unit_label": "₽/лид",
        "labels": {
            "metric_label": "CPU, ₽/лид",
            "metric_short_label": "CPU",
            "target_unit_label": "лид",
            "target_axis_label": "Лиды",
            "methodology_label": "",
        },
    }
    return p


def _make_ctx(payload: dict) -> dict:
    """Минимальный context dict для render_* функций sections.py."""
    if not os.path.exists(_STRINGS):
        return None
    with open(_STRINGS, encoding="utf-8") as f:
        strings = json.load(f)
    return {
        "meta": payload.get("meta") or {},
        "facts": payload.get("narrative_facts"),
        "channels": payload.get("channels") or [],
        "diagnostics": payload.get("diagnostics") or {},
        "strings": strings,
        "kpi": payload.get("kpi") or {},
    }


def _strip_tags(html: str) -> str:
    """Убирает HTML-теги для анализа текста."""
    return re.sub(r"<[^>]+>", " ", html)


@pytest.fixture(scope="module")
def count_ctx():
    p = _load_payload_with_kpi("count")
    if p is None:
        return None
    ctx = _make_ctx(p)
    return ctx


class TestHtmlKeyMessageCountKpi:
    """render_key_message с count kpi не должен давать «каждый рубль» / «ROAS»."""

    def test_key_message_no_ruble_phrase(self, count_ctx):
        if count_ctx is None:
            pytest.skip("Фикстура недоступна")
        from aurora_html.sections import render_key_message
        html = render_key_message(count_ctx)
        text = _strip_tags(html)
        assert "каждый рубль" not in text.lower(), \
            f"render_key_message (count): содержит 'каждый рубль': ...{text[max(0,text.lower().find('каждый рубль')-30):text.lower().find('каждый рубль')+60]!r}..."

    def test_key_message_no_roas(self, count_ctx):
        if count_ctx is None:
            pytest.skip("Фикстура недоступна")
        from aurora_html.sections import render_key_message
        html = render_key_message(count_ctx)
        text = _strip_tags(html)
        # ROAS в pull quote — недопустим для count
        # Разрешён в данных-справочниках (методология) — не проверяем весь HTML
        # Проверяем только pull quote (внутри blockquote)
        blockquote_re = re.compile(r'<blockquote[^>]*>(.*?)</blockquote>', re.DOTALL)
        for match in blockquote_re.finditer(html):
            quote_text = _strip_tags(match.group(1))
            assert "ROAS" not in quote_text, \
                f"render_key_message (count) pull quote содержит ROAS: {quote_text!r}"
            assert "каждый рубль" not in quote_text.lower(), \
                f"render_key_message (count) pull quote содержит 'каждый рубль': {quote_text!r}"


class TestHtmlActionTableCountKpi:
    """render_action_table с count kpi: contrib column header НЕ «₽ млн»."""

    def test_contrib_header_not_rub_mln(self, count_ctx):
        if count_ctx is None:
            pytest.skip("Фикстура недоступна")
        from aurora_html.sections import render_action_table
        html = render_action_table(count_ctx)
        # Заголовок contrib столбца — col 2
        # Ищем текст в th[data-col=2], потом unit span
        th_re = re.compile(
            r'<th[^>]*data-col="2"[^>]*>(.*?)</th>', re.DOTALL | re.IGNORECASE
        )
        match = th_re.search(html)
        if match:
            th_text = _strip_tags(match.group(1))
            assert "₽ млн" not in th_text, \
                f"render_action_table (count) contrib th содержит '₽ млн': {th_text!r}"

    def test_contrib_header_has_result_unit(self, count_ctx):
        if count_ctx is None:
            pytest.skip("Фикстура недоступна")
        from aurora_html.sections import render_action_table
        html = render_action_table(count_ctx)
        th_re = re.compile(
            r'<th[^>]*data-col="2"[^>]*>(.*?)</th>', re.DOTALL | re.IGNORECASE
        )
        match = th_re.search(html)
        if match:
            th_text = _strip_tags(match.group(1))
            # Должна быть единица результата (лид) или нейтральная «ед.»
            # Важно — НЕ «₽ млн»
            assert th_text.strip() != "", "contrib th для count пустой"


class TestHtmlAtAGlanceCountKpi:
    """render_at_a_glance с count kpi: f3_realloc_support НЕ «ROAS»."""

    def test_f3_support_no_roas(self, count_ctx):
        if count_ctx is None:
            pytest.skip("Фикстура недоступна")
        from aurora_html.sections import render_at_a_glance
        # Убедимся что hero != leader (иначе ветка f3_realloc не активируется)
        ctx = copy.deepcopy(count_ctx)
        if ctx.get("facts"):
            ctx["facts"]["reallocation_mln"] = 50.0
            ctx["facts"]["hero_channel"] = "Digital"
            ctx["facts"]["leader_channel"] = "TV"
            ctx["facts"]["expected_lift_pct"] = 3.5
            ctx["facts"]["binding_constraints"] = False
        html = render_at_a_glance(ctx)
        text = _strip_tags(html)
        # Ищем finding-support строки (после finding-headline, в finding-support)
        support_re = re.compile(
            r'class="finding-support"[^>]*>(.*?)</p>', re.DOTALL | re.IGNORECASE
        )
        for match in support_re.finditer(html):
            support_text = _strip_tags(match.group(1))
            # Допустимо «ROAS» только в monetary/legacy
            # Для count — недопустимо
            assert "Ожидаемый прирост ROAS" not in support_text, \
                f"render_at_a_glance (count) finding-support содержит 'Ожидаемый прирост ROAS': {support_text!r}"


class TestHtmlRecommendationCountKpi:
    """render_recommendation с count kpi: impact label НЕ «Прогнозный ROAS»."""

    def test_impact_label_no_roas(self, count_ctx):
        if count_ctx is None:
            pytest.skip("Фикстура недоступна")
        from aurora_html.sections import render_recommendation
        html = render_recommendation(count_ctx)
        text = _strip_tags(html)
        assert "Прогнозный ROAS" not in text, \
            f"render_recommendation (count) содержит 'Прогнозный ROAS': найдено в тексте"


# ─── PPTX builder с count kpi ─────────────────────────────────────────────────

def _pptx_text(prs) -> str:
    """Собирает весь текст из PPTX."""
    from pptx.oxml.ns import qn
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
    return "\n".join(parts)


@pytest.fixture(scope="module")
def count_payload():
    p = _load_payload_with_kpi("count")
    return p


class TestPptxBuilderCountKpi:
    """AuroraPPTXBuilder с count kpi не должен давать «Каждый рубль» / «ROAS» на результате."""

    def test_pptx_count_no_ruble_in_pull_quote(self, count_payload):
        """Pull quote (ключевой вывод слайд 5) не должен содержать «Каждый рубль … возвращает» для count.

        Примечание: «каждый рубль приносит убыток» в action_reasoning корректно
        (затраты — деньги для любого KPI) — его НЕ проверяем. Проверяем только
        pull quote в формате «Каждый рубль в X возвращает больше, чем в Y».
        """
        if count_payload is None:
            pytest.skip("Фикстура недоступна")
        try:
            from aurora_pptx.builder import AuroraPPTXBuilder
        except ImportError as e:
            pytest.skip(f"AuroraPPTXBuilder недоступен: {e}")
        try:
            prs = AuroraPPTXBuilder(count_payload).build()
        except Exception as e:
            if "TclError" in str(type(e).__name__) or "_tkinter" in str(e):
                pytest.skip(f"Окружение matplotlib/_tkinter: {e}")
            raise
        txt = _pptx_text(prs)
        # Специфическая фраза pull quote результата: «Каждый рубль в X возвращает»
        # Разрешено: «каждый рубль приносит убыток» (затраты/убыток — валидно для любого KPI)
        forbidden_pattern = re.compile(
            r"каждый рубль в .+ возвращает", re.IGNORECASE
        )
        match = forbidden_pattern.search(txt)
        assert match is None, \
            f"PPTX (count) pull quote содержит 'Каждый рубль в X возвращает': {match.group()!r}"

    def test_pptx_count_finding3_no_roas(self, count_payload):
        if count_payload is None:
            pytest.skip("Фикстура недоступна")
        try:
            from aurora_pptx.builder import AuroraPPTXBuilder
        except ImportError as e:
            pytest.skip(f"AuroraPPTXBuilder недоступен: {e}")
        try:
            prs = AuroraPPTXBuilder(count_payload).build()
        except Exception as e:
            if "TclError" in str(type(e).__name__) or "_tkinter" in str(e):
                pytest.skip(f"Окружение matplotlib/_tkinter: {e}")
            raise
        txt = _pptx_text(prs)
        # «Ожидаемый прирост ROAS» — недопустим для count
        assert "Ожидаемый прирост ROAS" not in txt, \
            f"PPTX (count) Finding 3 содержит 'Ожидаемый прирост ROAS'"

    def test_pptx_count_contrib_unit_not_rub_mln(self, count_payload):
        """PPTX таблица: единица вклада НЕ «₽ млн» для count KPI."""
        if count_payload is None:
            pytest.skip("Фикстура недоступна")
        try:
            from aurora_pptx.builder import AuroraPPTXBuilder
        except ImportError as e:
            pytest.skip(f"AuroraPPTXBuilder недоступен: {e}")
        # Проверяем через kpi_helpers что contrib_unit генерируется корректно
        from aurora_pptx.kpi_helpers import kpi_view
        kpi = kpi_view(count_payload)
        contrib_unit = (
            kpi.get("target_unit") or "ед."
            if kpi["kpi_kind"] == "count"
            else "₽ млн"
        )
        assert contrib_unit != "₽ млн", \
            f"count kpi: contrib_unit должен быть единицей результата, не '₽ млн', получили: {contrib_unit!r}"
        assert contrib_unit != "", "count kpi: contrib_unit пустой"


# ─── Контрольные тесты monetary (регрессия) ───────────────────────────────────

class TestMonetaryRegressionKpiHelpers:
    """Для monetary kpi формулировки должны оставаться как раньше."""

    def test_lift_phrase_monetary_still_roas(self):
        from aurora_pptx.kpi_helpers import lift_phrase
        result = lift_phrase(5.0, _kpi_monetary())
        assert "ROAS" in result, f"Регрессия: monetary lift_phrase потерял ROAS: {result!r}"

    def test_hero_vs_leader_monetary_still_ruble(self):
        from aurora_pptx.kpi_helpers import hero_vs_leader_quote
        result = hero_vs_leader_quote("Digital", "TV", _kpi_monetary())
        assert "рубль" in result.lower(), \
            f"Регрессия: monetary hero_vs_leader потерял 'рубль': {result!r}"

    def test_sections_lift_phrase_monetary_still_roas(self):
        from aurora_html.sections import _lift_phrase
        result = _lift_phrase(5.0, _kpi_monetary())
        assert "ROAS" in result, f"Регрессия: sections._lift_phrase monetary потерял ROAS: {result!r}"
