"""
Тесты секции «Прогноз на будущий период» (E5, 2026-07-10).

Покрывают:
  (a) PPTX с forecast → лишний слайд + проверка overflow
  (b) PPTX без forecast → кол-во слайдов не меняется (INV-50)
  (c) Оговорки попадают в XML
  (d) _page_shift == 3 при backtest + gen_compare + forecast
  (e) HTML render_forecast_plan возвращает "" при отсутствии forecast
  (f) HTML render_forecast_plan возвращает таблицу при наличии forecast
"""
import copy
import json
import os
import tempfile

import pytest

from aurora_pptx.builder import AuroraPPTXBuilder
from aurora_html.sections import render_forecast_plan

_HERE = os.path.dirname(os.path.abspath(__file__))
_FIXTURE = os.path.join(_HERE, "fixtures", "kagocel_builder_payload.json")

# ─── Минимальный валидный forecast ────────────────────────────────────────────

FORECAST_DATA = {
    "status": "ok",
    "scenarios": [
        {
            "name": "Базовый",
            "variant_id": "v1",
            "predictions": [100, 110, 120, 130],
            "ci_low": [90, 99, 108, 117],
            "ci_high": [110, 121, 132, 143],
            "total_kpi": 460.0,
            "total_spend_money": 5_000_000.0,
            "roas_money": 0.092,
            "period_labels": ["Янв", "Фев", "Мар", "Апр"],
            "disclaimers": ["Прогноз при неизменных условиях рынка"],
        }
    ],
    "historical_actual": [80, 85, 90, 95],
    "historical_dates": ["Сен", "Окт", "Ноя", "Дек"],
    "cutoff_index": 4,
    "accepted_variant": "v1",
    "disclaimers": ["Прогноз при неизменных условиях рынка"],
}


# ─── Fixture ──────────────────────────────────────────────────────────────────

@pytest.fixture(scope="module")
def base_payload():
    with open(_FIXTURE, encoding="utf-8") as f:
        return json.load(f)


def _build_deck(payload, path):
    prs = AuroraPPTXBuilder(payload).build()
    prs.save(path)
    return prs


def _xml_text(prs):
    """Собрать весь текст из XML всех слайдов."""
    from pptx.oxml.ns import qn
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
    return "\n".join(texts)


# ─── (a) С forecast → дополнительный слайд ───────────────────────────────────

def test_forecast_slide_added(base_payload, tmp_path):
    """При наличии forecast данных колода получает на 1 слайд больше."""
    # Базовая колода без forecast
    payload_base = copy.deepcopy(base_payload)
    payload_base.pop("forecast", None)
    out_base = str(tmp_path / "deck_base.pptx")
    prs_base = _build_deck(payload_base, out_base)
    n_base = len(prs_base.slides)

    # Колода с forecast
    payload_fc = copy.deepcopy(base_payload)
    payload_fc["forecast"] = FORECAST_DATA
    out_fc = str(tmp_path / "deck_forecast.pptx")
    prs_fc = _build_deck(payload_fc, out_fc)
    n_fc = len(prs_fc.slides)

    assert n_fc == n_base + 1, (
        f"Ожидали {n_base + 1} слайдов с forecast, получили {n_fc}"
    )


def test_forecast_slide_overflow_clean(base_payload, tmp_path):
    """Колода с forecast не должна иметь overflow."""
    from aurora_pptx.check_overflow import check

    payload_fc = copy.deepcopy(base_payload)
    payload_fc["forecast"] = FORECAST_DATA
    out = str(tmp_path / "deck_forecast_overflow.pptx")
    _build_deck(payload_fc, out)
    issues, n_slides = check(out)
    assert n_slides > 0
    detail = "\n".join(f"  слайд {s}: [{k}] {d}" for s, k, d in issues)
    assert issues == [], f"PPTX с forecast получил overflow/overlap:\n{detail}"


def test_forecast_slide_contains_title(base_payload, tmp_path):
    """Слайд прогноза должен содержать текст «Прогноз»."""
    payload_fc = copy.deepcopy(base_payload)
    payload_fc["forecast"] = FORECAST_DATA
    out = str(tmp_path / "deck_forecast_title.pptx")
    prs = _build_deck(payload_fc, out)
    all_text = _xml_text(prs)
    assert "Прогноз" in all_text, "Текст «Прогноз» не найден в XML колоды"


# ─── (b) Без forecast → кол-во слайдов не меняется (INV-50) ─────────────────

def test_no_forecast_no_extra_slide(base_payload, tmp_path):
    """Без forecast дека не получает лишний слайд-суррогат (INV-50)."""
    payload_no = copy.deepcopy(base_payload)
    payload_no.pop("forecast", None)

    out1 = str(tmp_path / "deck_no_fc_1.pptx")
    out2 = str(tmp_path / "deck_no_fc_2.pptx")
    prs1 = _build_deck(payload_no, out1)
    prs2 = _build_deck(payload_no, out2)
    assert len(prs1.slides) == len(prs2.slides), "Без forecast кол-во слайдов нестабильно"

    # С пустым / невалидным forecast — тот же эффект
    payload_invalid = copy.deepcopy(base_payload)
    payload_invalid["forecast"] = {"status": "error"}
    out3 = str(tmp_path / "deck_invalid_fc.pptx")
    prs3 = _build_deck(payload_invalid, out3)
    assert len(prs3.slides) == len(prs1.slides), (
        "Невалидный forecast породил лишний слайд-суррогат (INV-50)"
    )


# ─── (c) Оговорки попадают в XML ─────────────────────────────────────────────

def test_forecast_disclaimers_in_xml(base_payload, tmp_path):
    """Оговорки из forecast.disclaimers должны присутствовать в тексте слайда."""
    payload_fc = copy.deepcopy(base_payload)
    payload_fc["forecast"] = copy.deepcopy(FORECAST_DATA)
    payload_fc["forecast"]["disclaimers"] = ["неизменных условиях"]
    payload_fc["forecast"]["scenarios"][0]["disclaimers"] = ["неизменных условиях"]

    out = str(tmp_path / "deck_disclaimers.pptx")
    prs = _build_deck(payload_fc, out)
    all_text = _xml_text(prs)
    assert "неизменных условиях" in all_text, (
        "Оговорки из forecast не попали в XML колоды"
    )


# ─── (d) _page_shift == 3 при backtest + gen_compare + forecast ──────────────

def test_page_shift_all_three(base_payload):
    """При backtest + gen_compare + forecast _page_shift должен быть 3."""
    payload = copy.deepcopy(base_payload)

    # Минимальный валидный backtest
    payload["backtest"] = {
        "status": "ok",
        "windows": [{"window": "Q1", "actual_total": 100, "predicted_total": 98,
                     "pi_low_total": 80, "pi_high_total": 120, "hit_total": True}],
        "windows_hit_total": 1,
        "windows_with_interval": 1,
        "verdict": "validated",
        "granularity": "M",
        "horizon_periods": 3,
        "mape_model": 5.0,
        "mape_naive_best": 10.0,
        "coverage_per_period": 0.9,
        "n_holdout_points_with_interval": 3,
    }

    # Минимальный валидный gen_compare
    payload["generation_compare"] = {
        "status": "ok",
        "channels": [{"name": "TV", "roi_old": 1.2, "roi_new": 1.5,
                      "verdict": "stable", "verdict_ru": "стабильно"}],
        "summary": {"counts": {"stable": 1, "shift_within_ci": 0, "shift_strong": 0}},
        "baseline": {"timestamp": "2026-01-01T00:00:00"},
    }

    # Forecast
    payload["forecast"] = FORECAST_DATA

    builder = AuroraPPTXBuilder(payload)
    assert builder._page_shift == 3, (
        f"Ожидали _page_shift=3 при backtest+gen_compare+forecast, "
        f"получили {builder._page_shift}"
    )
    assert builder.total_slides == 15, (
        f"Ожидали total_slides=15, получили {builder.total_slides}"
    )


# ─── (e) HTML: нет forecast → render_forecast_plan возвращает "" ─────────────

def test_html_render_forecast_empty_without_data():
    """render_forecast_plan возвращает '' при отсутствии forecast (INV-50)."""
    assert render_forecast_plan({}) == ""
    assert render_forecast_plan({"forecast": None}) == ""
    assert render_forecast_plan({"forecast": {}}) == ""
    assert render_forecast_plan({"forecast": {"status": "error"}}) == ""
    assert render_forecast_plan({"forecast": {"status": "ok", "scenarios": []}}) == ""


# ─── (f) HTML: есть forecast → render_forecast_plan возвращает таблицу ───────

def test_html_render_forecast_with_data():
    """render_forecast_plan возвращает HTML-таблицу при наличии forecast."""
    ctx = {"forecast": FORECAST_DATA}
    html = render_forecast_plan(ctx)
    assert html != "", "render_forecast_plan вернул '' при валидных данных"
    assert "Прогноз" in html, "Заголовок «Прогноз» не найден в HTML"
    assert "Базовый" in html, "Название сценария «Базовый» не найден в HTML"
    assert "неизменных условиях" in html, "Оговорка не попала в HTML"
    assert "<table" in html, "Таблица сценариев не найдена в HTML"
    assert "★" in html, "Звёздочка accepted-сценария не найдена в HTML"


def test_html_render_forecast_multiple_scenarios():
    """Несколько сценариев: показываем до 4, остальные отбрасываем."""
    fc = copy.deepcopy(FORECAST_DATA)
    fc["scenarios"] = [
        {**fc["scenarios"][0], "name": f"Сценарий {i}", "variant_id": f"v{i}"}
        for i in range(6)
    ]
    fc["accepted_variant"] = "v2"
    html = render_forecast_plan({"forecast": fc})
    # Первые 4 сценария должны быть, 5-й и 6-й — нет
    assert "Сценарий 0" in html
    assert "Сценарий 3" in html
    assert "Сценарий 4" not in html
    assert "Сценарий 5" not in html


# ─── scenarios_comparison_chart ───────────────────────────────────────────────

def test_scenarios_comparison_chart_basic():
    """scenarios_comparison_chart генерирует непустой base64 PNG для 1 сценария."""
    from charts.generators import scenarios_comparison_chart
    import base64

    scenarios = [
        {
            "name": "Базовый",
            "total_kpi": 460.0,
            "total_spend_money": 5_000_000.0,
            "roas_money": 0.092,
        }
    ]
    result = scenarios_comparison_chart(scenarios)
    assert result != "", "scenarios_comparison_chart вернул пустую строку для 1 сценария"
    # Проверяем что это валидный base64
    decoded = base64.b64decode(result)
    assert len(decoded) > 100, "base64-строка слишком короткая, PNG невалиден"


def test_scenarios_comparison_chart_five_or_more():
    """При >=5 вариантах возвращается валидный base64 (горизонтальный layout)."""
    from charts.generators import scenarios_comparison_chart
    import base64

    scenarios = [
        {
            "name": f"Вариант {i + 1}",
            "total_kpi": float(400 + i * 20),
            "total_spend_money": float(5_000_000 + i * 500_000),
            "roas_money": round(0.08 + i * 0.01, 3),
        }
        for i in range(6)
    ]
    result = scenarios_comparison_chart(scenarios)
    assert result != "", "scenarios_comparison_chart вернул '' для 6 вариантов"
    decoded = base64.b64decode(result)
    assert len(decoded) > 100, "PNG для 6 вариантов невалиден"


def test_scenarios_comparison_chart_with_ci():
    """График с доверительными интервалами не падает, CI берётся из predictions_ci_low/high."""
    from charts.generators import scenarios_comparison_chart
    import base64

    scenarios = [
        {
            "name": "С интервалами",
            "total_kpi": 460.0,
            "predictions_ci_low": [90, 99, 108, 117],
            "predictions_ci_high": [110, 121, 132, 143],
        },
        {
            "name": "Без интервалов",
            "total_kpi": 440.0,
        },
    ]
    result = scenarios_comparison_chart(scenarios, kpi_label="Продажи")
    assert result != "", "scenarios_comparison_chart упал при смешанном наличии CI"
    base64.b64decode(result)  # не должно кидать исключение


def test_scenarios_comparison_chart_empty_returns_empty():
    """Пустой список сценариев возвращает пустую строку, не падает."""
    from charts.generators import scenarios_comparison_chart

    assert scenarios_comparison_chart([]) == "", "Пустой список должен возвращать ''"
    assert scenarios_comparison_chart(None) == "", "None должен возвращать ''"


# ─── HTML: scenarios_comparison_chart встраивается при ≥2 вариантах ──────────

def test_html_forecast_chart_present_for_two_or_more_scenarios():
    """При ≥2 сценариях render_forecast_plan должен содержать <img> с base64 PNG."""
    fc = copy.deepcopy(FORECAST_DATA)
    fc["scenarios"] = [
        {**fc["scenarios"][0], "name": "Базовый",  "variant_id": "v1", "total_kpi": 460.0},
        {**fc["scenarios"][0], "name": "Агрессивный", "variant_id": "v2", "total_kpi": 520.0},
    ]
    fc["accepted_variant"] = "v1"
    html = render_forecast_plan({"forecast": fc})
    assert '<img' in html, "При ≥2 сценариях ожидаем <img> в HTML"
    assert 'data:image/png;base64,' in html, "Ожидаем base64 PNG data-URI в <img>"


def test_html_forecast_chart_absent_for_one_scenario():
    """При одном сценарии <img> не добавляется."""
    fc = copy.deepcopy(FORECAST_DATA)
    assert len(fc["scenarios"]) == 1
    html = render_forecast_plan({"forecast": fc})
    assert html != "", "Секция должна рендериться даже при 1 сценарии"
    assert '<img' not in html, "При 1 сценарии <img> не должен присутствовать"


# ─── HTML: render_retro_insights ──────────────────────────────────────────────

def test_retro_insights_empty_for_reliable_model():
    """Блок «Что улучшить» пустой при reliable-модели."""
    from aurora_html.sections import render_retro_insights

    ctx = {"diagnostics": {"honesty_verdict": "reliable", "honesty_reasons": ["Всё хорошо"]}}
    assert render_retro_insights(ctx) == "", "При reliable-модели блок должен быть пустым"


def test_retro_insights_empty_without_diagnostics():
    """Блок «Что улучшить» пустой при отсутствии диагностики."""
    from aurora_html.sections import render_retro_insights

    assert render_retro_insights({}) == ""
    assert render_retro_insights({"diagnostics": {}}) == ""


def test_retro_insights_present_for_uncertain_model():
    """Блок «Что улучшить» появляется при uncertain-модели с reasons."""
    from aurora_html.sections import render_retro_insights

    ctx = {
        "diagnostics": {
            "honesty_verdict": "uncertain",
            "honesty_reasons": ["Мало наблюдений", "Широкие интервалы"],
            "thinness_cap": 50,
            "ratio": 1.8,
        }
    }
    html = render_retro_insights(ctx)
    assert html != "", "При uncertain-модели блок должен присутствовать"
    assert "Мало наблюдений" in html, "Причины honesty_reasons должны быть в тексте"
    assert "Что улучшить" in html, "Заголовок «Что улучшить» должен быть в тексте"


def test_retro_insights_present_for_unreliable_model():
    """Блок «Что улучшить» появляется при unreliable-модели."""
    from aurora_html.sections import render_retro_insights

    ctx = {
        "diagnostics": {
            "honesty_verdict": "unreliable",
            "honesty_reasons": ["R-hat > 1.05", "ESS < 100"],
            "r_squared": 0.45,
            "mape_pct": 28.0,
        }
    }
    html = render_retro_insights(ctx)
    assert html != "", "При unreliable-модели блок должен присутствовать"
    assert "R-hat" in html
    assert "R²" in html, "Должен быть пункт про низкий R²"
    assert "MAPE" in html, "Должен быть пункт про высокий MAPE"


def test_retro_insights_preflight_fail():
    """Блок включает пункт про provail приоров при prior_predictive_status=fail."""
    from aurora_html.sections import render_retro_insights

    ctx = {
        "diagnostics": {
            "honesty_verdict": "uncertain",
            "preflight": {
                "prior_predictive_status": "fail",
                "prior_predictive_coverage": 0.42,
            },
        }
    }
    html = render_retro_insights(ctx)
    assert html != ""
    assert "Априорные предположения расходятся" in html
    assert "42%" in html
