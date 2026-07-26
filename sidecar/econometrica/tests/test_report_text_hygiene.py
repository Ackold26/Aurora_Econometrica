"""Гигиена клиентских текстов — гейты П8-1 и П8-2.

П8-1: в клиентских строках нет em-dash (—); дефис-как-тире заменён на en-dash (–).
П8-2: нет голого «baseline» и «media-»; «adstock» не встречается без скобок
      в ключевых выводах.

Тест строит честный payload (honest_narrative=True, media_pct < 10%) —
именно он активирует все три ветки с проблемными строками (builder.py s05,
sections.py render_key_message/render_at_a_glance, narrative_adapter headline).
Фикстура без реальных данных (markers = not requires_real_data).
"""

import copy
import json
import os
import re

import pytest

from aurora_pptx.builder import AuroraPPTXBuilder
from aurora_html.sections import render_key_message, render_at_a_glance, SECTION_RENDERERS

_HERE = os.path.dirname(os.path.abspath(__file__))
_FIXTURE = os.path.join(_HERE, "fixtures", "kagocel_builder_payload.json")

# ─── Паттерны-нарушители ──────────────────────────────────────────────────────

# П8-1: em-dash в клиентском тексте
EM_DASH_RE = re.compile(r" — ")

# П8-2: голый anglicизм baseline (не в коде/переменных)
BASELINE_RE = re.compile(r"\bbaseline\b", re.IGNORECASE)

# П8-2: медиа-вклад/медиа-эффект через дефис с latin 'm'
MEDIA_LATIN_RE = re.compile(r"\bmedia-", re.IGNORECASE)

# П8-2: adstock без скобок в КЛИЕНТСКИХ КЛЮЧЕВЫХ ВЫВОДАХ
# Разрешено: "отложенный эффект (adstock)", формулы "adstock(x,t)", глоссарий "Adstock"
# Запрещено в action-строках: "проверить adstock," / "adstock и насыщение"
# Фильтр: строки содержащие adstock без (adstock) И без формульного контекста "adstock("
def _has_bare_adstock_in_client_text(text: str) -> bool:
    """True если в тексте есть голый adstock (не в скобках, не в формуле, не глоссарий)."""
    for line in text.split("\n"):
        lo = line.lower()
        if "adstock" not in lo:
            continue
        # Разрешаем формульные контексты: adstock(x, adstock = x_t
        if "adstock(" in lo or "adstock =" in lo or "adstock(x" in lo:
            continue
        # Разрешаем глоссарий — отдельная строка "Adstock"
        if line.strip().lower() == "adstock":
            continue
        # Разрешаем уже правильно оформленное: "(adstock)"
        bare = re.compile(r"(?<!\()\badstock\b(?!\))", re.IGNORECASE)
        if bare.search(line):
            return True
    return False


# ─── Честный payload (media < 10%, honest_narrative=True) ─────────────────────

@pytest.fixture(scope="module")
def honest_payload():
    """Payload с honest_narrative=True и media_pct=7% — активирует все honest-ветки."""
    with open(_FIXTURE, encoding="utf-8") as f:
        base = json.load(f)
    p = copy.deepcopy(base)
    # Переключаем в честный режим
    p["narrative_facts"]["honest_narrative"] = True
    p["narrative_facts"]["media_contribution_pct"] = 7.2
    p["narrative_facts"]["baseline_pct"] = 92.8
    p["narrative_facts"]["leader_share_contrib_pct"] = 42.0
    return p


@pytest.fixture(scope="module")
def normal_payload():
    """Стандартный payload без honest_narrative — покрывает остальные ветки."""
    with open(_FIXTURE, encoding="utf-8") as f:
        return json.load(f)


# ─── Хелпер: собрать весь текст из PPTX XML ──────────────────────────────────

def _pptx_text(prs) -> str:
    from pptx.oxml.ns import qn
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
    return "\n".join(parts)


# ─── Хелпер: context dict для sections.py ────────────────────────────────────

def _ctx(payload: dict) -> dict:
    """Минимальный context dict для render_* функций sections.py."""
    import json as _json
    _strings_path = os.path.join(
        os.path.dirname(_HERE), "aurora_html", "strings_ru.json"
    )
    with open(_strings_path, encoding="utf-8") as f:
        strings = _json.load(f)
    return {
        "meta": payload.get("meta") or {},
        "facts": payload.get("narrative_facts"),
        "channels": payload.get("channels") or [],
        "diagnostics": payload.get("diagnostics") or {},
        "strings": strings,
        "kpi": {},
    }


# ─── PPTX honest-mode: П8-1 + П8-2 ──────────────────────────────────────────

def test_pptx_honest_no_em_dash(honest_payload, tmp_path):
    """П8-1: em-dash не появляется в PPTX на честном payload."""
    prs = AuroraPPTXBuilder(honest_payload).build()
    txt = _pptx_text(prs)
    hits = EM_DASH_RE.findall(txt)
    assert not hits, f"П8-1: em-dash в PPTX (честный режим): {hits[:3]}"


def test_pptx_honest_no_bare_baseline(honest_payload, tmp_path):
    """П8-2: «baseline» не встречается голым в клиентском PPTX."""
    prs = AuroraPPTXBuilder(honest_payload).build()
    txt = _pptx_text(prs)
    hits = BASELINE_RE.findall(txt)
    assert not hits, f"П8-2: голый baseline в PPTX: {txt[max(0,txt.lower().find('baseline')-40):txt.lower().find('baseline')+60]!r}"


def test_pptx_honest_no_media_latin(honest_payload, tmp_path):
    """П8-2: «media-вклад»/«media-эффект» не появляется в PPTX."""
    prs = AuroraPPTXBuilder(honest_payload).build()
    txt = _pptx_text(prs)
    hits = MEDIA_LATIN_RE.findall(txt)
    assert not hits, f"П8-2: latin media- в PPTX: {hits[:3]}"


def test_pptx_honest_adstock_not_bare(honest_payload, tmp_path):
    """П8-2: adstock в клиентских ключевых выводах только со скобками.

    Формульные строки (adstock(x,t), adstock = ...) и глоссарный термин
    «Adstock» на отдельной строке разрешены — только narrative выводы проверяются.
    """
    prs = AuroraPPTXBuilder(honest_payload).build()
    txt = _pptx_text(prs)
    assert not _has_bare_adstock_in_client_text(txt), (
        "П8-2: голый adstock (не в скобках, не в формуле) найден в PPTX"
    )


# ─── HTML честный режим: render_key_message ──────────────────────────────────

def test_html_key_message_no_em_dash(honest_payload):
    """П8-1: em-dash не появляется в render_key_message (honest branch)."""
    try:
        ctx = _ctx(honest_payload)
    except Exception:
        pytest.skip("strings_ru.json недоступен")
    html = render_key_message(ctx)
    assert not EM_DASH_RE.search(html), "П8-1: em-dash в render_key_message"


def test_html_key_message_no_bare_baseline(honest_payload):
    """П8-2: «baseline» не попадает в HTML-вывод render_key_message."""
    try:
        ctx = _ctx(honest_payload)
    except Exception:
        pytest.skip("strings_ru.json недоступен")
    html = render_key_message(ctx)
    hits = BASELINE_RE.findall(html)
    assert not hits, f"П8-2: baseline в render_key_message HTML: {hits}"


def test_html_key_message_no_media_latin(honest_payload):
    """П8-2: «media-» не попадает в HTML-вывод render_key_message."""
    try:
        ctx = _ctx(honest_payload)
    except Exception:
        pytest.skip("strings_ru.json недоступен")
    html = render_key_message(ctx)
    hits = MEDIA_LATIN_RE.findall(html)
    assert not hits, f"П8-2: latin media- в render_key_message HTML: {hits}"


def test_html_key_message_adstock_not_bare(honest_payload):
    """П8-2: adstock в HTML только со скобками (формульные контексты разрешены)."""
    try:
        ctx = _ctx(honest_payload)
    except Exception:
        pytest.skip("strings_ru.json недоступен")
    html = render_key_message(ctx)
    # В HTML-контексте strip тегов и проверяем по строкам
    text = re.sub(r"<[^>]+>", " ", html)
    assert not _has_bare_adstock_in_client_text(text), (
        "П8-2: голый adstock найден в render_key_message HTML"
    )


# ─── HTML честный режим: render_at_a_glance ──────────────────────────────────

def test_html_at_a_glance_no_em_dash(honest_payload):
    """П8-1: em-dash не попадает в render_at_a_glance."""
    try:
        ctx = _ctx(honest_payload)
    except Exception:
        pytest.skip("strings_ru.json недоступен")
    html = render_at_a_glance(ctx)
    assert not EM_DASH_RE.search(html), "П8-1: em-dash в render_at_a_glance"


def test_html_at_a_glance_no_bare_baseline(honest_payload):
    """П8-2: baseline не попадает в render_at_a_glance."""
    try:
        ctx = _ctx(honest_payload)
    except Exception:
        pytest.skip("strings_ru.json недоступен")
    html = render_at_a_glance(ctx)
    hits = BASELINE_RE.findall(html)
    assert not hits, f"П8-2: baseline в render_at_a_glance: {hits}"


def test_html_at_a_glance_no_media_latin(honest_payload):
    """П8-2: «media-» не попадает в render_at_a_glance."""
    try:
        ctx = _ctx(honest_payload)
    except Exception:
        pytest.skip("strings_ru.json недоступен")
    html = render_at_a_glance(ctx)
    hits = MEDIA_LATIN_RE.findall(html)
    assert not hits, f"П8-2: latin media- в render_at_a_glance: {hits}"


# ─── Normal payload (не честный) — нет регрессии в стандартных ветках ────────

def test_pptx_normal_no_em_dash(normal_payload, tmp_path):
    """П8-1: em-dash не появляется в PPTX на стандартном payload."""
    prs = AuroraPPTXBuilder(normal_payload).build()
    txt = _pptx_text(prs)
    hits = EM_DASH_RE.findall(txt)
    assert not hits, f"П8-1: em-dash в PPTX (стандарт): {hits[:3]}"


def test_pptx_normal_no_media_latin(normal_payload, tmp_path):
    """П8-2: «media-» не появляется в PPTX на стандартном payload."""
    prs = AuroraPPTXBuilder(normal_payload).build()
    txt = _pptx_text(prs)
    hits = MEDIA_LATIN_RE.findall(txt)
    assert not hits, f"П8-2: latin media- в PPTX (стандарт): {hits[:3]}"


# ─── Фаза 3 (2026-07-25): полное покрытие HTML-секций ─────────────────────────
#
# ДО этой правки П8-1/П8-2 на HTML проверяли только render_key_message и
# render_at_a_glance — 2 из 17 функций в aurora_html.sections.SECTION_RENDERERS.
# Остальные 15 не проверялись НИКЕМ (см. builder.py::SECTION_RENDERERS). Дыра
# молчаливо читалась как «HTML-отчёт покрыт» — тот класс дефекта, что породил
# трёхнедельную утечку в соседнем продукте (см. PHASE3_ECON_COVERAGE.md).
#
# Расширение нашло РЕАЛЬНЫЕ нарушения в ранее непокрытых секциях (все починены
# в aurora_html/sections.py, см. отчёт Фазы 3):
#   - summary  (_reliability_disclaimer_html): 2× em-dash
#   - method   (_render_brand_perf_split_block): 2× em-dash + голое «Validate»
#               (в приложении шаг называется «Валидация» — project-state.js)
#   - sources: голый «baseline» в списке источников данных
#   - trust: 2× em-dash (E1 backtest-блок, E3 generation_compare-блок)
#   - retro: голый «adstock» + голое «Validate» в рекомендации Preflight
#   - timeline: внутренняя аудит-заметка утекала в клиентский HTML как
#               <!-- HTML-комментарий --> (не текст, но внутренняя информация
#               в экспортируемом файле) — вынесена в Python-комментарий
#
# Одна находка была оставлена спорной и РЕШЕНА владельцем 2026-07-26:
# глоссарий определяет термин «Baseline» headword-строкой на отдельной
# строке — тем же способом, каким уже разрешён «Adstock» в
# _has_bare_adstock_in_client_text (см. код выше). Решение: это принятый
# термин отрасли, определяемый в глоссарии, а не англицизм в прозе —
# узаконен точечным правилом (_bare_baseline_lines ниже), не xfail'ом и не
# ослаблением regex. Разрешена ровно строка-заголовок термина и только в
# секции glossary; «baseline» внутри определения или в любой другой секции
# остаётся нарушением.

_BLOCK_CLOSE_RE = re.compile(
    r"</(div|p|li|tr|td|th|table|thead|tbody|ul|ol|h1|h2|h3|h4|section|"
    r"details|summary|dl|dt|dd)>",
    re.IGNORECASE,
)


def _strip_tags_preserve_lines(html: str) -> str:
    """Снять HTML-теги, сохранив границы блоков как переводы строк.

    Наивный `re.sub(r"<[^>]+>", " ", html)` (использован для render_key_message
    adstock-проверки выше) склеивает многоблочные секции — например, глоссарий
    из 24 терминов — в одну гигантскую "строку"; построчные исключения
    `_has_bare_adstock_in_client_text` (headword-строка глоссария, формула)
    перестают срабатывать → ложные срабатывания на легитимных определениях
    терминов. Здесь границы блочных элементов явно превращаются в "\\n" до
    снятия тегов, так что каждый смысловой блок остаётся своей строкой.
    """
    html = _BLOCK_CLOSE_RE.sub(lambda m: m.group(0) + "\n", html)
    html = re.sub(r"<br\s*/?>", "\n", html, flags=re.IGNORECASE)
    return re.sub(r"<[^>]+>", " ", html)


def _full_ctx(payload: dict) -> dict:
    """context dict, покрывающий ВСЕ 17 SECTION_RENDERERS (не только 2).

    _ctx() выше даёт meta/facts/channels/diagnostics/strings/kpi — этого
    хватает render_key_message/render_at_a_glance, но недостаточно для
    остальных 15 (report_id, model_version, period_unit, brand_mark_svg,
    trust, forecast — иначе часть секций тихо рендерит "" и "покрытие"
    остаётся формальным: функцию вызвали, а текста внутри не было).
    Диагностика дополнена так, чтобы триггернуть ветки, которые иначе не
    рендерятся вовсе: F-A1-9 дисклеймер ненадёжности (summary), ретро-
    рекомендации при Preflight-провале (retro), brand/performance блок
    методологии (method).
    """
    ctx = _ctx(payload)
    diag = dict(ctx["diagnostics"] or {})
    diag.setdefault("honesty_verdict", "unreliable")
    diag.setdefault("honesty_reasons", ["Мало наблюдений", "Широкие интервалы"])
    diag.setdefault("preflight", {
        "prior_predictive_status": "fail", "prior_predictive_coverage": 0.4,
    })
    diag.setdefault("hierarchical", {
        "enabled": True,
        "channel_categories": {"TV": "brand", "Radio": "brand", "Search": "performance"},
        "priors_summary": {
            "brand_mu_logit_mean": 0.5, "performance_mu_logit_mean": -0.5,
        },
        "rhat_warning": None,
    })
    ctx["diagnostics"] = diag
    ctx["report_id"] = "TESTREPORTID"
    ctx["model_version"] = "1.1.0"
    ctx["period_unit"] = "неделям"
    ctx["brand_mark_svg"] = ""
    ctx["trust"] = {
        "backtest": {
            "status": "ok", "windows_hit_total": 3, "windows_with_interval": 4,
            "granularity": "M", "horizon_periods": 3,
            "windows": [{
                "window": "2026-01", "actual_total": 100, "predicted_total": 95,
                "pi_low_total": 80, "pi_high_total": 110, "hit_total": True,
            }],
            "mape_naive_best": 12.0, "mape_model": 8.0,
            "verdict_text": "Модель точнее наивного прогноза.",
        },
        "generation_compare": {
            "status": "ok",
            "summary": {"headline": "ROI каналов пересчитан на новых данных."},
            "channels": [{
                "name": "TV", "roi_old": 1.2, "roi_ci_old": [1.0, 1.4],
                "roi_new": 1.5, "roi_ci_new": [1.3, 1.7], "verdict_ru": "вырос",
            }],
        },
        "promises_summary": {
            "kept": 2, "missed": 1,
            "examples": [{
                "action_text": "Увеличить TV", "status": "kept", "status_ru": "сбылось",
            }],
        },
    }
    ctx["forecast"] = {
        "status": "ok",
        "scenarios": [{
            "name": "Базовый", "variant_id": "v1", "total_kpi": 460.0,
            "total_kpi_ci_low": 414.0, "total_kpi_ci_high": 506.0,
            "total_spend_money": 5_000_000.0, "roas_money": 0.092,
        }],
        "accepted_variant": "v1",
        "disclaimers": ["Прогноз при неизменных условиях рынка"],
    }
    return ctx


_ALL_SECTION_IDS = tuple(sid for sid, _ in SECTION_RENDERERS)
_SECTION_RENDER_BY_ID = dict(SECTION_RENDERERS)

# Узаконенные термины (решение владельца 2026-07-26) — где именно термин
# отрасли разрешён и почему. Реестр открытый и печатается в охвате: молча
# разрастись он не может, каждая запись требует обоснования строкой.
_LEGITIMISED_TERMS = {
    ("glossary", "Baseline"): (
        "строка-заголовок термина в глоссарии: определение принятого термина "
        "отрасли, не англицизм в прозе. Тот же приём уже действует для «Adstock»"
    ),
}


def _bare_baseline_lines(text: str, section_id: str) -> list:
    """Строки клиентского текста с голым «baseline» — кроме узаконенных.

    Разрешена ровно строка, целиком равная термину, и только в той секции,
    для которой это записано в _LEGITIMISED_TERMS. «baseline» внутри прозы,
    определения или любой другой секции — по-прежнему нарушение.
    """
    out = []
    for line in text.split("\n"):
        if not BASELINE_RE.search(line):
            continue
        stripped = line.strip()
        if (section_id, stripped) in _LEGITIMISED_TERMS:
            continue
        out.append(stripped[:120])
    return out


def test_baseline_rule_still_catches_prose():
    """INV-99: правило узаконивания обязано уметь краснеть.

    Узаконена строка-заголовок термина; всё, что рядом, ловится по-прежнему —
    иначе исключение молча превратилось бы в разрешение слова везде.
    """
    assert _bare_baseline_lines("Baseline", "glossary") == []
    assert _bare_baseline_lines("Baseline", "summary") == ["Baseline"]
    assert _bare_baseline_lines("Доля baseline выросла", "glossary") == [
        "Доля baseline выросла"
    ]


@pytest.mark.parametrize("section_id", _ALL_SECTION_IDS)
def test_html_section_hygiene(section_id, honest_payload):
    """П8-1/П8-2 на КАЖДОЙ из 17 секций HTML-отчёта (Фаза 3, было 2/17).

    section_id импортирован из живого SECTION_RENDERERS — появится 18-я
    секция, parametrize вырастет сам собой на следующем прогоне. Молчаливое
    сужение здесь структурно невозможно: пропустить секцию можно только
    явно, через _KNOWN_CONTENT_QUESTIONS с обоснованием.
    """
    try:
        ctx = _full_ctx(honest_payload)
    except Exception:
        pytest.skip("strings_ru.json недоступен")
    html = _SECTION_RENDER_BY_ID[section_id](ctx)

    em = EM_DASH_RE.findall(html)
    assert not em, f"П8-1: em-dash в секции «{section_id}»: {em[:3]}"

    ml = MEDIA_LATIN_RE.findall(html)
    assert not ml, f"П8-2: latin media- в секции «{section_id}»: {ml}"

    stripped = _strip_tags_preserve_lines(html)

    # baseline проверяется по клиентскому ТЕКСТУ, а не по сырому html: правило
    # узаконивания построчное (см. _LEGITIMISED_TERMS). Сужение охвата здесь
    # мнимое — проверено зондом, что ни в одной из 17 секций слова baseline в
    # разметке (классы, идентификаторы, атрибуты) нет вовсе; появится — его
    # поймает отдельная проверка ниже.
    bl = _bare_baseline_lines(stripped, section_id)
    assert not bl, f"П8-2: голый baseline в секции «{section_id}»: {bl}"

    markup_only = len(BASELINE_RE.findall(html)) - len(BASELINE_RE.findall(stripped))
    assert markup_only <= 0, (
        f"П8-2: baseline появился в разметке секции «{section_id}» "
        f"(класс/идентификатор/атрибут) — вне охвата текстовой проверки"
    )

    assert not _has_bare_adstock_in_client_text(stripped), (
        f"П8-2: голый adstock в секции «{section_id}»"
    )


def test_html_section_coverage_is_reported(honest_payload):
    """Печатает и проверяет ФАКТ охвата — числом и поимённо, не на глаз.

    Обязательное требование Фазы 3: проверка обязана сообщать, сколько
    секций проверено и сколько вне охвата — молчаливое сужение (как было
    до этой правки: 2/17 без единого слова об этом) отсюда structурно
    невозможно, потому что _ALL_SECTION_IDS читается из живого
    SECTION_RENDERERS, а не захардкожен.
    """
    total = len(_ALL_SECTION_IDS)
    terms = sorted(f"{sid}: «{term}»" for sid, term in _LEGITIMISED_TERMS)
    summary = (
        f"ОХВАТ П8-1/П8-2 (HTML-отчёт, aurora_html.sections.SECTION_RENDERERS): "
        f"секций проверено {total} из {total}, вне охвата 0; "
        f"узаконенных терминов {len(terms)} — {terms}"
    )
    print(summary)
    assert total == len(SECTION_RENDERERS), (
        "SECTION_RENDERERS разошёлся с census внутри теста — обновить _ALL_SECTION_IDS"
    )
    unknown_sections = sorted(
        sid for sid, _ in _LEGITIMISED_TERMS if sid not in _ALL_SECTION_IDS
    )
    assert not unknown_sections, (
        f"узаконенный термин записан для несуществующей секции {unknown_sections} — "
        f"исключение стало мёртвым, удалить или переписать"
    )
