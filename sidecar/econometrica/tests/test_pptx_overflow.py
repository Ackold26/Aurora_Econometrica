"""Overflow/overlap-гейт PPTX-отчёта Econometrica.

Портирован из Smart Analytica (2026-06-07). Econometrica прежде правил overflow
вручную магическими числами кегля (builder.py:1465) — без ассерта. Этот гейт
делает «текст не вылезает за низ контента и боксы не сталкиваются» проверяемым
фактом на ВСЕХ слайдах через реальные метрики шрифта (text_metrics, PIL).

Фикстура — реальный builder-payload проекта Кагоцел (…-0706-26, MQS 70 с
thinness_cap → включает оговорку F-DELIVERABLE-1). Детерминирована, не зависит от
APPDATA. Гейт регрессионно ловит будущие overflow на длинных RU-строках клиента.
"""
import json
import os
import tempfile

import pytest

from aurora_pptx.builder import AuroraPPTXBuilder
from aurora_pptx.check_overflow import check

_HERE = os.path.dirname(os.path.abspath(__file__))
_FIXTURE = os.path.join(_HERE, "fixtures", "kagocel_builder_payload.json")


def _build_deck(payload, path):
    prs = AuroraPPTXBuilder(payload).build()
    prs.save(path)
    return prs


@pytest.fixture(scope="module")
def payload():
    with open(_FIXTURE, encoding="utf-8") as f:
        return json.load(f)


def test_kagocel_deck_is_overflow_clean(payload, tmp_path):
    """Реальная колода Кагоцела — без переполнений и наложений (включая оговорку
    F-DELIVERABLE-1 на slide 11/13)."""
    out = str(tmp_path / "deck.pptx")
    _build_deck(payload, out)
    issues, n_slides = check(out)
    assert n_slides > 0
    detail = "\n".join(f"  слайд {s}: [{k}] {d}" for s, k, d in issues)
    assert issues == [], f"PPTX-отчёт получил overflow/overlap:\n{detail}"


def test_gate_catches_injected_overflow(payload, tmp_path):
    """Анти-вакуумность: гейт ОБЯЗАН поймать намеренно переполняющий текстбокс
    в контентной зоне (иначе зелёный тест ничего не значит)."""
    from pptx.util import Inches, Pt
    out = str(tmp_path / "deck_overflow.pptx")
    prs = _build_deck(payload, out)
    # На контентном слайде добавляем длинный абзац у самого низа зоны → вылезет за 7.05".
    slide = prs.slides[6]
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(6.6), Inches(6.0), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = ("Очень длинный сопроводительный текст, который заведомо не помещается "
                "в отведённую высоту и переносится на множество строк, выходя далеко "
                "за нижнюю границу контентной зоны слайда отчёта. " * 3)
    run.font.size = Pt(14)
    run.font.name = "Arial"
    prs.save(out)
    issues, _ = check(out)
    assert any(k == "OVERFLOW" for _, k, _ in issues), \
        "Гейт не поймал намеренное переполнение — проверка вакуумна"
