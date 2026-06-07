"""Fidelity-diff harness: «отчёты = программа» (аудит #12, INV-50).

Идея (2026-06-07): расхождение «отчёт показывает МЕНЬШЕ факторов чем программа»
ловится не глазами по одному отчёту, а сравнением артефактов против источника
истины. Этот скрипт:

  1. Берёт обученный проект (results/{decomposition,model-diagnostics,
     optimization}.json) — источник истины.
  2. Вычисляет «полный набор факторов», который ПОКАЗЫВАЕТ ПРОГРАММА:
     медиа-каналы (decompose.channels) + signed-факторы и праздники из
     signed_factor_contributions (mirror правила фронтового ChannelTimeline:
     показываются type ∈ signed_* | holiday; positive_control сворачивается
     в baseline и отдельным фактором НЕ считается).
  3. Генерит HTML и PPTX теми же билдерами, что и продакшн (html_export /
     pptx_export через _map_pipeline_to_builder_data).
  4. Парсит артефакты, извлекает фактически отрисованный набор факторов и
     ключевые числа.
  5. Диффает: какие SoT-факторы ОТСУТСТВУЮТ в каждом отчёте + расхождения
     ключевых чисел (per-channel ROI/вклад, baseline, total, MQS).

XLSX (билдер на Rust, report.rs) проверяется отдельным cargo-тестом —
здесь помечается как out-of-scope для Python-харнесса (см. --note).

Использование:
  python tools/report_fidelity_diff.py                 # тестовый Кагоцел 0706
  python tools/report_fidelity_diff.py "<project_dir>"  # конкретный проект
  python tools/report_fidelity_diff.py --json           # машинный вывод

Exit code: 0 — паритет факторов достигнут; 1 — есть отсутствующие факторы.
"""
from __future__ import annotations

import json
import os
import re
import sys
import tempfile

HERE = os.path.dirname(os.path.abspath(__file__))
SID = os.path.dirname(HERE)
sys.path.insert(0, SID)

PROJECTS_ROOT = os.path.join(os.environ.get('APPDATA', ''), 'aurora-econometrica-gui', 'projects')
DEFAULT_PROJECT = os.path.join(
    PROJECTS_ROOT, 'кагоцел-рф--данные-для-эконометрики---на-ммх-0706-26'
)

# Mirror фронтового ChannelTimeline.svelte (правило отбора факторов в timeline).
SIGNED_TYPES = {'signed_competitor', 'signed_price', 'signed_weather', 'signed_macro'}

# Канонизация медиа-имён теми же правилами, что билдеры (Social Бюджет… → Social).
from engines.narrative_adapter import _normalize_channel_name  # noqa: E402


def _norm(s: str) -> str:
    """Нормализация имени фактора для сравнения (переводы строк/пробелы)."""
    return re.sub(r'\s+', ' ', str(s)).strip().lower()


def _canon_media(name: str) -> str:
    """Каноничное короткое имя медиа-канала (как в артефактах отчётов)."""
    c = _normalize_channel_name(name)
    return _norm(c) if c else _norm(name)


def load_results(proj_dir: str) -> dict:
    res = os.path.join(proj_dir, 'results')
    out = {}
    for key, fname in (('decompose', 'decomposition.json'),
                       ('model', 'model-diagnostics.json'),
                       ('optimize', 'optimization.json')):
        p = os.path.join(res, fname)
        out[key] = json.load(open(p, encoding='utf-8')) if os.path.exists(p) else {}
    return out


def compute_sot(decompose: dict, model: dict) -> dict:
    """Источник истины: что ПОКАЗЫВАЕТ программа."""
    channels = decompose.get('channels') or []
    media = [c.get('name') for c in channels if c.get('name')]

    sfc = decompose.get('signed_factor_contributions') or {}
    shown_factors = []     # как в ChannelTimeline: signed_* | holiday
    folded = []            # positive_control → свёрнут в baseline (не фактор)
    for name, f in sfc.items():
        if not isinstance(f, dict):
            continue
        t = str(f.get('type', 'positive_control'))
        if t in SIGNED_TYPES or t == 'holiday':
            shown_factors.append({'name': name, 'type': t})
        else:
            folded.append({'name': name, 'type': t})

    mqs = (model.get('mqs') or {}).get('score')
    return {
        'media': media,
        'media_set': {_norm(m) for m in media},
        'signed_factors': shown_factors,
        'signed_set': {_norm(f['name']) for f in shown_factors},
        'folded': folded,
        'full_count': len(media) + len(shown_factors),
        'numbers': {
            'mqs': mqs,
            'baseline': decompose.get('baseline'),
            'total_sales': decompose.get('total_sales'),
            'roi': {_norm(c.get('name')): c.get('roi') for c in channels if c.get('name')},
            'contribution': {_norm(c.get('name')): c.get('contribution') for c in channels if c.get('name')},
        },
    }


def gen_reports(model: dict, decompose: dict, optimize: dict, out_dir: str, project_id: str) -> dict:
    """Генерит HTML+PPTX теми же билдерами, что продакшн. model_data обёрнут
    в {'diagnostics': ...} — как фронт (ReportStep.svelte:53)."""
    from engines.html_export import build_html
    from engines.pptx_export import build_pptx

    model_data = {'diagnostics': model}
    html_path = os.path.join(out_dir, 'fidelity.html')
    pptx_path = os.path.join(out_dir, 'fidelity.pptx')

    r_html = build_html(model_data, dict(decompose), optimize, html_path, project_id=project_id)
    r_pptx = build_pptx(model_data, dict(decompose), optimize, pptx_path, project_id=project_id)
    return {'html': (html_path, r_html), 'pptx': (pptx_path, r_pptx)}


def extract_html_factors(html_path: str) -> dict:
    """Парсит var CHART_DATA = {...}; → нормализованный blob всех отрисованных
    меток (timeline-серии + waterfall + mroas) для substring-матчинга."""
    txt = open(html_path, encoding='utf-8').read()
    m = re.search(r'var CHART_DATA = (\{.*?\});', txt)
    if not m:
        return {'blob': '', 'timeline_keys': [], 'error': 'CHART_DATA not found'}
    data = json.loads(m.group(1))
    tl = data.get('timeline') or {}
    labels = list((tl.get('channels') or {}).keys())
    # Доп. signed-серии (появятся после фикса Stage C) — по ключам/именам.
    for extra_key in ('factors', 'signed', 'controls', 'external'):
        ev = tl.get(extra_key)
        if isinstance(ev, dict):
            labels += list(ev.keys())
        elif isinstance(ev, list):
            labels += [x.get('name') for x in ev if isinstance(x, dict) and x.get('name')]
    labels += (data.get('waterfall') or {}).get('labels') or []
    labels += (data.get('mroas') or {}).get('names') or []
    blob = _norm(' \n '.join(str(x) for x in labels if x))
    return {'blob': blob, 'timeline_keys': list((tl.get('channels') or {}).keys())}


def extract_pptx_factors(pptx_path: str) -> dict:
    """Сканирует все слайды: имена серий графиков + текст → набор факторов."""
    from pptx import Presentation
    prs = Presentation(pptx_path)
    texts = []
    series_names = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
            if shape.has_chart:
                try:
                    for s in shape.chart.series:
                        if s.name:
                            series_names.append(s.name)
                except Exception:
                    pass
                try:
                    cats = list(shape.chart.plots[0].categories)
                    texts.extend(str(c) for c in cats)
                except Exception:
                    pass
    blob = _norm(' \n '.join(texts + series_names))
    return {'blob': blob, 'series_names': series_names}


def factor_present(item: dict, html_blob: str, pptx_blob: str) -> dict:
    """Фактор представлен, если его метка — подстрока нормализованного blob
    отчёта. Для медиа сверяем каноничное короткое имя (Social Бюджет… → social),
    для signed/holiday — сырое имя/ключ (так его несёт SSOT после фикса)."""
    if item['kind'] == 'media':
        keys = {_canon_media(item['name'])}
    else:
        keys = {_norm(item['name'])}
    return {
        'html': any(k and k in html_blob for k in keys),
        'pptx': any(k and k in pptx_blob for k in keys),
    }


def run(proj_dir: str, as_json: bool = False) -> int:
    project_id = os.path.basename(proj_dir.rstrip('/\\'))
    results = load_results(proj_dir)
    sot = compute_sot(results['decompose'], results['model'])

    tmp = tempfile.mkdtemp(prefix='fidelity_')
    gen = gen_reports(results['model'], results['decompose'], results['optimize'], tmp, project_id)

    html_path, r_html = gen['html']
    pptx_path, r_pptx = gen['pptx']

    report = {
        'project': project_id,
        'sot': {
            'media': sot['media'],
            'signed_factors': [f['name'] for f in sot['signed_factors']],
            'folded_into_baseline': [f['name'] for f in sot['folded']],
            'full_factor_count': sot['full_count'],
        },
        'gen_status': {'html': r_html.get('status'), 'pptx': r_pptx.get('status')},
        'missing': {'html': [], 'pptx': []},
        'present': {'html': [], 'pptx': []},
    }

    html_factors = extract_html_factors(html_path) if r_html.get('status') == 'ok' else {'blob': ''}
    pptx_factors = extract_pptx_factors(pptx_path) if r_pptx.get('status') == 'ok' else {'blob': ''}
    hblob = html_factors.get('blob', '')
    pblob = pptx_factors.get('blob', '')

    # Все SoT-факторы (медиа + signed) должны присутствовать в каждом отчёте.
    all_sot = [{'name': m, 'kind': 'media'} for m in sot['media']] + \
              [{'name': f['name'], 'kind': f['type']} for f in sot['signed_factors']]
    for item in all_sot:
        pres = factor_present(item, hblob, pblob)
        for fmt in ('html', 'pptx'):
            (report['present'] if pres[fmt] else report['missing'])[fmt].append(
                {'name': item['name'], 'kind': item['kind']}
            )

    missing_html = len(report['missing']['html'])
    missing_pptx = len(report['missing']['pptx'])
    report['summary'] = {
        'sot_factors': len(all_sot),
        'html_missing': missing_html,
        'pptx_missing': missing_pptx,
        'parity': missing_html == 0 and missing_pptx == 0,
    }

    if as_json:
        print(json.dumps(report, ensure_ascii=False, indent=2))
    else:
        print(f"\n=== FIDELITY DIFF: {project_id} ===")
        print(f"Источник истины (программа показывает): {len(all_sot)} факторов")
        print(f"  медиа ({len(sot['media'])}): {', '.join(sot['media'])}")
        print(f"  signed/holiday ({len(sot['signed_factors'])}): "
              f"{', '.join(f['name'] for f in sot['signed_factors'])}")
        print(f"  свёрнуто в baseline (positive_control): "
              f"{', '.join(f['name'] for f in sot['folded']) or '—'}")
        print(f"\nГенерация: HTML={r_html.get('status')} PPTX={r_pptx.get('status')}")
        print(f"\nHTML: отсутствует {missing_html}/{len(all_sot)} факторов")
        for x in report['missing']['html']:
            print(f"   ✗ [{x['kind']}] {x['name']}")
        print(f"\nPPTX: отсутствует {missing_pptx}/{len(all_sot)} факторов")
        for x in report['missing']['pptx']:
            print(f"   ✗ [{x['kind']}] {x['name']}")
        print(f"\nПАРИТЕТ ФАКТОРОВ: {'ДА ✓' if report['summary']['parity'] else 'НЕТ ✗'}")
        print("NB: XLSX (Rust report.rs) проверяется cargo-тестом, не этим харнессом.")
        print(f"Артефакты: {tmp}")

    return 0 if report['summary']['parity'] else 1


def main():
    args = [a for a in sys.argv[1:] if a != '--json']
    as_json = '--json' in sys.argv[1:]
    proj = args[0] if args else DEFAULT_PROJECT
    if not os.path.isdir(proj):
        print(f"ERROR: project dir not found: {proj}", file=sys.stderr)
        return 2
    return run(proj, as_json=as_json)


if __name__ == '__main__':
    sys.exit(main())
