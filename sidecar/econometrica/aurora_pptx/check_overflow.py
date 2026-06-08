# -*- coding: utf-8 -*-
"""Overflow-ассерт PPTX-вёрстки Econometrica: текст не вылезает за нижнюю границу
контента и текстовые боксы не накладываются.

Портирован из Smart Analytica (Dev/ROSST_AI_Media/tools/creative/check_overflow.py,
2026-06-07). Читает геометрию shape'ов из готового .pptx (python-pptx) и
ПЕРЕИЗМЕРЯЕТ реальную высоту текста той же функцией, что вёрстка (text_metrics,
PIL ImageFont). Заявленная height бокса не истина (autofit выключен — текст может
вылезать): высота считается из содержимого. Так «наложений нет» становится
проверяемым фактом для ВСЕХ слайдов, а не «посмотрели глазами».

Econometrica-геометрия (builder.py): слайд 13.333×7.5", footer-hairline y=7.05",
header-зона сверху 0.25". CONTENT_BOTTOM = 7.05 (низ контента = футер-hairline).

Запуск (CLI-гейт):
    python -m aurora_pptx.check_overflow <path.pptx> [--verbose]
Код возврата: 0 — чисто, 1 — есть наложения/overflow.
Также используется как pytest-гейт после build (tests/test_pptx_overflow.py).
"""
import sys
import os

from pptx import Presentation
from pptx.util import Inches

try:
    from . import text_metrics as TM
except ImportError:  # запуск как скрипт, не как модуль пакета
    sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
    import text_metrics as TM

# Econometrica: низ контента = footer-hairline 7.05" (builder.py:584).
CONTENT_BOTTOM = int(Inches(7.05))
FOOTER_ZONE_TOP = int(Inches(7.05))   # футер/номер страницы by-design ниже контента
TITLE_ZONE_TOP = int(Inches(0.25))    # выше — header/декор, не контент

TOL_OVERFLOW = int(Inches(0.04))      # допуск на округление при проверке низа
# Вертикальный допуск наложения. Econometrica намеренно стыкует стопки (eyebrow над
# заголовком, заголовок над подзаголовком) с микро-перекрытием боксов 0.06–0.08" —
# текст при этом НЕ сталкивается. 0.12" отсекает эти by-design стыки, оставляя
# реальные коллизии (обычно >0.15").
TOL_OVERLAP_X = int(Inches(0.06))
TOL_OVERLAP_Y = int(Inches(0.12))

# Декоративный огромный номер раздела («02») намеренно лежит за/рядом с текстом
# раздела. Детект: короткий текст крупным кеглем.
DECOR_MIN_SIZE_PT = 40.0
DECOR_MAX_LEN = 5

DEFAULT_SIZE_PT = 12.0
DEFAULT_LS = 1.2
DEFAULT_FONT = "Arial"


def _pt(length):
    if length is None:
        return None
    try:
        return length.pt
    except Exception:
        return float(length) / 12700.0


def _emu(length):
    if length is None:
        return 0
    try:
        return int(length)
    except Exception:
        return 0


def _para_line_spacing(para):
    ls = para.line_spacing
    if ls is None:
        return DEFAULT_LS
    if isinstance(ls, (int, float)):
        return float(ls)
    return None  # точная высота строки задана в EMU


def _para_exact_line_emu(para):
    ls = para.line_spacing
    if ls is not None and not isinstance(ls, (int, float)):
        return _emu(ls)
    return None


def _first_run_props(para):
    """size_pt, font_name, bold от ДОМИНИРУЮЩЕГО (самого длинного) run параграфа."""
    best, best_len = None, -1
    for run in para.runs:
        l = len(run.text or "")
        if l > best_len:
            best_len, best = l, run
    if best is None:
        return (DEFAULT_SIZE_PT, DEFAULT_FONT, False)
    return (_pt(best.font.size) or DEFAULT_SIZE_PT,
            best.font.name or DEFAULT_FONT,
            bool(best.font.bold))


def _textbox_real_height(shape):
    """Реальная высота текста в боксе (EMU) по содержимому, не по заявленной height."""
    tf = shape.text_frame
    width = _emu(shape.width) - _emu(tf.margin_left) - _emu(tf.margin_right)
    width = max(1, width)
    total = _emu(tf.margin_top) + _emu(tf.margin_bottom)
    for para in tf.paragraphs:
        text = para.text
        size_pt, font_name, bold = _first_run_props(para)
        n = TM.wrap_lines(text if text.strip() else " ", width, size_pt, font_name, bold)
        exact = _para_exact_line_emu(para)
        if exact is not None:
            lh = exact
        else:
            ls = _para_line_spacing(para) or DEFAULT_LS
            lh = TM.line_height_emu(size_pt, ls, font_name, bold)
        total += lh * n
        total += _emu(para.space_after)
    return total


def _table_real_height(shape):
    """Реальная высота ТАБЛИЦЫ (EMU) по содержимому (auto-grow строк под текст)."""
    tbl = shape.table
    col_w = [_emu(c.width) for c in tbl.columns]
    ncols = len(col_w)
    total = 0
    for r in range(len(tbl.rows)):
        row_h = int(Inches(0.30))
        for j in range(ncols):
            cell = tbl.cell(r, j)
            tf = cell.text_frame
            cw = max(1, col_w[j] - _emu(cell.margin_left) - _emu(cell.margin_right))
            cell_h = _emu(cell.margin_top) + _emu(cell.margin_bottom)
            for para in tf.paragraphs:
                text = para.text
                size_pt, font_name, bold = _first_run_props(para)
                n = TM.wrap_lines(text if text.strip() else " ", cw, size_pt, font_name, bold)
                ls = _para_line_spacing(para)
                ls = 1.15 if ls is None else max(ls, 1.0)
                cell_h += TM.line_height_emu(size_pt, ls, font_name, bold) * n
                cell_h += _emu(para.space_after)
            row_h = max(row_h, cell_h)
        total += row_h
    return total


def _short(text, n=46):
    t = " ".join(text.split())
    return (t[:n] + "...") if len(t) > n else t


def _rect_overlap(a, b):
    ax0, ay0, ax1, ay1 = a
    bx0, by0, bx1, by1 = b
    ox = min(ax1, bx1) - max(ax0, bx0)
    oy = min(ay1, by1) - max(ay0, by0)
    return ox > TOL_OVERLAP_X and oy > TOL_OVERLAP_Y, ox, oy


def _is_decorative_overlap_exempt(shape, top):
    """Намеренные дизайн-слои Econometrica, исключаемые из OVERLAP-проверки.

    НЕ исключаются из OVERFLOW (туда не попадают по зоне). Это:
      - футер-зона (wordmark «AURORA AI» + номер страницы, y≥7.05);
      - dot-leader строки TOC (текст из точек/пробелов);
      - огромный декор-номер раздела («02» кеглем ≥40 на divider-слайдах).
    Эти наложения by-design (текст визуально не сталкивается)."""
    if top >= FOOTER_ZONE_TOP:
        return True
    tf = getattr(shape, "text_frame", None)
    if tf is None:
        return False
    raw = tf.text or ""
    stripped = raw.strip()
    # dot-leader: только точки/пробелы/middot
    if stripped and all(c in ".· \t…" for c in stripped):
        return True
    # декоративный крупный короткий номер
    if len(stripped) <= DECOR_MAX_LEN:
        size_pt = max((_pt(r.font.size) or 0) for p in tf.paragraphs for r in p.runs) \
            if any(p.runs for p in tf.paragraphs) else 0
        if size_pt >= DECOR_MIN_SIZE_PT:
            return True
    return False


def check(pptx_path, verbose=False):
    """Возвращает (issues, total_slides). issues = [(slide_no, kind, detail)]."""
    prs = Presentation(pptx_path)
    total_slides = len(prs.slides._sldIdLst)
    issues = []

    for sidx, slide in enumerate(prs.slides, start=1):
        boxes = []
        for shape in slide.shapes:
            if getattr(shape, "has_table", False) and shape.has_table:
                top = _emu(shape.top)
                tbottom = top + _table_real_height(shape)
                if TITLE_ZONE_TOP <= top < FOOTER_ZONE_TOP and tbottom > CONTENT_BOTTOM + TOL_OVERFLOW:
                    over_in = (tbottom - CONTENT_BOTTOM) / 914400.0
                    issues.append((sidx, "OVERFLOW", f'таблица: низ +{over_in:.2f}" за границу'))
                continue
            if not shape.has_text_frame:
                continue
            if not shape.text_frame.text.strip():
                continue
            top = _emu(shape.top)
            left = _emu(shape.left)
            width = _emu(shape.width)
            real_h = _textbox_real_height(shape)
            real_bottom = top + real_h
            label = _short(shape.text_frame.text)
            # OVERLAP только для реальных контент-боксов (не декор-слои by-design).
            if not _is_decorative_overlap_exempt(shape, top):
                boxes.append(((left, top, left + width, real_bottom), label, top))

            if TITLE_ZONE_TOP <= top < FOOTER_ZONE_TOP:
                if real_bottom > CONTENT_BOTTOM + TOL_OVERFLOW:
                    over_in = (real_bottom - CONTENT_BOTTOM) / 914400.0
                    issues.append((sidx, "OVERFLOW",
                                   f'низ +{over_in:.2f}" за границу | "{label}"'))

        for i in range(len(boxes)):
            for j in range(i + 1, len(boxes)):
                ra, la, _ = boxes[i]
                rb, lb, _ = boxes[j]
                hit, ox, oy = _rect_overlap(ra, rb)
                if hit:
                    issues.append((sidx, "OVERLAP",
                                   f'{ox/914400.0:.2f}"x{oy/914400.0:.2f}" | '
                                   f'"{la}" <> "{lb}"'))

        if verbose:
            print(f"  слайд {sidx}: текстбоксов {len(boxes)}")

    return issues, total_slides


def report(pptx_path, verbose=False):
    """CLI-обёртка: печатает отчёт, возвращает код возврата 0/1."""
    issues, total_slides = check(pptx_path, verbose=verbose)
    print(f"\n=== OVERFLOW-АССЕРТ: {os.path.basename(pptx_path)} ===")
    print(f"слайдов: {total_slides} | content_bottom: {CONTENT_BOTTOM / 914400.0:.2f}\"")
    if not issues:
        print("РЕЗУЛЬТАТ: ЧИСТО — наложений и переполнений нет.")
        return 0
    by_slide = {}
    for sno, kind, detail in issues:
        by_slide.setdefault(sno, []).append((kind, detail))
    n_over = sum(1 for _, k, _ in issues if k == "OVERFLOW")
    n_lap = sum(1 for _, k, _ in issues if k == "OVERLAP")
    print(f"РЕЗУЛЬТАТ: {len(issues)} проблем(ы) на {len(by_slide)} слайдах "
          f"(OVERFLOW={n_over}, OVERLAP={n_lap}):\n")
    for sno in sorted(by_slide):
        print(f"слайд {sno}:")
        for kind, detail in by_slide[sno]:
            print(f"    [{kind}] {detail}")
    return 1


def main():
    args = [a for a in sys.argv[1:] if not a.startswith("--")]
    verbose = "--verbose" in sys.argv
    if not args:
        print("usage: python -m aurora_pptx.check_overflow <path.pptx> [--verbose]")
        return 2
    return report(args[0], verbose=verbose)


if __name__ == "__main__":
    raise SystemExit(main())
