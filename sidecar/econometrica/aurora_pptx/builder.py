"""
aurora_pptx.builder - production PPTX builder (13-slide tier-1 MMM report).

Ported from build_wireframe.py (finalized 2026-04-24 after iterative review with Антон).
Uses aurora_tokens (generated from Standards/tokens/tokens.json) for brand consistency.

Session 3 scope: initial port with default Kagocel data (pilot).
Session 4 (M4): параметризовать через data dict from Econometrica pipeline.

Public API:
    from econometrica.aurora_pptx.builder import AuroraPPTXBuilder
    builder = AuroraPPTXBuilder()
    prs = builder.build()
    prs.save(output_path)

Or via convenience function:
    from econometrica.aurora_pptx import build_pptx
    prs = build_pptx(data, lang='ru')

Brand principles applied (per Standards/CLIENT_READY_ANATOMY.md):
  - One gold accent per slide (strict discipline)
  - Running header + minimal footer (no duplication)
  - Sacred lime #CCFF00 under action titles
  - Pull quotes in Georgia italic
  - SCQAR structure в executive summary
  - Methodology includes explicit limitations
  - Closing slide inspirational, not technical
  - No em dash anywhere (Aurora rule)
  - Safe zones: top 0.25", bottom 0.20" (footer hairline 7.05)
  - Sources left-bottom quadrant, prижаты to footer hairline
  - Line_spacing=1.0 for headings
"""

from __future__ import annotations

import math
import zlib
from datetime import datetime

try:
    from econometrica.engines.narrative_adapter import (
        compute_report_id,
        derive_action_headline,
    )
except ImportError:
    from engines.narrative_adapter import (
        compute_report_id,
        derive_action_headline,
    )

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

try:
    from aurora_tokens import COLORS, TYPOGRAPHY, SIZING
except ImportError:
    # v2.1.0 (пилот 2026-05-17 audit C-1): fallback tokens когда build script
    # не был запущен (Standards/tokens/build.py отсутствует в распакованном
    # репо). Раньше выбрасывался RuntimeError с raw devops-сообщением,
    # которое попадало в красный banner на frontend Отчёта.
    # Минимальный валидный набор для PPTX генерации.
    COLORS = {
        'navy_deep': '#0A1628',
        'gold': '#C5A46D',
        'lime': '#CCFF00',
        'white': '#FFFFFF',
        'off_white': '#F7F7F7',
        'gray_text': '#475569',
        'gray_muted': '#94A3B8',
        'gray_line': '#E5E7EB',
        'red': '#DC2626',
        'green': '#22C55E',
        'blue': '#3B82F6',
    }
    TYPOGRAPHY = {
        'font_sans': 'Inter',
        'font_serif': 'Georgia',
        'font_mono': 'JetBrains Mono',
        'size_h1': 32,
        'size_h2': 24,
        'size_h3': 18,
        'size_body': 11,
        'size_caption': 9,
    }
    SIZING = {
        'slide_width_in': 13.333,
        'slide_height_in': 7.5,
        'margin_left_in': 0.5,
        'margin_right_in': 0.5,
        'margin_top_in': 0.4,
        'margin_bottom_in': 0.4,
    }


def hex_to_rgb(h):
    return RGBColor.from_string(h.lstrip("#").upper())


def _fmt_pct(v, fallback="-"):
    """N1 (Phase 0.1 fix-session 2026-04-25): conditional precision - never lies via rounding to 0%.
    Mirrors aurora_html/sections.py:_fmt_pct. See that file for behavior table.
    """
    if v is None:
        return fallback
    try:
        f = float(v)
    except (TypeError, ValueError):
        return fallback
    if f == 0:
        return "0%"
    av = abs(f)
    if av < 0.1:
        return "<0.1%" if f > 0 else ">-0.1%"
    if av < 1.0:
        return f"{f:.1f}%"
    return f"{round(f)}%"


# ─── KPI/mode-aware helpers (v1.3.2) ────────────────────────────────────────
#
# Helpers extracted to aurora_pptx.kpi_helpers - импортируются без aurora_tokens
# зависимости (тесты могут юзать без production token generation). Aliases в
# private form для legacy callers внутри builder.py.

from .kpi_helpers import (
    kpi_view as _kpi_view,
    fmt_metric as _fmt_metric_pptx,
    fmt_metric_with_ci_text as _fmt_metric_with_ci_text,
    weighted_summary_phrase as _weighted_summary_phrase_pptx,
    under_breakeven_phrase as _under_breakeven_phrase_pptx,
    table_metric_header as _table_metric_header_pptx,
    lift_phrase as _lift_phrase_pptx,
    hero_vs_leader_quote as _hero_vs_leader_quote_pptx,
    contrib_scale as _contrib_scale_pptx,
    fmt_contrib as _fmt_contrib_pptx,
)


class AuroraPPTXBuilder:
    """Production PPTX builder for Aurora AI Econometrica MMM reports.

    Currently uses default pilot data (Kagocel). Session 4 will parametrize
    via data dict passed from Econometrica pipeline.
    """

    def __init__(self, data=None, lang="ru"):
        if lang != "ru":
            raise NotImplementedError(
                f"lang={lang!r} not yet supported. RU-only pilot; EN scheduled for v1.0.12."
            )
        self.data = data or {}
        self.lang = lang
        # Tokens now come from generated aurora_tokens module (not JSON load)
        self.t = {"color": COLORS, "typography": TYPOGRAPHY, "sizing": SIZING}
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

        # v2.1.0 (пилот 2026-05-17): defensive .get() chain - aurora_tokens fallback
        # может не содержать nested ключи (brand.sig.lime, brand.bg.quiet). PPTX
        # export ранее падал KeyError 'sig' → frontend показывал red banner.
        c = self.t["color"]
        brand = c.get("brand", {}) if isinstance(c, dict) else {}
        deep = brand.get("deep", {}) if isinstance(brand, dict) else {}
        gold = brand.get("gold", {}) if isinstance(brand, dict) else {}
        bg = brand.get("bg", {}) if isinstance(brand, dict) else {}
        sig = brand.get("sig", {}) if isinstance(brand, dict) else {}
        self.deep_100 = hex_to_rgb(deep.get("100", "#0A1628"))
        self.deep_80  = hex_to_rgb(deep.get("80",  "#1E293B"))
        self.deep_60  = hex_to_rgb(deep.get("60",  "#475569"))
        self.deep_40  = hex_to_rgb(deep.get("40",  "#94A3B8"))
        self.deep_20  = hex_to_rgb(deep.get("20",  "#CBD5E1"))
        self.gold       = hex_to_rgb(gold.get("primary", "#C5A46D"))
        self.gold_muted = hex_to_rgb(gold.get("muted",   "#B8975D"))
        self.rule_color = hex_to_rgb(brand.get("rule", "#E5E7EB"))
        self.lime       = hex_to_rgb(sig.get("lime", "#CCFF00"))
        self.white      = hex_to_rgb(bg.get("white", "#FFFFFF"))
        self.bg_quiet   = hex_to_rgb(bg.get("quiet", "#F7F7F7"))

        # v2.1.0 (пилот 2026-05-17): defensive access - fallback TYPOGRAPHY мог не
        # содержать nested ключи. PPTX export ранее падал KeyError 'fontFamily'.
        ty = self.t.get("typography") or {}
        fam = (ty.get("fontFamily") or {}) if isinstance(ty, dict) else {}
        self.serif = fam.get("serif", "Georgia")
        self.sans  = fam.get("sans",  "Arial")
        self.mono  = fam.get("mono",  "Consolas")

        self.safe = 0.4
        self.w = 13.333
        self.h = 7.5

        # --- Meta (Session 4 M4: parametrized via data.meta, fallback = Kagocel pilot) ---
        meta = self.data.get("meta") or {}
        # B1-fix R-01/R-02 (2026-07-03): live-гейт. Реальный клиентский экспорт
        # (адаптер прислал channels) НЕ имеет права на wireframe-дефолты: пробел
        # в данных обязан рендериться честным «—»/«н/д», а не демо-значением
        # Kagocel-пилота (зонд: ESS 1247 при реальном tail 382.7; период
        # «W01 W13 2026» при данных 2023–2025). Preview-режим (build без data)
        # сохраняет прежние wireframe-дефолты — это легитимная демка.
        self.is_live = bool(self.data.get("channels"))
        self.client = meta.get("client", "Kagocel")
        self.project_id = meta.get("project_id", "KAGOCEL-Q1-2026")
        self.version = meta.get("version", "1.0.11")
        self.report_date = meta.get("report_date", "24 апреля 2026")
        self.period_label = meta.get(
            "period_label", None if self.is_live else "Q1 2026")
        self.forecast_period_label = meta.get(
            "forecast_period_label", None if self.is_live else "Q3-Q4 2026")
        self.data_window_label = meta.get(
            "data_window_label", None if self.is_live else "W01 W13 2026")
        # Stage C.6.3: TOC shrunk to 5 real sections with content (Option B
        # symmetric tier-1 structure). Previous 8-section layout promised
        # sections 4/6/7 (Модель / Оптимизация / Рекомендации) with no
        # dedicated slides; now each section-divider slide precedes real
        # content. Header "01 / 05" honest across the deck.
        self.section_names = meta.get("section_names", [
            "Главное",                  # 1: TOC + ataglance + keymsg + SCQAR (B4: было Executive summary)
            "Декомпозиция вкладов",     # 2: chart + table + timeline
            "Методология",              # 3: methodology
            "Данные и качество",        # 4: sources
            "Приложение и источники",   # 5: glossary + colophon
        ])
        self.total_sections = len(self.section_names)
        # B4-2 (2026-07-03): 12 слайдов — отдельные слайды-дивайдеры убраны
        # (4 полупустых из 16; стайлгайд §1), секции открывает _section_intro.
        self.total_slides = meta.get("total_slides", 12)
        # Physical page where each section begins (first content slide).
        self.toc_page_refs = meta.get("toc_page_refs", [3, 6, 9, 10, 11])
        # Physical-slide → (section_idx, section_label) map (12-slide layout).
        self.slide_to_section = meta.get("slide_to_section") or {
            2:  (1, "Главное"),                    # TOC
            3:  (1, "Главное"),                    # At a glance
            4:  (1, "Главное"),                    # Key message
            5:  (1, "Главное"),                    # SCQAR
            6:  (2, "Декомпозиция вкладов"),       # Action chart (mROAS) + плашка
            7:  (2, "Декомпозиция вкладов"),       # Action table (portfolio)
            8:  (2, "Декомпозиция вкладов"),       # Action timeline
            9:  (3, "Методология"),                # Methodology + плашка
            10: (4, "Данные и качество"),          # Sources + плашка
            11: (5, "Приложение и источники"),     # Glossary + плашка
            12: (5, "Приложение и источники"),     # Colophon
        }
        # Header center label (shown on every content slide)
        self.header_project_label = meta.get(
            "header_project_label",
            (f"{self.client.upper()} . ОТЧЁТ MMM . {self.period_label}"
             if self.period_label else f"{self.client.upper()} . ОТЧЁТ MMM"),
        )
        # Copyright footer on cover - dynamic year for future-proofing
        _year = datetime.now().year
        self.copyright_line = meta.get(
            "copyright_line",
            f"© {_year} Aurora AI. Конфиденциально",
        )
        # Source-note label template (client name substituted)
        self.sources_client_label = meta.get("sources_client_label", self.client)

        # --- Diagnostics (Phase 2 numbers - parametrized metric callouts) ---
        # B1-fix R-01: в live-режиме отсутствующая метрика = None (рендерится
        # «н/д» через _mstr), НЕ wireframe-число. Демо-значения — только preview.
        diag = self.data.get("diagnostics") or {}
        _d = (lambda key, preview: diag.get(key, None if self.is_live else preview))
        self.mqs_score = _d("mqs_score", 87)
        self.mqs_tier_label = _d("mqs_tier_label", "GOOD - готовность к production")
        self.r_squared = _d("r_squared", 0.872)
        self.mape_pct = _d("mape_pct", 8.3)
        self.r_hat_max = _d("r_hat_max", 1.008)
        self.ess_min = _d("ess_min", 1247)
        # B1-fix R-16: honesty-контур (вердикт надёжности + preflight) — доставлен
        # адаптером из движковой диагностики; рендерится на слайде «Данные и качество».
        self.honesty_verdict = diag.get("honesty_verdict")
        self.honesty_reasons = diag.get("honesty_reasons") or []
        self.preflight = diag.get("preflight") or {}
        # E2 (2026-07-03): калибровка lift-тестами — [CALIBRATED]-пометка у
        # канала в таблице + строки честности (включая расхождение модель-vs-
        # тест) на слайде «Данные и качество». Только live (из диагностики).
        self.calibration = diag.get("calibration") or {}
        self.calibrated_channels = {
            str(a.get("channel")) for a in (self.calibration.get("applied") or [])
        }
        # E4 (2026-07-03): сбывшиеся рекомендации (только сверенные kept/missed
        # из адаптера; live-only — wireframe-режима нет).
        self.promises_summary = (
            self.data.get("promises_summary") if self.is_live else None
        ) or None
        # B1-fix R-03/R-04: честное покрытие данных (n наблюдений / частота /
        # разрывы) от адаптера; None → строки рендерятся «—» или скрываются.
        self.data_coverage = self.data.get("data_coverage") or {}
        # E1+E3 (2026-07-03): вставные слайды честности в секции «ГЛАВНОЕ»
        # сразу после SCQAR — самые убедительные для продления подписки:
        #  · «Проверка на истории» (E1, models/backtest.json);
        #  · «Что изменилось с прошлого квартала» (E3, generation_compare.json).
        # Оба рисуются ТОЛЬКО с живыми данными — wireframe-режима НЕТ по
        # построению (урок B1: «замаскированная дефолтом честность»).
        # Дека 12 (+1 за каждый вставной); хвост и TOC сдвигаются формулой.
        _bt = self.data.get("backtest") or {}
        self.backtest = (
            _bt if (self.is_live and _bt.get("status") == "ok" and _bt.get("windows"))
            else None
        )
        _gc = self.data.get("generation_compare") or {}
        self.gen_compare = (
            _gc if (self.is_live and _gc.get("status") == "ok" and _gc.get("channels"))
            else None
        )
        _fc = self.data.get("forecast") or {}
        self.forecast = (
            _fc if (self.is_live and _fc.get("status") == "ok" and _fc.get("scenarios"))
            else None
        )
        self._page_shift = int(bool(self.backtest)) + int(bool(self.gen_compare)) + int(bool(self.forecast))
        if self._page_shift:
            _s = self._page_shift
            if "total_slides" not in meta:
                self.total_slides = 12 + _s
            if "toc_page_refs" not in meta:
                self.toc_page_refs = [3, 6 + _s, 9 + _s, 10 + _s, 11 + _s]
            if not meta.get("slide_to_section"):
                _map = {n: (1, "Главное") for n in range(2, 6 + _s)}
                for n in range(6 + _s, 9 + _s):
                    _map[n] = (2, "Декомпозиция вкладов")
                _map[9 + _s] = (3, "Методология")
                _map[10 + _s] = (4, "Данные и качество")
                _map[11 + _s] = (5, "Приложение и источники")
                _map[12 + _s] = (5, "Приложение и источники")
                self.slide_to_section = _map
        # INV-50 F-DELIVERABLE-1 (2026-06-07): data-thinness disclosure.
        # Default None → preview/wireframe (Кагоцел fallback) НЕ фабрикует
        # оговорку о переобучении; она появляется только когда backend реально
        # применил cap (thinness_cap=50/70). ratio_eff — эффективный (драйвит cap).
        self.thinness_cap = diag.get("thinness_cap")
        self.ratio_eff = diag.get("ratio")
        # v2.1.0 (Pilot C): engine detection. 'ols' для small-data fallback,
        # 'bayesian' (default) для production v1.2/v1.3 pickles. Determines
        # methodology labels (MCMC/NUTS vs closed-form/bootstrap).
        self.engine = diag.get("engine") or (
            "ols" if self.data.get("model_version") == "1.0-ols" else "bayesian"
        )
        self.is_ols = self.engine == "ols"

        # --- Channels + narrative facts (Session C - Path C parametrization) ---
        # If adapter supplied channels + facts: real client data drives slide
        # templates. Otherwise builder keeps its Kagocel pilot narrative
        # (wireframe / preview mode) - single source-of-truth guard lives
        # with `self.facts is None` checks in each slide method.
        self.channels = self.data.get("channels") or []
        self.facts = self.data.get("narrative_facts")
        self.time_series = self.data.get("time_series") or None
        # Аудит #12 (2026-06-07): канонический набор серий timeline-декомпозиции
        # (baseline_reduced + media + вынесенные signed/holiday) — тот же, что в
        # программе и остальных отчётах.
        self.decomposition_series = self.data.get("decomposition_series") or None

        # v1.3.2: KPI metadata (kpi_kind, derived_mode, labels) per ADR-016.
        # narrative_adapter populates data['kpi']; legacy callers (v1.2) get
        # defaults через _kpi_view fallback. self.kpi.is_legacy guards keep
        # все hardcoded ROI/mROAS strings backward-compat.
        self.kpi = _kpi_view(self.data)

        # --- Report ID (Stage B.4, post-audit unified 2026-04-25) ---
        # Deterministic trace hash shared with aurora_html. Uses the raw
        # diagnostics dict from data (NOT the Kagocel fallback attrs like
        # self.mqs_score=87), so partial-data runs don't leak pilot values
        # into the hash. Same pipeline output → same ID in HTML and PPTX.
        raw_diagnostics = self.data.get("diagnostics") or {}
        self.report_id = self.data.get("report_id") or compute_report_id(
            self.client, self.project_id, self.channels, raw_diagnostics,
        )

    # ---------- Primitives ----------

    # B4-3: семантическая привязка секция → пиктограмма (assets/icons/*.png,
    # генератор tools/gen_report_icons.py; линейный стиль deep/gold).
    _SECTION_ICONS = {2: "growth", 3: "lens", 4: "db", 5: "book"}

    def _section_intro(self, slide, num, title, takeaway=None):
        """B4-2 (2026-07-03, стайлгайд §1): секцию открывает компактная плашка
        на первом содержательном слайде — вместо отдельного слайда-дивайдера
        (4 полупустых слайда из 16 = 25% деки; Knaflic: каждый элемент экрана
        тратит внимание читателя). Строка: «[пиктограмма] 0X · НАЗВАНИЕ — вывод»."""
        from pathlib import Path
        text_x = self.safe
        icon_name = self._SECTION_ICONS.get(num)
        if icon_name:
            icon_path = Path(__file__).parent / "assets" / "icons" / f"{icon_name}.png"
            if icon_path.exists():
                # graceful: файла нет (дистрибутив без ассетов) — плашка без иконки
                slide.shapes.add_picture(
                    str(icon_path), Inches(self.safe), Inches(0.545),
                    height=Inches(0.20),
                )
                text_x = self.safe + 0.28
        runs = [
            (f"{num:02d} · {title.upper()}",
             {"font": self.sans, "size": 9, "bold": True, "color": self.gold}),
        ]
        if takeaway:
            _t = takeaway if len(takeaway) <= 120 else takeaway[:117] + "…"
            runs.append((f"  –  {_t}",  # П8-1: en-dash вместо em-dash в section_intro
                         {"font": self.sans, "size": 9, "italic": True, "color": self.deep_60}))
        self._rich(slide, text_x, 0.56, self.w - text_x - self.safe, 0.2, runs=runs)

    def _mstr(self, value, fmt="{}", na="н/д"):
        """B1-fix R-01: None-безопасный рендер метрики. Отсутствующее значение
        в live-режиме — честное «н/д», не wireframe-число."""
        if value is None:
            return na
        try:
            return fmt.format(value)
        except (ValueError, TypeError):
            return na

    def _period_unit_label(self) -> str:
        """Русская единица периода для заголовков timeline («по неделям» /
        «по месяцам» / «по периодам»). Реиспользует detect_granularity из
        aurora_html — П2-паттерн: не врать «недельной» на месячных данных."""
        ts = (self.time_series or {}) if isinstance(self.time_series, dict) else {}
        dates = ts.get("dates") or ts.get("weeks") or []
        if len(dates) < 2:
            return "по периодам"
        try:
            from utils.forecast_validation import detect_granularity
            g = detect_granularity(dates)
            if g.get("confidence", 0.0) < 0.5:
                return "по периодам"
            return {
                "D": "по дням", "W": "по неделям", "M": "по месяцам",
                "Q": "по кварталам", "Y": "по годам",
            }.get(g.get("granularity"), "по периодам")
        except Exception:
            return "по периодам"

    def _blank(self):
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def _text(self, slide, x, y, w, h, text, *,
              font=None, size=12, bold=False, italic=False,
              color=None, align=None, line_spacing=None):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0); tf.margin_right = Inches(0)
        tf.margin_top = Inches(0); tf.margin_bottom = Inches(0)
        p = tf.paragraphs[0]
        if align is not None:
            p.alignment = align
        if line_spacing is not None:
            p.line_spacing = line_spacing
        r = p.add_run()
        r.text = text
        f = r.font
        f.name = font or self.sans
        f.size = Pt(size)
        f.bold = bold
        f.italic = italic
        if color:
            f.color.rgb = color
        return tb

    def _rich(self, slide, x, y, w, h, runs, *, line_spacing=1.3, align=None):
        """Multi-run paragraph. runs = list of (text, opts_dict)."""
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0); tf.margin_right = Inches(0)
        tf.margin_top = Inches(0); tf.margin_bottom = Inches(0)
        p = tf.paragraphs[0]
        p.line_spacing = line_spacing
        if align is not None:
            p.alignment = align
        for text, opts in runs:
            r = p.add_run()
            r.text = text
            f = r.font
            f.name = opts.get("font", self.sans)
            f.size = Pt(opts.get("size", 12))
            f.bold = opts.get("bold", False)
            f.italic = opts.get("italic", False)
            f.color.rgb = opts.get("color", self.deep_100)
        return tb

    def _paragraphs(self, slide, x, y, w, h, lines, *, line_spacing=1.35):
        """Paragraph list. lines = list of (text, opts)."""
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0); tf.margin_right = Inches(0)
        tf.margin_top = Inches(0); tf.margin_bottom = Inches(0)
        for i, item in enumerate(lines):
            text, opts = item if isinstance(item, tuple) else (item, {})
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.line_spacing = opts.get("line_spacing", line_spacing)
            if "align" in opts:
                p.alignment = opts["align"]
            r = p.add_run()
            r.text = text
            f = r.font
            f.name = opts.get("font", self.sans)
            f.size = Pt(opts.get("size", 11))
            f.bold = opts.get("bold", False)
            f.italic = opts.get("italic", False)
            f.color.rgb = opts.get("color", self.deep_100)
        return tb

    def _hairline(self, slide, x, y, width, *, weight=0.5, color=None):
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(x), Inches(y), Inches(x + width), Inches(y)
        )
        line.line.color.rgb = color or self.rule_color
        line.line.width = Pt(weight)
        return line

    def _vline(self, slide, x, y, height, *, weight=0.5, color=None):
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(x), Inches(y), Inches(x), Inches(y + height)
        )
        line.line.color.rgb = color or self.rule_color
        line.line.width = Pt(weight)
        return line

    def _lime_under(self, slide, x, y, width):
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(x), Inches(y), Inches(x + width), Inches(y)
        )
        line.line.color.rgb = self.lime
        line.line.width = Pt(2)

    def _vbar(self, slide, x, y, height, *, weight=2, color=None):
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(x), Inches(y), Inches(x), Inches(y + height)
        )
        line.line.color.rgb = color or self.gold
        line.line.width = Pt(weight)

    def _rect(self, slide, x, y, w, h, *, fill=None, line=False):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y), Inches(w), Inches(h)
        )
        if fill is not None:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill
        else:
            shape.fill.background()
        if not line:
            shape.line.fill.background()
        return shape

    def _arrow(self, slide, x1, y1, x2, y2, *, weight=0.75, color=None):
        """Thin annotation arrow."""
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(x1), Inches(y1), Inches(x2), Inches(y2)
        )
        line.line.color.rgb = color or self.deep_80
        line.line.width = Pt(weight)
        # Arrowhead on the end
        try:
            from pptx.oxml.ns import qn
            from lxml import etree
            ln = line.line._get_or_add_ln()
            tailEnd = etree.SubElement(ln, qn('a:tailEnd'))
            tailEnd.set('type', 'triangle')
            tailEnd.set('w', 'sm')
            tailEnd.set('len', 'sm')
        except Exception:
            pass
        return line

    # ---------- Brand mark (sigil PNG variants) ----------

    def _brand_mark_compact_path(self):
        """Compact «cuted» variant: sigil + AURORA AI текст в одной композиции.
        Для cover slide (top-left small box). 1024×768 4:3 aspect ratio.
        """
        from pathlib import Path
        p = Path(__file__).parent / 'templates' / 'brand_mark_compact.png'
        return p if p.exists() else None

    def _brand_mark_full_path(self):
        """Full gold-accent variant: large sigil с centered AURORA AI wordmark.
        Для back cover slide (centered, hero size). 1024×1024 square.
        """
        from pathlib import Path
        p = Path(__file__).parent / 'templates' / 'brand_mark_full.png'
        return p if p.exists() else None

    # ---------- Wordmark (strict typographic logo) ----------

    def _wordmark(self, slide, x, y, *, size=12, color=None):
        """'AURORA | AI' as strict typographic lockup with gold divider."""
        col = color or self.deep_100
        # Scale factor - wordmark width grows with size
        scale = size / 12.0
        w_aurora = 0.95 * scale
        w_gap = 0.08 * scale
        w_ai = 0.28 * scale
        height = 0.22 * scale

        # AURORA in serif small caps
        tb = slide.shapes.add_textbox(
            Inches(x), Inches(y), Inches(w_aurora + 0.1), Inches(height)
        )
        tf = tb.text_frame
        tf.margin_left = Inches(0); tf.margin_right = Inches(0)
        tf.margin_top = Inches(0); tf.margin_bottom = Inches(0)
        tf.word_wrap = False
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = "AURORA"
        r.font.name = self.serif
        r.font.size = Pt(size)
        r.font.bold = False
        r.font.color.rgb = col

        # Gold vertical bar divider
        div_x = x + w_aurora + w_gap
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(div_x), Inches(y + height * 0.15),
            Inches(div_x), Inches(y + height * 0.85)
        )
        line.line.color.rgb = self.gold
        line.line.width = Pt(max(0.75, size / 12))

        # AI in bold
        ai_x = div_x + w_gap
        tb2 = slide.shapes.add_textbox(
            Inches(ai_x), Inches(y), Inches(w_ai + 0.1), Inches(height)
        )
        tf2 = tb2.text_frame
        tf2.margin_left = Inches(0); tf2.margin_right = Inches(0)
        tf2.margin_top = Inches(0); tf2.margin_bottom = Inches(0)
        tf2.word_wrap = False
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = "AI"
        r2.font.name = self.serif
        r2.font.size = Pt(size)
        r2.font.bold = True
        r2.font.color.rgb = col

    # ---------- Running header (top of content slides) ----------

    def _header(self, slide, *, slide_num=None, section_idx=None, section_label=None,
                include_confidential=True, include_project=False):
        """Slim running header: closer to top edge (y=0.25) to reclaim content space.

        Stage B.2: center "project identifier" textbox removed - it was
        leaking internal project slug + MMM REPORT + period into every
        slide (35 chars of visual noise). Tier-1 minimalism: left section
        label + right CONFIDENTIAL only. `include_project` kept as a
        parameter (default False) for explicit opt-in if ever needed.

        Stage C.6.2: prefer `slide_num` which resolves (section_idx,
        section_label) via self.slide_to_section. Explicit section_idx /
        section_label still accepted for backward compat and edge overrides.
        """
        if slide_num is not None and slide_num in self.slide_to_section:
            mapped_idx, mapped_label = self.slide_to_section[slide_num]
            if section_idx is None:
                section_idx = mapped_idx
            if section_label is None:
                section_label = mapped_label
        if section_idx is None or section_label is None:
            raise ValueError(
                "_header requires either slide_num (mapped in slide_to_section) "
                "or explicit section_idx + section_label"
            )
        y = 0.25  # compact safe_top zone
        # Left: section tag like "03 / 08 . Methodology"
        self._text(
            slide, self.safe, y, 5.0, 0.2,
            f"{section_idx:02d} / {self.total_sections:02d} . {section_label.upper()}",
            font=self.sans, size=8, bold=True, color=self.deep_100, align=PP_ALIGN.LEFT,
        )
        # Center: project identifier - DISABLED by default (Stage B.2)
        if include_project:
            self._text(
                slide, 4.5, y, 4.333, 0.2,
                self.header_project_label,
                font=self.sans, size=8, color=self.deep_60,
                align=PP_ALIGN.CENTER,
            )
        # Right: classification
        if include_confidential:
            self._text(
                slide, self.w - self.safe - 4.0, y, 4.0, 0.2,
                "КОНФИДЕНЦИАЛЬНО",
                font=self.sans, size=8, italic=True, color=self.deep_60,
                align=PP_ALIGN.RIGHT,
            )
        # Hairline under header (at y=0.50, was 0.70)
        self._hairline(slide, self.safe, y + 0.25, self.w - 2 * self.safe, weight=0.25)

    # ---------- Footer (minimal: mini wordmark + page num) ----------

    def _footer(self, slide, page_num, *, show_page=True, show_wordmark=True):
        """Footer: hairline at y=7.05 (fixed safe_bottom boundary).
        Elements (wordmark + page num) visually centered between hairline and slide bottom (7.50).
        Midpoint = 7.275 → element_y = 7.20.

        `show_wordmark=False` suppresses the left-rail wordmark (used on the
        colophon which renders a hero wordmark in its body - avoids the
        tiny-wordmark-plus-hero-wordmark visual repetition).
        """
        self._hairline(slide, self.safe, 7.05, self.w - 2 * self.safe, weight=0.25)
        element_y = 7.20  # visual center of (7.05, 7.50) zone
        # Left: mini wordmark (suppressed on colophon)
        if show_wordmark:
            self._wordmark(slide, self.safe, element_y, size=8, color=self.deep_60)
        # Center: page num в формате "N/Total"
        if show_page:
            self._text(
                slide, 0, element_y, self.w, 0.18,
                f"{page_num}/{self.total_slides}",
                font=self.sans, size=9, color=self.deep_60,
                align=PP_ALIGN.CENTER,
            )

    # ---------- Section progress bar (little dots with current highlighted) ----------

    def _section_progress(self, slide, x, y, *, current):
        """8 numbered pills showing section progression, current is gold."""
        pill_w = 0.35
        pill_gap = 0.08
        total = self.total_sections
        for i in range(total):
            px = x + i * (pill_w + pill_gap)
            is_current = (i + 1) == current
            # Number
            self._text(
                slide, px, y, pill_w, 0.18,
                f"{i+1:02d}",
                font=self.sans, size=7, bold=is_current,
                color=self.gold if is_current else self.deep_40,
                align=PP_ALIGN.CENTER,
            )
            # Underline hairline
            self._hairline(
                slide, px, y + 0.2, pill_w,
                weight=1.25 if is_current else 0.25,
                color=self.gold if is_current else self.rule_color,
            )

    # ---------- Action title block (full pattern with lime) ----------

    def _action_title(self, slide, text, *, show_lime=True, y=1.0, height=0.95,
                      size=22):
        """Stage C.4: size made overridable so long titles (4+ lines) can
        shrink to 18-20pt to fit. Default 22pt preserved for short titles.
        """
        left = self.safe
        width = self.w - 2 * self.safe
        self._text(
            slide, left, y, width, height,
            text, font=self.serif, size=size, bold=True,
            color=self.deep_100, line_spacing=1.0,
        )
        if show_lime:
            self._lime_under(slide, left, y + height, width)

    # ---------- Pull quote (serif italic 20pt with gold vertical bar) ----------

    def _pull_quote(self, slide, x, y, w, h, text, *, size=18, with_bar=True):
        if with_bar:
            self._vbar(slide, x - 0.1, y, h, weight=2.5, color=self.gold)
        self._text(
            slide, x, y, w, h, text,
            font=self.serif, size=size, italic=True,
            color=self.deep_80, line_spacing=1.35,
        )

    # ---------- Big number (hero callout) ----------

    def _big_number(self, slide, x, y, number, *, label, support=None, size=96):
        self._text(
            slide, x, y - 0.3, 5.0, 0.25, label.upper(),
            font=self.sans, size=8, bold=True, color=self.gold,
        )
        self._hairline(slide, x, y - 0.05, 1.2, weight=0.75, color=self.gold)
        self._text(
            slide, x, y, 5.0, size / 50,
            number, font=self.serif, size=size, color=self.deep_100,
        )
        if support:
            self._text(
                slide, x, y + size / 55, 5.0, 0.3, support,
                font=self.sans, size=10, italic=True, color=self.deep_60,
            )

    # ---------- Source footnote (hairline + italic source) ----------

    def _source(self, slide, y, *, text, width=None):
        """Source footnote. Left-bottom quadrant positioning:
        - No own hairline (footer hairline is single bottom divider)
        - Max width 8.3 inch (left 2/3 of slide) - right side reserved
        - Font size 7 italic deep_60 - minimal weight
        - y=6.67 default: max close to footer hairline 6.85 (0.03 gap above)"""
        width = width or 8.3
        self._text(
            slide, self.safe, y, width, 0.2, text,
            font=self.sans, size=7, italic=True, color=self.deep_60,
        )

    # ---------- Category tag (small gold label above content) ----------

    def _category(self, slide, x, y, text, *, w=6.0, color=None):
        self._text(
            slide, x, y, w, 0.2, text.upper(),
            font=self.sans, size=8, bold=True,
            color=color or self.gold,
        )

    # ----------------------------------------------------------------
    # Narrative helpers (Session C - Path C parametrization)
    # ----------------------------------------------------------------

    def _build_at_a_glance_findings(self):
        """Five findings for s02, slot-filled from self.channels + self.facts.
        Returns list of (num, finding, support) tuples. Safe when some fields
        are None - missing data is elided rather than showing '-'.
        """
        f = self.facts or {}
        leader = f.get("leader_channel") or "Лидер"
        hero = f.get("hero_channel") or leader
        leader_contrib_pct = f.get("leader_share_contrib_pct")
        leader_spend_pct = f.get("leader_share_spend_pct")
        weighted_roi = f.get("weighted_roi")
        reallocation_mln = f.get("reallocation_mln")
        expected_lift_pct = f.get("expected_lift_pct")
        honest = bool(f.get("honest_narrative"))
        media_pct = f.get("media_contribution_pct")
        baseline_pct = f.get("baseline_pct")

        by_mroas = sorted(self.channels, key=lambda c: float(c.get("mroas") or 0), reverse=True)
        hero_ch = by_mroas[0] if by_mroas else {}
        hero_mroas = float(hero_ch.get("mroas") or 0)
        hero_spend_pct = None
        if self.facts and self.facts.get("total_budget_mln"):
            hero_spend_mln = float(hero_ch.get("spend") or 0) / 1_000_000
            total_mln = self.facts.get("total_budget_mln") or 0
            if total_mln > 0:
                hero_spend_pct = hero_spend_mln / total_mln * 100

        # Finding 1 - leader contribution vs budget share
        # Honest mode: baseline-dominated → disclose actual media share, not
        # leader's share-of-media (misleading "X% sales" phrasing).
        if honest and media_pct is not None and baseline_pct is not None:
            f1 = f"Медиа-вклад {_fmt_pct(media_pct)}, базовый спрос {_fmt_pct(baseline_pct)} – модель объясняет продажи через organic"  # П8-2
            s1 = f"{leader} – лидер среди медиа ({_fmt_pct(leader_contrib_pct)} медиа-вклада)" if leader_contrib_pct is not None else f"{leader} – лидер среди медиа"  # П8-2
        else:
            if leader_contrib_pct is not None and leader_spend_pct is not None:
                # B1-fix R-12 (2026-07-03): leader_share_contrib_pct — доля в
                # МЕДИА-вкладе, не в продажах (Kagocel: «32% продаж» реально
                # 32% медиа-вклада = 12% продаж при медиа 38%). Квалификатор
                # обязателен — иначе клиент завышает роль канала втрое.
                f1 = f"{leader} - {_fmt_pct(leader_contrib_pct)} медиа-вклада при {_fmt_pct(leader_spend_pct)} бюджета"
            else:
                f1 = f"{leader} - максимальный медиа-вклад в продажи"
            # v1.3.2: KPI-aware portfolio metric (ROI×/CPU/доля).
            if weighted_roi is not None:
                if self.kpi["is_legacy"]:
                    s1 = f"ROI {weighted_roi:.1f}× средневзвешенный по каналам"
                else:
                    s1 = f"{_weighted_summary_phrase_pptx(weighted_roi, self.kpi)} средневзвешенный по каналам"
            else:
                s1 = "Основной драйвер портфеля"

        # Finding 2 - hero channel by mROAS
        # v1.3.2: KPI-aware metric short label + breakeven phrase.
        hero_metric_short = self.kpi["metric_short"]
        hero_metric_fmt = (
            f"{hero_mroas:.1f}×" if self.kpi["is_legacy"]
            else _fmt_metric_pptx(hero_mroas, self.kpi)
        )
        # INV-50 (2026-06-03 synthetic-truth аудит): sub-breakeven hero НИКОГДА не
        # «самый эффективный» (противоречит вердикту «убыточный»). Гейт расцеплён
        # от honest_narrative (media<10%). effectiveness-mode исключён (метрика —
        # доля, breakeven неприменим). Зеркалит decomposer + HTML-hero.
        if hero_mroas < 1.0 and self.kpi["mode"] != "effectiveness":
            f2 = f"{hero} - лучший среди медиа, но под breakeven ({hero_metric_short} {hero_metric_fmt})"
            s2 = f"{_under_breakeven_phrase_pptx(self.kpi)} означает что канал тратит больше чем приносит"
        elif hero_mroas > 0:
            f2 = f"{hero} - самый эффективный канал, {hero_metric_short} {hero_metric_fmt}"
            s2 = f"Текущий бюджет на нём {_fmt_pct(hero_spend_pct)}" if hero_spend_pct is not None else "Потенциал для перераспределения бюджета"
        else:
            f2 = f"{hero} - наиболее эффективный канал по {hero_metric_short}"
            s2 = "Потенциал для перераспределения бюджета"

        # Finding 3 - reallocation / honest disclosure
        def _fmt_mln(v):
            if v is None:
                return "0"
            return f"{v:.1f}" if v < 10 else f"{v:.0f}"

        # v1.3.2 audit fix (M1): для effectiveness mode skip - shares always
        # sum to 100%, «under breakeven» semantically meaningless.
        all_below_breakeven = (
            bool(self.channels)
            and self.kpi["mode"] != "effectiveness"
            and all((float(c.get("mroas") or c.get("roi") or 0) < 1.0) for c in self.channels)
        )
        if honest and all_below_breakeven:
            f3 = "Рекомендация: все каналы под breakeven - сократить медиа или диагностика данных"
            if self.kpi["is_legacy"]:
                s3 = "При weighted ROI < 1× оптимизация перераспределением не вернёт прибыльность"
            else:
                s3 = (
                    f"Когда у всех каналов {_under_breakeven_phrase_pptx(self.kpi)} - "
                    "оптимизация перераспределением не вернёт прибыльность"
                )
        elif reallocation_mln and reallocation_mln >= 0.5 and (
            f.get("cut_source_channel") and f.get("scale_destination_channel")
        ):
            # L15 (math-fix v1.4 Section C): action-driven reallocation subjects
            cut_source = f.get("cut_source_channel")
            scale_dest = f.get("scale_destination_channel")
            f3 = f"Рекомендация: перераспределить {_fmt_mln(reallocation_mln)} млн из {cut_source} в {scale_dest}"
            s3 = _lift_phrase_pptx(expected_lift_pct, self.kpi) if expected_lift_pct is not None else "Ожидаемый эффект - положительный"
        elif reallocation_mln and reallocation_mln >= 0.5 and hero != leader:
            # Legacy fallback when cut_source/scale_destination not yet populated
            f3 = f"Рекомендация: перераспределить {_fmt_mln(reallocation_mln)} млн из {leader} в {hero}"
            s3 = _lift_phrase_pptx(expected_lift_pct, self.kpi) if expected_lift_pct is not None else "Ожидаемый эффект - положительный"
        else:
            f3 = "Рекомендация: сохранить текущую аллокацию по лидеру портфеля"
            # B1-fix R-09/R-13: «+0.0 пп» — пустое обещание; при незначимом lift
            # честная формулировка вместо нулевого числа.
            if expected_lift_pct is not None and expected_lift_pct >= 0.5:
                s3 = _lift_phrase_pptx(expected_lift_pct, self.kpi)
            else:
                s3 = "Перераспределение не даёт ощутимого прироста - портфель у оптимума"

        # Finding 4 - verdict distribution (how portfolio looks)
        # Stage C.3: idiomatic Russian plural forms (no lazy "канал(ов)" hack).
        verdicts = [c.get("verdict") for c in self.channels]
        scale_n = sum(1 for v in verdicts if v == "Scale")
        cut_n = sum(1 for v in verdicts if v in ("Cut", "Reduce"))
        uncertain_n = sum(1 for v in verdicts if v == "Uncertain")
        def _ru_channels(n: int) -> str:
            if n == 0:
                return "ни одного канала"
            if n % 10 == 1 and n % 100 != 11:
                return f"{n} канал"
            if n % 10 in (2, 3, 4) and n % 100 not in (12, 13, 14):
                return f"{n} канала"
            return f"{n} каналов"
        # B1-fix R-14 (2026-07-03): при неопределённых вердиктах (широкие
        # posterior-интервалы) «чёткая рекомендация по каждому» — ложь.
        # Kagocel-зонд: все 4 канала Uncertain + текст обещал чёткость.
        if uncertain_n > 0 and scale_n == 0 and cut_n == 0:
            f4 = f"Портфель: вердикты {_ru_channels(uncertain_n)} неопределённые - интервалы эффективности широки"
            s4 = "Для уверенных действий нужно больше данных или длиннее история"
        elif uncertain_n > 0:
            f4 = f"Портфель: {_ru_channels(scale_n)} к росту, {_ru_channels(cut_n)} к сокращению"
            s4 = f"Из {len(self.channels)} каналов у {_ru_channels(uncertain_n)} вердикт неопределённый"
        else:
            f4 = f"Портфель: {_ru_channels(scale_n)} к росту, {_ru_channels(cut_n)} к сокращению"
            s4 = f"Из {len(self.channels)} активных каналов - чёткая рекомендация по каждому"

        # Finding 5 - MQS quality signal.
        # 2026-07-25: нет числа - нет подписи. self.mqs_score=None (метрика не
        # рассчитана для этого прогона) раньше тихо превращался в mqs=0.0 и
        # рендерился как «MQS 0/100 - требует доработки» - правдоподобный
        # приговор модели, хотя её просто не оценивали. Честное отсутствие
        # оценки вместо фиктивного нуля - зеркалит aurora_html render_at_a_glance.
        try:
            mqs = float(self.mqs_score) if self.mqs_score is not None else None
        except (TypeError, ValueError):
            mqs = None
        if mqs is None:
            f5 = "Оценка качества модели (MQS) не выполнялась для этого расчёта"
            s5 = "Итоговый балл недоступен - диагностические метрики см. на слайде «Данные и качество»"
        elif mqs >= 80:
            f5 = f"Качество модели: MQS {mqs:.0f}/100 - готовность к использованию"
            s5 = "Можно опираться на рекомендации в планировании"
        elif mqs >= 60:
            f5 = f"Качество модели: MQS {mqs:.0f}/100 - приемлемо"
            s5 = "Рекомендации валидны с учётом диагностических метрик"
        else:
            f5 = f"Качество модели: MQS {mqs:.0f}/100 - требует доработки"
            s5 = "Следует расширить период данных перед финальными решениями"

        return [
            ("01", f1, s1),
            ("02", f2, s2),
            ("03", f3, s3),
            ("04", f4, s4),
            ("05", f5, s5),
        ]

    def _build_action_table_rows(self, channels):
        """Format merged channels list into the (name, budget, contrib, mROAS,
        share_pct, verdict, footnote) tuples consumed by s07 action table.
        Auto-generates footnote superscripts only for Reduce/Cut verdicts,
        keyed by order (max 3 footnotes to fit bottom-block layout).
        """
        total_contrib = sum(float(c.get("contribution") or 0) for c in channels) or 1.0
        # Assign footnote numbers to the first 3 flagged channels (Reduce/Cut).
        # Filter within channels[:10] (same slice the row loop uses) so the
        # bottom-block footnote text and the row superscript always pair up.
        # (Pre-fix bug: flagged could include channel 15+ which has no rendered
        # row, leaving the bottom-block footnote orphaned.)
        visible = channels[:10]
        # Масштаб+единица столбца «Вклад» — согласованно с шапкой в s07 (fix
        # 2026-07-13, INV-50). Та же логика по тем же данным = детерминированно.
        contrib_scale, _ = _contrib_scale_pptx(
            self.kpi, [c.get("contribution") for c in visible]
        )
        flagged = [c for c in visible if c.get("verdict") in ("Reduce", "Cut")][:3]
        fn_by_name = {c["name"]: str(i + 1) for i, c in enumerate(flagged) if c.get("name")}

        rows = []
        for c in channels[:10]:
            name = c.get("name") or "-"
            spend = float(c.get("spend") or 0)
            contrib = float(c.get("contribution") or 0)
            mroas = c.get("mroas")
            verdict = c.get("verdict", "Watch")

            budget_str = f"{spend / 1_000_000:.0f}" if spend else "0"
            contrib_str = _fmt_contrib_pptx(contrib, contrib_scale) if contrib else "0"
            # v1.3.2 audit fix (B4): KPI-aware metric cell - backend convention
            # c.mroas = mathematical units/₽; для count display invert к CPU.
            # legacy = '1.5'× / count = '80' (CPU; unit в header «₽/ед.») /
            # effectiveness = '25' (percentage; unit «%»).
            if self.kpi["is_legacy"]:
                roi_str = f"{float(mroas):.1f}" if mroas else "-"
            elif self.kpi["mode"] == "effectiveness":
                if mroas:
                    val = float(mroas)
                    roi_str = f"{val * 100:.1f}" if abs(val) <= 1.0 else f"{val:.0f}"
                else:
                    roi_str = "-"
            elif self.kpi["kpi_kind"] == "count":
                # Invert units/₽ → CPU ₽/ед. (B4 audit fix).
                if mroas and float(mroas) > 0:
                    roi_str = f"{1.0 / float(mroas):.0f}"
                else:
                    roi_str = "-"
            else:
                roi_str = f"{float(mroas):.1f}" if mroas else "-"
            share_pct = int(round(contrib / total_contrib * 100))
            share_str = f"{share_pct}" if share_pct > 0 else "0"
            footnote = fn_by_name.get(name, "")

            # Phase 1.9: posterior 90% HDI bracket on mROAS - None when v1.0/v1.1 pickle.
            # v1.3.2 audit fix (B4): для count mode CI inverts (units/₽ → CPU)
            # и swaps ordering: lo_mroas → hi_cpu, hi_mroas → lo_cpu.
            ci_low = c.get("mroas_ci_low")
            ci_high = c.get("mroas_ci_high")
            ci_str = ""
            if ci_low is not None and ci_high is not None:
                if self.kpi["is_legacy"]:
                    ci_str = f"[{float(ci_low):.1f} - {float(ci_high):.1f}]"
                elif self.kpi["mode"] == "effectiveness":
                    lo = float(ci_low) * (100 if abs(float(ci_low)) <= 1.0 else 1)
                    hi = float(ci_high) * (100 if abs(float(ci_high)) <= 1.0 else 1)
                    ci_str = f"[{lo:.1f} - {hi:.1f}]"
                elif self.kpi["kpi_kind"] == "count":
                    try:
                        lo_raw = float(ci_low)
                        hi_raw = float(ci_high)
                        if lo_raw > 0 and hi_raw > 0:
                            lo_cpu = 1.0 / hi_raw  # invert + swap
                            hi_cpu = 1.0 / lo_raw
                            ci_str = f"[{lo_cpu:.0f} - {hi_cpu:.0f}]"
                    except (TypeError, ValueError, ZeroDivisionError):
                        pass
                else:
                    ci_str = f"[{float(ci_low):.1f} - {float(ci_high):.1f}]"

            rows.append((name, budget_str, contrib_str, roi_str, share_str, verdict, footnote, ci_str))
        return rows

    # ----------------------------------------------------------------
    # SLIDE 01 - COVER
    # ----------------------------------------------------------------

    def s01_cover(self):
        slide = self._blank()

        # Full-height thin gold vertical line (right margin) - visual anchor
        self._vline(
            slide, self.w - self.safe - 0.1,
            self.safe, self.h - 2 * self.safe,
            weight=1.0, color=self.gold,
        )

        # 5d (2026-05-04 + audit-iter 2026-05-04): cuted compact PNG в верхнем левом.
        # PNG already содержит sigil + AURORA AI text (one composition) → НЕ нужен
        # отдельный _wordmark() call который дублировал бы текст. Customer-approved
        # layout per mmm_report_20260504_024312.pptx reference.
        # Position (0.40, 0.15) - close к slide edge, height 0.78" (matches reference).
        brand_path = self._brand_mark_compact_path()
        if brand_path is not None:
            slide.shapes.add_picture(
                str(brand_path),
                Inches(self.safe), Inches(0.15),
                height=Inches(0.78),
            )
        else:
            # Fallback: typographic wordmark for legacy environments без bundled PNG.
            self._wordmark(slide, self.safe, self.safe, size=14, color=self.deep_100)

        # Horizontal hairline full-width (margin to right-gold-line)
        self._hairline(
            slide, self.safe, self.safe + 0.55,
            self.w - 2 * self.safe - 0.15,
            weight=0.5,
        )

        # Category tag
        self._text(
            slide, self.safe, self.safe + 0.65, 10.0, 0.25,
            "MARKETING MIX MODEL REPORT",
            font=self.sans, size=9, color=self.deep_60,
        )

        # Main title - left aligned, serif regular (not bold - confidence)
        self._text(
            slide, self.safe, 2.6, self.w - 2 * self.safe, 1.3,
            "Декомпозиция медиабюджета",
            font=self.serif, size=48, bold=False, color=self.deep_100,
            line_spacing=1.0,
        )
        # Subtitle in italic
        self._text(
            slide, self.safe, 4.0, self.w - 2 * self.safe, 0.55,
            "и рекомендации по оптимизации",
            font=self.serif, size=24, italic=True, color=self.deep_80,
        )

        # Sacred lime under title
        self._lime_under(slide, self.safe, 4.65, 1.4)

        # 4-column metadata grid bottom
        grid_y = 6.1
        grid_x = self.safe
        cols = [
            ("ПОДГОТОВЛЕНО ДЛЯ", self.client),
            ("ДАТА",              self.report_date),
            ("REPORT ID",         self.report_id),
            ("КЛАССИФИКАЦИЯ",    "Confidential"),
        ]
        col_w = (self.w - 2 * self.safe - 0.2) / 4
        for i, (label, val) in enumerate(cols):
            cx = grid_x + i * col_w
            self._text(
                slide, cx, grid_y, col_w, 0.2, label,
                font=self.sans, size=8, bold=True, color=self.deep_60,
            )
            self._text(
                slide, cx, grid_y + 0.28, col_w, 0.4, val,
                font=self.serif, size=16, color=self.deep_100,
            )
            # Column divider (except last)
            if i < len(cols) - 1:
                self._vline(
                    slide, cx + col_w - 0.05, grid_y + 0.05,
                    0.65, weight=0.25,
                )

        # Bottom hairline + confidentiality (text vertically centered between hairline 7.10 and slide bottom 7.50)
        self._hairline(slide, self.safe, 7.10, self.w - 2 * self.safe - 0.15, weight=0.5)
        # Midpoint 7.30; text box height 0.15 → y = 7.225
        self._text(
            slide, self.safe, 7.225, self.w - 2 * self.safe - 0.15, 0.15,
            self.copyright_line,
            font=self.sans, size=7, italic=True, color=self.deep_60,
            align=PP_ALIGN.CENTER,
        )

    # ----------------------------------------------------------------
    # SLIDE 02 - AT A GLANCE (5 key findings)
    # ----------------------------------------------------------------

    def s02_at_a_glance(self):
        slide = self._blank()
        self._header(slide, slide_num=3)

        self._category(slide, self.safe, 0.60, "ОТЧЁТ ЗА 60 СЕКУНД")
        self._text(
            slide, self.safe, 0.70, 10.0, 0.7,
            "Пять находок",
            font=self.serif, size=32, color=self.deep_100,
        )
        self._hairline(slide, self.safe, 1.50, 1.2, weight=0.75, color=self.gold)

        # Five findings - slot-fill from facts when channels present,
        # else Kagocel pilot text (preview / wireframe mode).
        findings = self._build_at_a_glance_findings() if (self.facts and self.channels) else [
            ("01", "TV обеспечивает 42% инкрементальных продаж при 28% доли бюджета",
             "ROI 1.8× выше среднего по каналам"),
            ("02", "Насыщение на TV начинается с 80 TRP/нед",
             "Предельный ROI падает на 22% относительно IV кв. 2025"),
            ("03", "Digital video - самый эффективный канал с mROAS 1.9×",
             "Текущий бюджет на нём меньше 15%"),
            ("04", "Базовый уровень растёт на 8% год к году – кампании работают на долгосроке",  # П8-1
             "Бренд-эффект виден в динамике"),
            ("05", "Рекомендация: перераспределить 25 млн из TV в digital video",
             "Ожидаемый прирост ROAS: +12 пп к III-IV кв. 2026"),
        ]
        # Stage C.4: findings_y 1.80 → 1.95 (more breathing room below title);
        # finding height 0.45 → 0.55, support y-offset 0.42 → 0.55,
        # step 0.92 → 1.02 (italic subtitle was overlapping next finding).
        y = 1.95
        for i, (num, finding, support) in enumerate(findings):
            # Number
            self._text(
                slide, self.safe, y, 0.7, 0.5, num,
                font=self.serif, size=28, color=self.gold,
            )
            # Finding in Georgia bold
            self._text(
                slide, self.safe + 0.9, y, 9.0, 0.55, finding,
                font=self.serif, size=15, bold=True, color=self.deep_100,
            )
            # Support text (moved lower so italic subtitle doesn't overlap)
            self._text(
                slide, self.safe + 0.9, y + 0.55, 9.0, 0.3, support,
                font=self.sans, size=10, italic=True, color=self.deep_60,
            )
            # Hairline between items only (not after last - footer has its own rule)
            if i < len(findings) - 1:
                self._hairline(slide, self.safe, y + 0.92, self.w - 2 * self.safe, weight=0.25)
            y += 1.02

        self._footer(slide, 3)  # Stage C.6.1: was 2 (TOC moved ahead of Executive Summary)

    # ----------------------------------------------------------------
    # SLIDE 03 - TOC
    # ----------------------------------------------------------------

    def s03_toc(self):
        slide = self._blank()
        self._header(slide, slide_num=2)

        # B4 (стайлгайд §2): категория «AGENDA» удалена — англицизм и дубль
        # заголовка «Содержание отчёта» строкой ниже (Ильяхов: убрать лишнее).
        self._text(
            slide, self.safe, 0.70, 10.0, 0.7, "Содержание отчёта",
            font=self.serif, size=32, color=self.deep_100,
        )
        self._hairline(slide, self.safe, 1.50, 1.2, weight=0.75, color=self.gold)

        # Main list (5 sections - Stage C.6.3 honest TOC)
        y = 1.85
        for i, (name, pg) in enumerate(zip(self.section_names, self.toc_page_refs), start=1):
            # Number
            self._text(
                slide, self.safe, y, 0.8, 0.4,
                f"{i:02d}", font=self.serif, size=20, color=self.gold,
            )
            # Section name
            self._text(
                slide, self.safe + 0.9, y + 0.05, 7.5, 0.3,
                name, font=self.serif, size=17, color=self.deep_100,
            )
            # Leader dots (dense, tier-1 TOC style)
            leader = "." * 80
            self._text(
                slide, 5.5, y + 0.08, 3.8, 0.3, leader,
                font=self.sans, size=9, color=self.deep_40,
            )
            # Page num (just number, no "стр.")
            self._text(
                slide, 9.5, y + 0.05, 0.6, 0.3,
                f"{pg:02d}",
                font=self.serif, size=14, color=self.deep_100,
                align=PP_ALIGN.RIGHT,
            )
            # Hairline
            self._hairline(slide, self.safe, y + 0.5, 9.7, weight=0.25)
            y += 0.5
            # P-2 fix (2026-07-16): вставной слайд «Прогноз на будущий период»
            # живёт внутри секции «Главное» и не имел следа в оглавлении —
            # приёмка «раздел планирования не найден». Подстрока делает его
            # находимым, не перекраивая сетку из 5 секций.
            if i == 1 and self.forecast:
                fc_pg = 6 + int(bool(self.backtest)) + int(bool(self.gen_compare))
                # y уже за hairline секции (+0.5 выше); +0.02 — чтобы линия
                # не прочерчивала подпись (находка аудита 2026-07-16).
                self._text(
                    slide, self.safe + 0.9, y + 0.02, 7.5, 0.25,
                    f"в том числе «Прогноз на будущий период» — стр. {fc_pg:02d}",
                    font=self.sans, size=10, color=self.deep_60,
                )
                y += 0.32

        # Right metadata sidebar
        side_x = 10.8
        side_y = 1.85
        side_w = self.w - self.safe - side_x
        # Card bg
        self._rect(slide, side_x, side_y, side_w, 4.0, fill=self.bg_quiet)
        self._text(
            slide, side_x + 0.2, side_y + 0.2, side_w - 0.4, 0.25,
            "ЧТЕНИЕ", font=self.sans, size=8, bold=True, color=self.gold,
        )
        self._hairline(slide, side_x + 0.2, side_y + 0.47, 1.0, weight=0.5, color=self.gold)

        # B1-fix R-11 (2026-07-03): «Время/Слов» — декоративные wireframe-оценки,
        # в live-режиме не показываем (не выдумываем факты о собственном отчёте).
        meta = [
            *([] if self.is_live else [("Время", "~12 мин")]),
            ("Страниц",          f"{self.total_slides}"),
            ("Разделов",         f"{self.total_sections}"),
            *([] if self.is_live else [("Таблиц", "3"), ("Графиков", "4"), ("Слов", "~2 800")]),
            ("MQS модели",       self._mstr(self.mqs_score, "{:.0f} / 100")),
            ("Данные",           self.data_window_label or "—"),
        ]
        my = side_y + 0.7
        for label, val in meta:
            self._text(
                slide, side_x + 0.2, my, 0.9, 0.2, label,
                font=self.sans, size=9, color=self.deep_60,
            )
            self._text(
                slide, side_x + 1.1, my, side_w - 1.3, 0.2, val,
                font=self.sans, size=9, bold=True, color=self.deep_100,
                align=PP_ALIGN.RIGHT,
            )
            my += 0.3

        self._footer(slide, 2)  # Stage C.6.1: was 3 (TOC now physical slide 2)

    # ----------------------------------------------------------------
    # SLIDE 04 - SECTION DIVIDER WITH TAKEAWAY
    # ----------------------------------------------------------------

    def _render_section_divider(self, slide_num, *, takeaway, topics):
        """Tier-1 section-divider slide with big number, takeaway, topic list.

        Stage C.6.3: extracted from s04 for reuse by 3 new dividers
        (Методология / Данные / Приложение). `takeaway` may be any string;
        `topics` is a list of short phrases rendered as bullets.
        """
        slide = self._blank()
        self._header(slide, slide_num=slide_num, include_confidential=True)
        section_idx, section_label = self.slide_to_section[slide_num]

        # Big number
        self._text(
            slide, self.safe, 1.3, 4.5, 3.5,
            f"{section_idx:02d}",
            font=self.serif, size=220, color=self.deep_20,
        )

        # Section label
        self._text(
            slide, self.safe + 4.3, 2.1, 8.0, 0.3,
            f"РАЗДЕЛ {section_idx:02d} / {self.total_sections:02d}",
            font=self.sans, size=9, bold=True, color=self.gold,
        )
        # Section name
        self._text(
            slide, self.safe + 4.3, 2.5, 8.0, 0.8,
            section_label,
            font=self.serif, size=40, color=self.deep_100,
        )
        self._lime_under(slide, self.safe + 4.3, 3.35, 2.5)

        # KEY TAKEAWAY
        self._text(
            slide, self.safe + 4.3, 3.7, 8.0, 0.3, "ОСНОВНОЙ ВЫВОД",
            font=self.sans, size=8, bold=True, color=self.gold,
        )
        self._text(
            slide, self.safe + 4.3, 4.05, 8.0, 1.3,
            takeaway,
            font=self.serif, size=20, italic=True, color=self.deep_80,
            line_spacing=1.3,
        )

        # In this section
        self._text(
            slide, self.safe + 4.3, 5.6, 8.0, 0.25, "В ЭТОМ РАЗДЕЛЕ",
            font=self.sans, size=8, bold=True, color=self.deep_60,
        )
        ty = 5.9
        for topic in topics:
            self._text(
                slide, self.safe + 4.3, ty, 8.0, 0.3, f"·  {topic}",
                font=self.sans, size=11, color=self.deep_100,
            )
            ty += 0.28

        # Progress bar
        self._section_progress(slide, self.safe, 6.3, current=section_idx)

        self._footer(slide, slide_num)
        return slide

    def s04_section_divider(self):
        """Divider for section 2 'Декомпозиция вкладов'.

        B4-2 (2026-07-03): ВНЕ ПОТОКА build() — отдельные слайды-дивайдеры
        заменены компактными плашками _section_intro (стайлгайд §1).
        Метод сохранён для возможного возврата/кастомных сборок."""
        if self.facts:
            leader = self.facts.get("leader_channel") or "Лидер"
            cpct = self.facts.get("leader_share_contrib_pct")
            spct = self.facts.get("leader_share_spend_pct")
            if cpct is not None and spct is not None:
                # B1-fix R-12: cpct — доля в МЕДИА-вкладе, не в продажах
                # (без квалификатора клиент завышает роль канала кратно).
                takeaway = (
                    f"{leader} даёт {_fmt_pct(cpct)} медиа-вклада при {_fmt_pct(spct)} бюджета - "
                    "основная точка оптимизации портфеля"
                )
            else:
                takeaway = f"{leader} - основной драйвер портфеля и точка оптимизации"
        else:
            takeaway = "TV генерирует 42% продаж при 28% бюджета - основная точка оптимизации портфеля"
        self._render_section_divider(
            slide_num=6,
            takeaway=takeaway,
            topics=[
                "Индивидуальные вклады каналов и ROI ранжирование",
                "Портфельная таблица с вердиктами по каналам",
                f"Декомпозиция динамики продаж {self._period_unit_label()}",
            ],
        )

    def s_divider_methodology(self):
        """NEW divider for section 3 'Методология'."""
        if self.is_ols:
            self._render_section_divider(
                slide_num=10,
                takeaway=(
                    "Линейная регрессия с отложенным эффектом (adstock) и Hill-насыщением – "  # П8-2 П8-1
                    "прозрачная математическая модель с bootstrap-интервалами"
                ),
                topics=[
                    "Спецификация модели и уравнение отклика",
                    "Параметры отложенного эффекта (adstock) и насыщения (фиксированные)",  # П8-2
                    "Диагностика OLS · closed-form · bootstrap CI",
                ],
            )
        else:
            self._render_section_divider(
                slide_num=10,
                takeaway=(
                    "Байесовский MMM с адстоком и Hill-насыщением - "
                    "прозрачная математическая модель с интервалами доверия"
                ),
                topics=[
                    "Спецификация модели и уравнение отклика",
                    "Априорные распределения и ограничения параметров",
                    "Диагностика сходимости MCMC",
                ],
            )

    def s_divider_data(self):
        """NEW divider for section 4 'Данные и качество'."""
        tb_mln = self.facts.get("total_budget_mln") if self.facts else None
        # B1-fix R-02: период — только когда реально известен; MQS через _mstr;
        # «готовность к production» не заявляем без основания (wireframe-ветка).
        _mqs_part = f"MQS модели {self._mstr(self.mqs_score, '{:.0f}')}/100"
        _window_part = (
            f"Данные охватывают {self.data_window_label}" if self.data_window_label
            else "Данные загружены из файла клиента"
        )
        if tb_mln:
            takeaway = f"{_window_part}, бюджет {tb_mln:.0f} млн руб - {_mqs_part}"
        elif self.is_live:
            takeaway = f"{_window_part} - {_mqs_part}"
        else:
            takeaway = f"{_window_part} - {_mqs_part}, готовность к production"
        self._render_section_divider(
            slide_num=12,
            takeaway=takeaway,
            topics=[
                "Источники данных и период наблюдения",
                "Процедуры очистки и обработка выбросов",
                "Метрики качества модели и спецификация",
            ],
        )

    def s_divider_appendix(self):
        """NEW divider for section 5 'Приложение и источники'."""
        self._render_section_divider(
            slide_num=14,
            takeaway=(
                "Справочные материалы отчёта - глоссарий терминов, "
                "методологические ссылки и контактная информация"
            ),
            topics=[
                "Глоссарий ключевых терминов и метрик",
                "Источники данных и референсы",
                "Контакты команды Aurora AI",
            ],
        )

    # ----------------------------------------------------------------
    # SLIDE 05 - KEY MESSAGE (Big Number + Pull Quote)
    # ----------------------------------------------------------------

    def s05_key_message(self):
        slide = self._blank()
        self._header(slide, slide_num=4)

        self._category(slide, self.safe, 0.60, "КЛЮЧЕВОЙ ВЫВОД")

        # Title + big number + pull quote - slot-fill from facts when present,
        # else Kagocel pilot (preview / wireframe mode).
        if self.facts and self.channels:
            leader = self.facts.get("leader_channel") or "Лидер"
            hero = self.facts.get("hero_channel") or leader
            cpct = self.facts.get("leader_share_contrib_pct")
            spct = self.facts.get("leader_share_spend_pct")
            wr = self.facts.get("weighted_roi")
            honest = bool(self.facts.get("honest_narrative"))
            media_pct = self.facts.get("media_contribution_pct")
            baseline_pct = self.facts.get("baseline_pct")

            # v1.3.2: KPI-aware portfolio metric phrase.
            portfolio_phrase = (
                f"ROI портфеля {wr:.2f}× средневзвешенный" if (wr is not None and self.kpi["is_legacy"])
                else (f"{_weighted_summary_phrase_pptx(wr, self.kpi)} средневзвешенный" if wr is not None else "")
            )

            if honest and media_pct is not None and baseline_pct is not None:
                # Honest narrative: baseline-dominated model (media < 10%).
                # Disclose this rather than leading with leader's media-share.
                title = "Модель преимущественно отражает базовый спрос – медиа-вклад ограничен"  # П8-2 П8-1
                big_number = _fmt_pct(media_pct)
                big_label = "Медиа-вклад в продажи"
                big_support = (
                    f"Базовый спрос: {_fmt_pct(baseline_pct)}. {portfolio_phrase}."  # П8-2
                    if portfolio_phrase else f"Базовый спрос: {_fmt_pct(baseline_pct)}."  # П8-2
                )
                quote_txt = (
                    f"{leader} – лидер среди медиа ({_fmt_pct(cpct)} медиа-вклада), "  # П8-2 П8-1
                    f"но абсолютный медиа-эффект {_fmt_pct(media_pct)} от продаж. "  # П8-2
                    "Низкий вклад медиа – проверить отложенный эффект (adstock), насыщение, качество данных."  # П8-2 П8-1
                )
            else:
                # Action title - leader's position statement
                title = f"{leader} остаётся основным драйвером, но эффективность требует проверки насыщения"
                # Big number - leader contribution share
                big_number = _fmt_pct(cpct) if cpct is not None else "-"
                big_label = f"Доля {leader} в инкрементальных продажах"
                if spct is not None and portfolio_phrase:
                    big_support = f"При {_fmt_pct(spct)} доли бюджета. {portfolio_phrase}."
                else:
                    big_support = "Лидер по вкладу в продажи"
                # Pull quote - hero outperforms leader, reallocate signal
                # Пласт 2 (2026-07-11): KPI-aware, не хардкодить «рубль/ROAS» для count/effectiveness.
                if hero != leader:
                    quote_txt = (
                        f"{_hero_vs_leader_quote_pptx(hero, leader, self.kpi)} "
                        "Сигнал к перераспределению части бюджета."
                    )
                else:
                    quote_txt = (
                        f"{leader} - единственный лидер и по вкладу, и по эффективности. "
                        "Бюджет следует сохранить до признаков насыщения."
                    )
        else:
            title = "TV остаётся основным драйвером, но эффективность достигла локального максимума"
            big_number = "42%"
            big_label = "Доля TV в инкрементальных продажах, I кв."
            big_support = "При 28% доли бюджета. ROI 1.8× выше среднего."
            quote_txt = (
                "Каждый рубль в TV возвращает в 1.8 раза больше, "
                "чем среднее по каналам. Однако начиная с 80 TRP/нед "
                "маржинальный возврат падает - сигнал к перераспределению в digital."
            )

        self._action_title(
            slide,
            title,
            show_lime=True, y=0.80, height=0.90,
        )

        # Left: big number
        self._big_number(
            slide, self.safe, 3.3,
            number=big_number,
            label=big_label,
            support=big_support,
            size=140,
        )

        # Right: pull quote
        quote_x = 6.8
        quote_y = 3.1
        quote_w = self.w - self.safe - quote_x
        self._pull_quote(
            slide, quote_x, quote_y, quote_w, 2.6,
            quote_txt,
            size=18, with_bar=True,
        )

        # Attribution
        self._text(
            slide, quote_x, quote_y + 2.7, quote_w, 0.25,
            "ВЫВОД ИЗ MMM-ДЕКОМПОЗИЦИИ",
            font=self.sans, size=8, bold=True, color=self.deep_60,
        )

        # Source footnote at bottom
        # B1-fix R-02/R-06: период — только реальный; «откалибрована под
        # FMCG-бенчмарки» — недоказуемое заявление, в live не пишем.
        _engine_label = "OLS MMM" if self.is_ols else "Bayesian MMM"
        _window_sfx = f" {self.data_window_label}" if self.data_window_label else ""
        _src_text = f"Источник: {_engine_label} · {self.report_id}; данные {self.sources_client_label}{_window_sfx}."
        if not self.is_live:
            _src_text += " Модель откалибрована под FMCG-бенчмарки."
        self._source(slide, 6.87, text=_src_text)

        self._footer(slide, 4)

    # ----------------------------------------------------------------
    # SLIDE 06 - ACTION + CHART + COMMENTARY (with annotation)
    # ----------------------------------------------------------------

    def s06_action_chart(self):
        slide = self._blank()
        self._header(slide, slide_num=6 + self._page_shift)

        # B4-2: плашка секции вместо отдельного слайда-дивайдера.
        if self.facts:
            _l = self.facts.get("leader_channel") or "Лидер"
            _c = self.facts.get("leader_share_contrib_pct")
            _s = self.facts.get("leader_share_spend_pct")
            _tk = (f"{_l} даёт {_fmt_pct(_c)} медиа-вклада при {_fmt_pct(_s)} бюджета"
                   if _c is not None and _s is not None
                   else f"{_l} - основной драйвер портфеля")
        else:
            _tk = None
        self._section_intro(slide, 2, "Декомпозиция вкладов", _tk)

        # v1.3.2: KPI-aware section category + fallback action title.
        if self.kpi["mode"] == "effectiveness":
            category_text = "ДОЛЯ КАНАЛОВ В ЭФФЕКТЕ"
            fallback_action = "Сбалансировать портфель по доле эффекта"
        elif self.kpi["kpi_kind"] == "count":
            category_text = "CPU ПО КАНАЛАМ"
            fallback_action = "Сбалансировать портфель по CPU"
        else:
            category_text = "ROI ПО КАНАЛАМ"
            fallback_action = "Сбалансировать портфель по mROAS"

        # B4-2: категория-слово заменена секционной плашкой выше (KPI-нюанс
        # категории живёт в заголовке графика chart_title_text ниже).
        _ = category_text

        # Stage C.5: McKinsey action-first headline via narrative_adapter.
        # Shared helper keeps PPTX+HTML in sync and applies zero-effect guard.
        action_title = (
            derive_action_headline(self.channels, self.facts, "mroas")
            or fallback_action
        )

        # Stage C.4: action title on s06 can run 4+ lines when channel
        # names are long. Height bumped to 1.10 (was 0.80) + font 20pt
        # (was 22pt) to prevent overflow into label row below.
        self._action_title(
            slide,
            action_title,
            show_lime=True, y=0.80, height=1.10, size=20,
        )

        # Chart zone (left 58%)
        chart_x = self.safe
        chart_y = 1.95
        chart_w = (self.w - 2 * self.safe) * 0.58
        chart_h = 3.7

        # Chart title & subtitle - v1.3.2 KPI-aware.
        if self.kpi["mode"] == "effectiveness":
            chart_title_text = "ДОЛЯ КАНАЛОВ В ЭФФЕКТЕ / %"
        elif self.kpi["kpi_kind"] == "count":
            chart_title_text = "CPU ПО КАНАЛАМ / ₽ ЗА ЕДИНИЦУ"
        else:
            chart_title_text = "MROAS ПО КАНАЛАМ / МУЛЬТИПЛИКАТОР"
        self._text(
            slide, chart_x, chart_y, chart_w, 0.25,
            chart_title_text,
            font=self.sans, size=9, bold=True, color=self.deep_80,
        )
        # v1.3.2: chart subtitle adapts per KPI.
        # B1-fix R-02: период только когда известен (live без wireframe «Q1 2026»).
        _period_sfx = f", {self.period_label}" if self.period_label else ""
        if self.kpi["mode"] == "effectiveness":
            chart_subtitle_text = f"Доли каналов в продажах{_period_sfx}"
        elif self.kpi["kpi_kind"] == "count":
            chart_subtitle_text = f"Стоимость следующей единицы (incremental cost-per-unit){_period_sfx}"
        else:
            chart_subtitle_text = f"Marginal ROI последнего вложенного рубля{_period_sfx}"
        self._text(
            slide, chart_x, chart_y + 0.27, chart_w, 0.22,
            chart_subtitle_text,
            font=self.sans, size=9, italic=True, color=self.deep_60,
        )
        # Hairline removed per brand rule - minimize horizontal lines

        # NATIVE PPTX CHART (production-ready) - editable в PowerPoint
        # Data source: adapter-supplied channels sorted by mROAS desc (up to 10),
        # else Kagocel pilot bars for preview / wireframe mode.
        #
        # v1.3.2 audit fix (B1+B4): bar values transformed per KPI mode:
        # - monetary roi: raw mROAS× (no transform; sort desc, hero = highest).
        # - count: invert units/₽ → CPU ₽/ед. (sort ascending by CPU = lowest
        #   cost first = hero is best).
        # - effectiveness: keep as fraction (0..1); native '0.0%' format
        #   auto-multiplies by 100. Sort desc - лидер по доле.
        if self.channels:
            kpi_mode = self.kpi["mode"]
            kpi_kind = self.kpi["kpi_kind"]
            valid = [c for c in self.channels if float(c.get("mroas") or 0) > 0]
            if kpi_kind == "count" and kpi_mode != "effectiveness":
                # CPU semantics: lowest CPU = best. Compute CPU = 1/raw_mroas,
                # sort ascending (best first → hero at index 0).
                annotated = [
                    (c.get("name") or "-", 1.0 / float(c.get("mroas") or 1.0))
                    for c in valid
                ]
                annotated.sort(key=lambda t: t[1])  # ascending CPU
                annotated = annotated[:10]
                bar_labels = [name for name, _ in annotated]
                bar_values = [cpu for _, cpu in annotated]
            else:
                by_mroas = sorted(
                    valid,
                    key=lambda c: float(c.get("mroas") or 0),
                    reverse=True,
                )[:10]
                bar_labels = [c.get("name") or "-" for c in by_mroas]
                bar_values = [float(c.get("mroas") or 0) for c in by_mroas]
        else:
            bar_labels = ["Digital video", "Search", "TV", "OOH", "Social", "Print"]
            bar_values = [1.9, 1.7, 1.5, 1.2, 1.0, 0.7]
        hero_idx = 0  # always best after KPI-aware sort - hero at index 0

        chart_data = CategoryChartData()
        # Reversed: PPTX bar chart renders first category at bottom
        chart_data.categories = list(reversed(bar_labels))
        # v1.3.2: series label per KPI.
        chart_data.add_series(self.kpi["metric_short"], list(reversed(bar_values)))

        bar_area_x = chart_x
        bar_area_y = chart_y + 0.75
        bar_area_w = chart_w
        bar_area_h = 2.9

        graphic_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED,
            Inches(bar_area_x), Inches(bar_area_y),
            Inches(bar_area_w), Inches(bar_area_h),
            chart_data,
        )
        chart = graphic_frame.chart
        chart.has_legend = False
        chart.has_title = False

        # Color discipline: ONE gold hero bar (Digital video = last after reverse), others muted
        series = chart.plots[0].series[0]
        # reversed order - hero is at index len-1
        reversed_hero = len(bar_values) - 1 - hero_idx
        for i, point in enumerate(series.points):
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = self.gold if i == reversed_hero else self.deep_40
            point.format.line.fill.background()

        # Data labels on bar ends - v1.3.2 audit fix (B1): format per KPI.
        # - effectiveness: native '0.0%' format auto-multiplies fraction by 100
        #   (PPTX/Excel built-in % handling). Pre-fix '0.0"%"' literal "%" suffix
        #   showed 0.25 → "0.25%" instead of 25.0%.
        # - count: bars already inverted to CPU (₽/ед.) above. Plain integer
        #   format с unit suffix in literal text.
        # - monetary: × multiplier, 1-decimal.
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        if self.kpi["mode"] == "effectiveness":
            data_labels.number_format = '0.0%'  # native percent - auto-multiplies
        elif self.kpi["kpi_kind"] == "count":
            data_labels.number_format = '0" ₽/ед."'
        else:
            data_labels.number_format = '0.0"×"'
        data_labels.font.size = Pt(10)
        data_labels.font.name = self.sans
        data_labels.font.color.rgb = self.deep_80
        from pptx.enum.chart import XL_DATA_LABEL_POSITION
        try:
            data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
        except Exception:
            pass

        # Axis styling - minimalist (tier-1 MBB spec)
        cat_axis = chart.category_axis
        cat_axis.tick_labels.font.size = Pt(10)
        cat_axis.tick_labels.font.name = self.sans
        cat_axis.tick_labels.font.color.rgb = self.deep_100
        cat_axis.format.line.fill.background()  # no axis line

        val_axis = chart.value_axis
        val_axis.visible = False  # hide numeric value axis (labels are direct)
        val_axis.minimum_scale = 0
        # v1.3.2 audit fix: adaptive axis scale + tick interval per KPI.
        # - monetary: mROAS× values 0..5+; tick 0.5, min span 2.2.
        # - count: CPU ₽/ед. values 10..200+; tick auto-derived.
        # - effectiveness: fractions 0..1 native; tick 0.2, max ~1.15.
        _max_v = max(bar_values) if bar_values else 2.0
        if self.kpi["mode"] == "effectiveness":
            val_axis.major_unit = 0.2
            val_axis.maximum_scale = max(1.0, _max_v * 1.15)
        elif self.kpi["kpi_kind"] == "count":
            # CPU domain: tick ≈ 1/10 of max; max accommodate clipping headroom.
            val_axis.major_unit = max(10, _max_v / 10)
            val_axis.maximum_scale = max(20, _max_v * 1.15)
        else:
            val_axis.major_unit = 0.5
            val_axis.maximum_scale = max(2.2, _max_v * 1.15)

        # Gap between bars
        from pptx.oxml.ns import qn
        ser = series._element
        ser_pr = ser.find(qn('c:spPr'))
        # Set gap width via XML (python-pptx limitation for gap_width)
        try:
            bar_chart = chart.plots[0]._element
            gap = bar_chart.find(qn('c:gapWidth'))
            if gap is not None:
                gap.set('val', '60')
        except Exception:
            pass

        # Breakeven reference note - v1.3.2: text adapts per KPI/mode.
        if self.kpi["mode"] == "effectiveness":
            breakeven_note = "Доля > среднего по портфелю = недо-инвестирован"
        elif self.kpi["kpi_kind"] == "count":
            vpcu = self.kpi.get("vpcu")
            if vpcu:
                breakeven_note = f"CPU < {float(vpcu):.0f} ₽/ед. = окупается  ·  ниже = прибыльно"
            else:
                breakeven_note = "Чем ниже CPU, тем дешевле каждая единица"
        else:
            breakeven_note = "1.0× = безубыточность  ·  выше = прибыльно"
        self._text(
            slide, chart_x + chart_w - 2.5, bar_area_y - 0.22, 2.5, 0.18,
            breakeven_note,
            font=self.sans, size=7, italic=True, color=self.deep_60,
            align=PP_ALIGN.RIGHT,
        )

        # Source at bottom (unified position max low to footer hairline)
        _src_text = (
            f"Источник: OLS MMM Aurora AI · {self.report_id}; точечные оценки + bootstrap CI"
            if self.is_ols
            else f"Источник: Bayesian MMM Aurora AI · {self.report_id}; медианы апостериорного распределения"
        )
        self._source(slide, 6.87, text=_src_text)

        # Right rail: commentary (42%)
        right_x = chart_x + chart_w + 0.4
        right_w = self.w - self.safe - right_x

        self._text(
            slide, right_x, chart_y, right_w, 0.25,
            "ЧТО ЭТО ЗНАЧИТ",
            font=self.sans, size=9, bold=True, color=self.gold,
        )
        self._hairline(slide, right_x, chart_y + 0.3, 1.0, weight=0.75, color=self.gold)

        # Commentary - math-fix v1.0.14.1 B refactor (2026-04-28).
        # Pre-fix: hardcoded «явный потенциал для наращивания» / «потенциал
        # удержания» / «топ-2 канала» blocks based на mROAS rank - independent
        # от derive_verdict в action table → contradictions.
        # Post-fix: action-driven. Reads ch['action_label']/['action_reasoning']
        # populated narrative_adapter via single-source-of-truth
        # engines.channel_action.compute_channel_action. Action в s07 table cell
        # + s06 commentary lead garanteed identical per channel.
        if self.channels and self.facts:
            # Sort by action priority (Scale=5 first, ..., Cut=0). De-dup по
            # action key so same action не shown 3× для 3 Scale channels.
            by_priority = sorted(
                self.channels,
                key=lambda c: (
                    -int(c.get("action_priority") or 0),
                    -float(c.get("mroas") or 0),
                ),
            )
            commentary: list[tuple[str, str]] = []
            seen_actions: set[str] = set()
            for ch in by_priority:
                ch_action = ch.get("action") or "Watch"
                if ch_action == "Uncertain":
                    continue
                if ch_action in seen_actions:
                    continue
                seen_actions.add(ch_action)
                ch_name = ch.get("name") or "-"
                label = ch.get("action_label") or ch_action
                if ch.get("action_reasoning"):
                    reasoning = ch["action_reasoning"]
                elif self.kpi["is_legacy"]:
                    reasoning = f"mROAS {float(ch.get('mroas') or 0):.1f}× - рекомендация по портфелю."
                else:
                    metric_val = _fmt_metric_pptx(ch.get('mroas') or 0, self.kpi)
                    reasoning = f"{self.kpi['metric_short']} {metric_val} - рекомендация по портфелю."
                commentary.append((
                    f"{ch_name} - {label}.",
                    f" {reasoning}",
                ))
                if len(commentary) >= 3:
                    break
            # Fallback (legacy callers без action decoration)
            if not commentary:
                hero = by_priority[0] if by_priority else {}
                hero_name = hero.get("name") or "Лидер"
                hero_m = float(hero.get("mroas") or 0)
                if self.kpi["is_legacy"]:
                    commentary.append((
                        f"{hero_name} - лидер по mROAS.",
                        f" mROAS {hero_m:.1f}×. Бюджет следует пересмотреть с учётом насыщения.",
                    ))
                else:
                    metric_short = self.kpi["metric_short"]
                    metric_val = _fmt_metric_pptx(hero_m, self.kpi)
                    commentary.append((
                        f"{hero_name} - лидер по {metric_short}.",
                        f" {metric_short} {metric_val}. Бюджет следует пересмотреть с учётом насыщения.",
                    ))
        else:
            # Wireframe placeholder when no channels (preview mode)
            commentary = [
                ("Digital video - Масштабировать.",
                 " mROAS 1.9× - Optimizer рекомендует +50%, недо-инвестирован."),
                ("Search - Удерживать.",
                 " mROAS 1.7× стабилен, gap +1пп - баланс."),
                ("Print и Radio - Сократить.",
                 " mROAS 0.7-0.75× ниже breakeven - бюджет приносит убыток."),
            ]
        cy = chart_y + 0.55
        for lead, body in commentary:
            self._rich(
                slide, right_x, cy, right_w, 1.3,
                runs=[
                    (lead, {"font": self.sans, "size": 11, "bold": True, "color": self.deep_100}),
                    (body, {"font": self.sans, "size": 11, "color": self.deep_80}),
                ],
                line_spacing=1.35,
            )
            self._vbar(slide, right_x - 0.15, cy + 0.03, 1.2, weight=2, color=self.gold)
            cy += 1.35

        self._footer(slide, 6 + self._page_shift)

    # ----------------------------------------------------------------
    # SLIDE 07 - ACTION + TABLE (with conditional formatting & footnotes)
    # ----------------------------------------------------------------

    def s07_action_table(self):
        slide = self._blank()
        self._header(slide, slide_num=7 + self._page_shift)

        self._category(slide, self.safe, 0.60, "ПОРТФЕЛЬ КАНАЛОВ")

        # Action title - data-driven: top-N contributors covering >=85% of
        # incremental sales. Handles edge cases (single channel, all-zero
        # contributions, balanced portfolio) with idiomatic Russian phrasing.
        # Fallback ("Пять каналов...") applies only in preview / wireframe
        # mode when no channels supplied.
        # Stage C.5: McKinsey action-first headline via narrative_adapter.
        # Covers all edge cases (all-zero / single-channel / consolidate / balanced).
        s07_title = (
            derive_action_headline(self.channels, self.facts, "portfolio")
            or "Консолидировать до топ-5 каналов - они обеспечивают 87% продаж"
        )
        self._action_title(
            slide, s07_title,
            show_lime=True, y=0.80, height=0.80,
        )

        # Custom table
        table_x = self.safe
        table_y = 1.95
        table_w = self.w - 2 * self.safe
        # Proportional columns: Channel / Budget / Contribution / ROI / Share / Verdict
        col_weights = [3.2, 1.8, 1.8, 1.4, 1.5, 1.8]
        total_w = sum(col_weights)
        col_widths = [w * (table_w / total_w) for w in col_weights]

        # Header - v1.3.2: main metric column adapts per KPI (mROAS/CPU/Доля).
        # Пласт 2 (2026-07-11): contrib column unit — для count вклад НЕ в рублях.
        metric_col_hdr, metric_col_unit = _table_metric_header_pptx(self.kpi)
        # Единица столбца «Вклад» — согласованно с масштабом ячеек в
        # _build_action_table_rows (fix 2026-07-13, INV-50).
        contrib_scale_s07, contrib_unit_pptx = _contrib_scale_pptx(
            self.kpi, [c.get("contribution") for c in (self.channels or [])[:10]]
        )
        headers = ["Канал", "Бюджет", "Вклад", metric_col_hdr, "Доля эффекта", "Вердикт"]
        units =   ["",      "₽ млн",  contrib_unit_pptx,  metric_col_unit, "%",  ""]
        x = table_x
        for i, (hdr, cw) in enumerate(zip(headers, col_widths)):
            align = PP_ALIGN.LEFT if i in (0, 5) else PP_ALIGN.RIGHT
            self._text(
                slide, x + (0.05 if align == PP_ALIGN.LEFT else 0),
                table_y, cw - 0.1, 0.25,
                hdr.upper(),
                font=self.sans, size=9, bold=True, color=self.deep_100, align=align,
            )
            if units[i]:
                self._text(
                    slide, x + (0.05 if align == PP_ALIGN.LEFT else 0),
                    table_y + 0.27, cw - 0.1, 0.2, units[i],
                    font=self.sans, size=8, italic=True, color=self.deep_60, align=align,
                )
            x += cw
        # Header hairline (thick)
        self._hairline(slide, table_x, table_y + 0.55, table_w, weight=0.75, color=self.deep_100)

        # Data rows - sourced from adapter-supplied self.channels when present;
        # fallback to Kagocel pilot rows for preview / wireframe mode.
        if self.channels:
            rows = self._build_action_table_rows(self.channels)
        else:
            rows = [
                ("Digital video", "65",  "110",  "1.9", "26",  "Scale",  ""),
                ("TV",            "120", "180",  "1.5", "42",  "Scale",  "1"),
                ("Search",        "28",  "48",   "1.7", "11",  "Scale",  ""),
                ("OOH",           "35",  "52",   "1.5", "12",  "Hold",   ""),
                ("Social",        "18",  "24",   "1.3", "6",   "Watch",  "2"),
                ("Print",         "12",  "8",    "0.7", "2",   "Cut",    "3"),
                ("Radio",         "8",   "6",    "0.8", "1",   "Cut",    ""),
            ]
        verdict_colors = {
            "Scale":  (self.deep_100, True),
            "Hold":   (self.deep_60, False),
            "Watch":  (self.deep_60, False),
            "Reduce": (self.gold_muted, True),
            "Cut":    (self.gold, True),
        }
        # Stage C.3: display verdicts in Russian. Enum keys stay English
        # internally so derive_verdict() and downstream narrative helpers
        # remain unchanged.
        verdict_ru = {
            "Scale":  "Увеличить",
            "Hold":   "Держать",
            "Watch":  "Наблюдать",
            "Reduce": "Сократить",
            "Cut":    "Остановить",
        }
        row_y = table_y + 0.65
        for row in rows:
            # Phase 1.9: rows are 8-tuples (added ci_str). Backward compat for fallback
            # 7-tuple wireframe rows: pad with empty ci_str.
            if len(row) == 7:
                row = row + ("",)
            channel, budget, contrib, roi, share, verdict, footnote, ci_str = row
            # Channel name (+E2: пометка калиброванного тестом канала)
            x = table_x
            if channel in self.calibrated_channels:
                self._rich(
                    slide, x + 0.05, row_y, col_widths[0] - 0.1, 0.3,
                    runs=[
                        (channel, {"font": self.sans, "size": 11, "bold": True,
                                   "color": self.deep_100}),
                        ("  [CALIBRATED]", {"font": self.sans, "size": 7,
                                            "color": self.gold}),
                    ],
                )
            else:
                self._text(
                    slide, x + 0.05, row_y, col_widths[0] - 0.1, 0.3,
                    channel,
                    font=self.sans, size=11, bold=True, color=self.deep_100,
                )
            x += col_widths[0]
            # Budget
            self._text(
                slide, x, row_y, col_widths[1] - 0.1, 0.3, budget,
                font=self.sans, size=11, color=self.deep_100,
                align=PP_ALIGN.RIGHT,
            )
            x += col_widths[1]
            # Contribution
            self._text(
                slide, x, row_y, col_widths[2] - 0.1, 0.3, contrib,
                font=self.sans, size=11, color=self.deep_100,
                align=PP_ALIGN.RIGHT,
            )
            x += col_widths[2]
            # mROAS (with footnote and/or Phase 1.9 90% HDI bracket if present)
            mroas_runs = [(roi, {"font": self.sans, "size": 11, "color": self.deep_100})]
            if ci_str:
                mroas_runs.append(
                    (" " + ci_str, {"font": self.sans, "size": 8, "color": self.deep_60})
                )
            if footnote:
                mroas_runs.append(
                    (footnote, {"font": self.sans, "size": 7, "color": self.gold})
                )
            if len(mroas_runs) > 1:
                self._rich(
                    slide, x, row_y, col_widths[3] - 0.1, 0.3,
                    runs=mroas_runs,
                    align=PP_ALIGN.RIGHT,
                )
            else:
                self._text(
                    slide, x, row_y, col_widths[3] - 0.1, 0.3, roi,
                    font=self.sans, size=11, color=self.deep_100,
                    align=PP_ALIGN.RIGHT,
                )
            x += col_widths[3]
            # Share %
            self._text(
                slide, x, row_y, col_widths[4] - 0.1, 0.3, share,
                font=self.sans, size=11, color=self.deep_100,
                align=PP_ALIGN.RIGHT,
            )
            x += col_widths[4]
            # Verdict - fallback to Hold styling if unknown key
            vcolor, vbold = verdict_colors.get(verdict, (self.deep_60, False))
            verdict_label = verdict_ru.get(verdict, verdict)
            self._text(
                slide, x + 0.05, row_y, col_widths[5] - 0.1, 0.3, verdict_label,
                font=self.sans, size=11, bold=vbold, color=vcolor,
            )
            # Row hairline
            self._hairline(slide, table_x, row_y + 0.32, table_w, weight=0.25, color=self.deep_20)
            row_y += 0.35

        # Total row (compact) - data-driven when facts present, else Kagocel pilot
        self._text(
            slide, table_x + 0.05, row_y + 0.02, col_widths[0] - 0.1, 0.25,
            "ИТОГО",
            font=self.sans, size=10, bold=True, color=self.deep_100,
        )
        if self.facts:
            tb = self.facts.get("total_budget_mln") or 0
            tc_raw = (self.facts.get("total_contrib_mln") or 0) * 1_000_000.0
            wr = self.facts.get("weighted_roi")
            # v1.3.2: aggregate cell adapts per KPI (× / ₽/ед. inverse / 100%).
            if wr is None:
                wr_cell = "-"
            elif self.kpi["is_legacy"]:
                wr_cell = f"{wr:.2f}"
            elif self.kpi["mode"] == "effectiveness":
                wr_cell = "100"  # доли suммируются в 100% по построению
            elif self.kpi["kpi_kind"] == "count":
                try:
                    cpu = 1.0 / float(wr) if float(wr) > 0 else None
                    wr_cell = f"{cpu:.0f}" if cpu else "-"
                except (TypeError, ValueError, ZeroDivisionError):
                    wr_cell = "-"
            else:
                wr_cell = f"{wr:.2f}"
            totals = [
                f"{tb:.0f}",
                _fmt_contrib_pptx(tc_raw, contrib_scale_s07),
                wr_cell,
                "100",
            ]
        else:
            totals = ["286", "428", "1.50", "100"]
        aligns = [PP_ALIGN.RIGHT] * 4
        x = table_x + col_widths[0]
        for w, v, al in zip(col_widths[1:5], totals, aligns):
            self._text(
                slide, x, row_y + 0.02, w - 0.1, 0.25, v,
                font=self.sans, size=11, bold=True, color=self.deep_100,
                align=al,
            )
            x += w
        self._hairline(slide, table_x, row_y + 0.32, table_w, weight=0.75, color=self.deep_100)

        # Unified bottom block: ПРИМЕЧАНИЯ label + up to 3 data-driven footnotes
        # generated from flagged channels (Reduce/Cut) - matches footnote
        # numbers already assigned in _build_action_table_rows. Preview /
        # wireframe mode falls back to Kagocel pilot footnotes below.
        if self.channels:
            # Filter flagged within the same [:10] window the table rows use,
            # so bottom-block text pairs with the rendered row superscripts.
            flagged = [c for c in self.channels[:10] if c.get("verdict") in ("Reduce", "Cut")][:3]
            # v1.3.2: KPI-aware verdict reasons.
            if self.kpi["is_legacy"]:
                reason_by_verdict = {
                    "Cut": "ниже точки безубыточности по mROAS; рекомендовано остановить или перевести в другие каналы.",
                    "Reduce": "достигнуто насыщение; маржинальный возврат от дополнительного рубля ниже среднего по портфелю.",
                }
            else:
                metric_short = self.kpi["metric_short"]
                if self.kpi["mode"] == "effectiveness":
                    # M1 audit fix: «низкая доля канала» - natural phrasing
                    # для effectiveness mode (no «breakeven» metaphor).
                    reason_by_verdict = {
                        "Cut": "низкая доля в портфеле; рекомендовано остановить или перевести в другие каналы.",
                        "Reduce": "вклад канала в долю эффекта ниже среднего по портфелю.",
                    }
                elif self.kpi["kpi_kind"] == "count":
                    reason_by_verdict = {
                        "Cut": f"{metric_short} превышает ценность единицы; рекомендовано остановить или перевести в другие каналы.",
                        "Reduce": "достигнуто насыщение; стоимость следующей единицы выше среднего по портфелю.",
                    }
                else:
                    reason_by_verdict = {
                        "Cut": f"ниже точки безубыточности по {metric_short}; рекомендовано остановить или перевести в другие каналы.",
                        "Reduce": "достигнуто насыщение; маржинальный возврат ниже среднего по портфелю.",
                    }
            footnotes = [
                (str(i + 1), f"{c.get('name') or '-'}: {reason_by_verdict.get(c.get('verdict'), 'рекомендовано пересмотреть аллокацию.')}")
                for i, c in enumerate(flagged)
            ]
            if not footnotes:
                ok_metric = "mROAS" if self.kpi["is_legacy"] else self.kpi["metric_short"]
                footnotes = [("1", f"Все каналы портфеля в рабочем диапазоне {ok_metric}; критических рекомендаций нет.")]
        else:
            footnotes = [
                ("1", "TV: mROAS считается при текущих 85 TRP/нед; выше 100 TRP/нед ROI падает ниже 1.2×."),
                ("2", "Social volatile - mROAS 1.3× median, но CI 0.8 1.8× (высокая неопределённость)."),
                ("3", "Print: ниже точки безубыточности; рекомендация основана на 3-квартальном тренде."),
            ]
        line_h = 0.14
        # Layout (bottom-up): source (6.87) · fn3 (6.73) · fn2 (6.59) · fn1 (6.45) · label (6.30)
        label_y = 6.30
        self._text(
            slide, self.safe, label_y, 1.5, line_h, "ПРИМЕЧАНИЯ",
            font=self.sans, size=7, bold=True, color=self.deep_60,
        )
        fy = label_y + 0.15
        for num, txt in footnotes:
            self._rich(
                slide, self.safe, fy, 8.3, line_h,
                runs=[
                    (num, {"font": self.sans, "size": 7, "bold": True, "color": self.gold}),
                    ("  " + txt, {"font": self.sans, "size": 7, "italic": True, "color": self.deep_60}),
                ],
                line_spacing=1.0,
            )
            fy += line_h
        # Source at y=6.67 (max close to footer hairline 6.85)
        # B1-fix R-07 (2026-07-03): фактический уровень интервалов — 90% HDI
        # (utils.posterior_propagation.DEFAULT_HDI_PROB=0.9; OLS bootstrap тот же
        # compute_ci_hdi, n_boot=200). Подпись «95%» была враньём семейства F-18.
        _src_text2 = (
            f"Источник: OLS MMM Aurora AI · {self.report_id}; bootstrap 90% HDI (n=200), точечные оценки."
            if self.is_ols
            else f"Источник: Bayesian MMM Aurora AI · {self.report_id}; 90% HDI-интервалы, медианы апостериорного распределения."
        )
        self._source(slide, 6.87, text=_src_text2)

        self._footer(slide, 7 + self._page_shift)

    # ----------------------------------------------------------------
    # SLIDE 08 - ACTION + FULL TIMELINE (with annotations)
    # ----------------------------------------------------------------

    def s08_action_timeline(self):
        slide = self._blank()
        self._header(slide, slide_num=8 + self._page_shift)

        self._category(slide, self.safe, 0.60, "ДИНАМИКА")

        # Stage C.5: timeline action headline = schedule recommendation.
        title = (
            derive_action_headline(self.channels, self.facts, "timeline")
            or "Пульсирующее размещение вместо непрерывного - экономия 15-20% без потери охвата"
        )
        self._action_title(
            slide,
            title,
            show_lime=True, y=0.80, height=0.80,
        )

        chart_x = self.safe
        chart_y = 1.95
        chart_w = self.w - 2 * self.safe
        chart_h = 3.7

        # Subtitle: when real time_series is present we know exact period
        # range; otherwise fall back to data_window_label preview text.
        ts = self.time_series if isinstance(self.time_series, dict) else None
        if ts and ts.get("dates"):
            dates_list = list(ts["dates"])
            period_label = f"{dates_list[0]} - {dates_list[-1]}" if dates_list else (self.data_window_label or "—")
        else:
            period_label = self.data_window_label or "—"

        # v1.3.2: timeline axis title per KPI unit (₽ / упак).
        if self.kpi["kpi_kind"] == "count":
            timeline_axis_label = "ПРОДАЖИ ПО ПЕРИОДАМ / УПАК"
        else:
            timeline_axis_label = "ПРОДАЖИ ПО ПЕРИОДАМ / ₽"
        self._text(
            slide, chart_x, chart_y, chart_w, 0.25,
            timeline_axis_label,
            font=self.sans, size=9, bold=True, color=self.deep_80,
        )
        # Б-2 (аудит Т3+): свёртку считаем ДО подписи — подпись перечисляет
        # ФАКТИЧЕСКИЙ состав полос («База + Медиа + … · период»), а не устаревшее
        # «вклад каждого канала» (после П2 полосы = 4 верхние группы).
        ds = self.decomposition_series
        _collapsed = None
        if ds and ds.get("series"):
            from engines.decomposer import collapse_series_to_top_groups
            _collapsed = collapse_series_to_top_groups(ds["series"])
        if _collapsed:
            _tl_sub = " + ".join(c["name"] for c in _collapsed) + f" · {period_label}"
        else:
            _tl_sub = f"Базовый уровень + вклад каждого канала · {period_label}"
        self._text(
            slide, chart_x, chart_y + 0.27, chart_w, 0.22,
            _tl_sub,
            font=self.sans, size=9, italic=True, color=self.deep_60,
        )
        # Hairline removed per brand rule - minimize horizontal lines

        # ── Real stacked area chart from decomposition.time_series ────────
        # If pipeline provided per-period series, render a native python-pptx
        # AREA_STACKED chart (real data, real categorical x-axis). Otherwise
        # fall back to legacy preview/wireframe mode that paints stylized
        # rectangles week-by-week (Kagocel pilot bands).
        if (ds and ds.get("series")) or (ts and ts.get("dates") and ts.get("baseline")):
            from .charts import make_timeline_area

            chart_inner_x = chart_x
            chart_inner_y = chart_y + 0.8
            chart_inner_w = chart_w
            chart_inner_h = 2.8

            factor_series = None
            baseline_label = "Baseline"
            if _collapsed:
                # Т3-плюс (двухуровневость): обзорный timeline свёрнут в 4 верхние
                # группы (БАЗА·МЕДИА·ВНЕШНИЕ·КОНКУРЕНТЫ) — паритет с новым дефолтом
                # программы (ChannelTimeline тоже свёрнут). Детализация вкладов —
                # в waterfall (s06) и таблице каналов (s07). Тождество свёртки ==
                # исходных гарантирует collapse_series_to_top_groups (SSOT).
                tl_dates = list(ds.get("dates") or (ts or {}).get("dates") or [])
                _c_base = next((c for c in _collapsed if c["top_group"] == "БАЗА"), None)
                _c_media = next((c for c in _collapsed if c["top_group"] == "МЕДИА"), None)
                tl_baseline = (
                    list(_c_base["data"]) if _c_base
                    else list((ts or {}).get("baseline") or [])
                )
                if _c_base:
                    baseline_label = _c_base["name"]  # «База» — русская легенда (Б-2)
                channel_dict = {_c_media["name"]: list(_c_media["data"])} if _c_media else {}
                # ВНЕШНИЕ/КОНКУРЕНТЫ — вынесенными полосами с explicit-цветом группы
                # (зеркалит GROUP_COLORS фронта: amber / red).
                _GROUP_RGB = {"ВНЕШНИЕ ФАКТОРЫ": "F59E0B", "КОНКУРЕНТЫ": "DC2626"}
                factor_series = [
                    {"name": c["name"], "type": None, "rgb": _GROUP_RGB[c["top_group"]],
                     "side": c["side"], "data": list(c["data"])}
                    for c in _collapsed if c["top_group"] in _GROUP_RGB
                ]
            else:
                # Legacy fallback (нет decomposition_series): медиа-only как раньше.
                tl_dates = list(ts["dates"])
                tl_baseline = list(ts["baseline"])
                channel_series = ts.get("channels") or {}
                ranked = sorted(
                    channel_series.items(),
                    key=lambda kv: sum(float(v) for v in kv[1] or []),
                    reverse=True,
                )[:6]
                channel_dict = {name: list(values) for name, values in ranked}

            make_timeline_area(
                slide,
                chart_inner_x, chart_inner_y, chart_inner_w, chart_inner_h,
                dates=tl_dates,
                baseline=tl_baseline,
                channel_series=channel_dict,
                factor_series=factor_series,
                baseline_label=baseline_label,
            )

            # Source at bottom (unified position, real-data variant)
            _engine_decomp = "OLS MMM" if self.is_ols else "Bayesian MMM"
            self._source(
                slide, 6.87,
                text=f"Источник: {self.sources_client_label}, продажи за период {period_label}; декомпозиция {_engine_decomp} · {self.report_id}",
            )

            self._footer(slide, 8 + self._page_shift)
            return

        # ── Legacy preview/wireframe path (no real time_series) ───────────
        weeks = 13
        area_x = chart_x + 0.7  # leave space for y-axis labels
        area_y = chart_y + 0.8
        area_w = chart_w - 2.8  # leave space for legend on right
        area_h = 2.6

        # Bands - data-driven from top contributors when channels present,
        # else Kagocel pilot bands (preview / wireframe mode).
        # Target peak-sum under area_h (2.6") accounting for leader seasonal
        # multiplier up to 1.6×.
        palette = [self.gold, self.deep_80, self.deep_60, self.deep_40, self.deep_20, self.deep_20]
        if self.channels:
            by_contrib = sorted(
                self.channels,
                key=lambda c: float(c.get("contribution") or 0),
                reverse=True,
            )[:5]
            total_c = sum(float(c.get("contribution") or 0) for c in by_contrib) or 1.0
            # Channel bands occupy ~55% of total height; baseline ~45%.
            channel_h_budget = 1.10  # inches, fits sum × peak mod ≤ area_h
            baseline_h = 0.85
            bands = [("Базовый уровень", self.deep_40, baseline_h)]
            for i, c in enumerate(by_contrib):
                share = float(c.get("contribution") or 0) / total_c
                bands.append((
                    c.get("name") or "-",
                    palette[i % len(palette)],
                    max(0.05, share * channel_h_budget),
                ))
            # Leader name = first channel after baseline (highest contribution)
            leader_name = bands[1][0] if len(bands) > 1 else None
        else:
            bands = [
                ("Базовый уровень", self.deep_40, 0.85),
                ("TV",            self.gold,    0.55),
                ("Digital video", self.deep_80, 0.35),
                ("Search",        self.deep_60, 0.17),
                ("OOH",           self.deep_20, 0.12),
                ("Social",        self.deep_20, 0.07),
            ]
            leader_name = "TV"
        seg_w = area_w / weeks

        # Draw stacked with seasonal modulation.
        # Preview mode retains W06/W11 TV flights; real-data mode uses generic
        # sinusoidal seasonality without synthetic flight spikes.
        preview_mode = not self.channels
        for w_idx in range(weeks):
            week_x = area_x + w_idx * seg_w
            base_mod = 0.9 + 0.15 * math.sin(w_idx / 13 * 2 * math.pi)
            if preview_mode and w_idx in (5, 10):
                leader_mod = 1.6
            elif preview_mode:
                leader_mod = 0.85 if w_idx in (4, 6, 9, 11) else 0.7
            else:
                leader_mod = 0.95 + 0.15 * math.sin(w_idx / 4.2 + 0.7)

            y_top = area_y + area_h
            for name, col, base_h in bands:
                if name == leader_name:
                    h_band = base_h * leader_mod
                elif name == "Baseline":
                    h_band = base_h * base_mod
                else:
                    # zlib.crc32 is deterministic across processes (unlike
                    # built-in hash() which is salted with PYTHONHASHSEED).
                    # Keeps band modulation phase stable between builds so
                    # slide XML content is reproducible for diffing / review.
                    # (Outer PPTX SHA still varies due to python-pptx save-time
                    # timestamps in docProps/core.xml - that is outside scope.)
                    jitter = zlib.crc32(name.encode("utf-8")) % 7
                    h_band = base_h * (0.95 + 0.1 * math.sin(w_idx / 5 + jitter))
                y_top -= h_band
                self._rect(slide, week_x, y_top, seg_w, h_band, fill=col)

        # Y-axis label ticks (3 values)
        for i, (v, y_pos) in enumerate([(5, 0.1), (3, area_h / 2), (0, area_h - 0.1)]):
            self._text(
                slide, chart_x, area_y + y_pos - 0.1, 0.6, 0.2,
                str(v), font=self.sans, size=8, color=self.deep_60,
                align=PP_ALIGN.RIGHT,
            )
        self._vline(slide, area_x, area_y, area_h, weight=0.5, color=self.deep_60)

        # X-axis labels
        self._hairline(slide, area_x, area_y + area_h, area_w, weight=0.5, color=self.deep_60)
        for wk in [1, 3, 5, 7, 9, 11, 13]:
            xp = area_x + (wk - 1) * seg_w
            self._text(
                slide, xp - 0.25, area_y + area_h + 0.08, 0.5, 0.18,
                f"W{wk:02d}",
                font=self.sans, size=7, color=self.deep_60, align=PP_ALIGN.CENTER,
            )

        # Annotations - only in preview mode (TV flight narrative is Kagocel-
        # specific). Real-data runs show clean bands without flight callouts;
        # per-week peak detection is deferred to XLSX "Динамика" sheet.
        if preview_mode:
            w06_x = area_x + 5 * seg_w + seg_w / 2
            ann_y_top = area_y - 0.3
            self._text(
                slide, w06_x - 0.9, ann_y_top, 1.8, 0.2,
                "TV FLIGHT . 95 TRP/нед",
                font=self.sans, size=7, bold=True, color=self.gold, align=PP_ALIGN.CENTER,
            )
            self._vline(slide, w06_x, area_y, 0.08, weight=1.5, color=self.gold)

            w11_x = area_x + 10 * seg_w + seg_w / 2
            self._text(
                slide, w11_x - 0.9, ann_y_top, 1.8, 0.2,
                "HOLIDAY PUSH . DIGITAL",
                font=self.sans, size=7, bold=True, color=self.gold, align=PP_ALIGN.CENTER,
            )
            self._vline(slide, w11_x, area_y, 0.08, weight=1.5, color=self.gold)

        # Legend on the right
        leg_x = area_x + area_w + 0.3
        leg_y = area_y + 0.1
        leg_w = chart_x + chart_w - leg_x
        self._text(
            slide, leg_x, leg_y - 0.3, leg_w, 0.2, "КАНАЛЫ",
            font=self.sans, size=8, bold=True, color=self.deep_80,
        )
        self._hairline(slide, leg_x, leg_y - 0.05, 0.8, weight=0.5, color=self.rule_color)
        for i, (name, col, _) in enumerate(bands):
            sw = self._rect(
                slide, leg_x, leg_y + i * 0.35, 0.15, 0.15, fill=col,
            )
            self._text(
                slide, leg_x + 0.25, leg_y + i * 0.35 - 0.02, leg_w - 0.3, 0.2, name,
                font=self.sans, size=9, color=self.deep_100,
            )

        # Source at bottom (unified position)
        # B1-fix R-02: период в подписи — timeline уже показывает реальные даты;
        # используем ту же period_label-логику (реальные даты > мета > «—»).
        _engine_decomp2 = "OLS MMM" if self.is_ols else "Bayesian MMM"
        _tl_period = period_label if period_label and period_label != "—" else None
        _tl_period_part = f", продажи за период {_tl_period}" if _tl_period else ""
        self._source(
            slide, 6.87,
            text=f"Источник: {self.sources_client_label}{_tl_period_part}; декомпозиция {_engine_decomp2} · {self.report_id}",
        )

        self._footer(slide, 8 + self._page_shift)

    # ----------------------------------------------------------------
    # SLIDE 09 - EXECUTIVE SUMMARY (SCQAR)
    # ----------------------------------------------------------------

    def s09_scqar(self):
        slide = self._blank()
        self._header(slide, slide_num=5)

        self._category(slide, self.safe, 0.60, "РЕКОМЕНДАЦИЯ")

        # Build SCQAR blocks from facts when present, else Kagocel pilot.
        if self.facts and self.channels:
            f = self.facts
            leader = f.get("leader_channel") or "Лидер"
            hero = f.get("hero_channel") or leader
            tb = f.get("total_budget_mln") or 0
            n_ch = f.get("n_active_channels") or len(self.channels)
            wr = f.get("weighted_roi")
            leader_spend_pct = f.get("leader_share_spend_pct")
            realloc = f.get("reallocation_mln") or 0
            lift = f.get("expected_lift_pct")
            # Hero mROAS
            hero_ch = next((c for c in self.channels if c.get("name") == hero), {})
            hero_m = float(hero_ch.get("mroas") or 0)
            # Underperformer names. L23 fix (2026-04-29): dedup от cut_source -
            # narrative_adapter уже исключил cut_source из underperformer_names
            # в facts dict, но self.channels - raw merged list. Apply here
            # locally for PPTX consistency (avoid «из TRPs ... остановить TRPs»).
            cut_source_local = (self.facts or {}).get("cut_source_channel")
            underperf = [
                c.get("name") for c in self.channels
                if c.get("verdict") in ("Cut", "Reduce") and c.get("name") != cut_source_local
            ]
            underperf_str = ", ".join(underperf) if underperf else "отстающие каналы"

            # Stage C.5: McKinsey 3-scenario SCQAR action headline.
            # narrative_adapter.derive_action_headline handles:
            #   - Risk (all underperf) → "Сократить X и сфокусировать на Y"
            #   - Rebalance (lift ≥ 0.5pp) → "Перераспределить N млн ₽ в Y - +X пп к ROAS"
            #   - Hold + control (weak lift) → "Портфель сбалансирован - A/B тест"
            action_title = (
                derive_action_headline(self.channels, self.facts, "scqar")
                or f"Пересмотреть аллокацию {leader}"
            )

            # v1.3.2: KPI-aware portfolio metric phrase в situation.
            if wr is not None:
                if self.kpi["is_legacy"]:
                    wr_segment = f"Средневзвешенный ROI {wr:.1f}×, "
                else:
                    wr_segment = f"{_weighted_summary_phrase_pptx(wr, self.kpi)}, "
            else:
                wr_segment = ""
            # B1-fix R-02: «в квартал» — wireframe-предположение о периодичности;
            # честно: бюджет за период данных. MQS через _mstr (None → «н/д»).
            _period_phrase = (
                f"за период {self.data_window_label}" if self.data_window_label
                else "за анализируемый период"
            )
            # Склонение: 1 канал / 2-4 канала / 5+ каналов.
            _ch_word = ("канал" if n_ch % 10 == 1 and n_ch % 100 != 11
                        else "канала" if n_ch % 10 in (2, 3, 4) and n_ch % 100 not in (12, 13, 14)
                        else "каналов")
            _ch_adj = ("активный" if _ch_word == "канал"
                       else "активных")
            situation_body = (
                f"{self.client} размещает {tb:.0f} млн ₽ {_period_phrase} через {n_ch} {_ch_adj} {_ch_word}. "
                + wr_segment
                + f"MQS модели {self._mstr(self.mqs_score, '{:.0f}')}/100."
            )

            # L14 (math-fix v1.4 Section C, 2026-04-29): use budget_dominator
            # вместо leader (top contribution). budget_dominator = max spend.
            budget_dom = f.get("budget_dominator_channel") or leader
            bd_spend_pct = f.get("budget_dominator_spend_pct") or leader_spend_pct
            bd_contrib_pct = f.get("budget_dominator_contrib_pct") or 0.0
            complication_parts = []
            if budget_dom and bd_spend_pct is not None and abs((bd_spend_pct or 0) - (bd_contrib_pct or 0)) >= 5.0:
                complication_parts.append(
                    f"{budget_dom} занимает {_fmt_pct(bd_spend_pct)} бюджета, "
                    f"но даёт {_fmt_pct(bd_contrib_pct)} эффекта"
                )
            if hero != leader and hero_m >= 1.0:
                # v1.3.2: KPI-aware metric label в complication.
                if self.kpi["is_legacy"]:
                    complication_parts.append(f"по mROAS {hero} опережает ({hero_m:.1f}×)")
                else:
                    metric_short = self.kpi["metric_short"]
                    metric_fmt = _fmt_metric_pptx(hero_m, self.kpi)
                    complication_parts.append(f"по {metric_short} {hero} опережает ({metric_fmt})")
            if underperf:
                complication_parts.append(f"{underperf_str} тянут портфель вниз")
            # B1-fix R-13 (2026-07-03): хвост «Портфель требует перебалансировки»
            # противоречил ответу «Сохранить аллокацию» при lift≈0 (Kagocel-зонд:
            # ПРОБЛЕМА требует, ОТВЕТ сохраняет). Хвост — из фактов оптимизатора.
            _lift_f = f.get("expected_lift_pct")
            _realloc_potential = bool(
                (f.get("cut_source_channel") or f.get("scale_destination_channel"))
                and (f.get("reallocation_mln") or 0) >= 1
            )
            _has_real_lift = _lift_f is not None and _lift_f >= 0.5
            if _realloc_potential and _has_real_lift:
                _compl_tail = "Портфель требует перебалансировки."
            elif f.get("binding_constraints"):
                _compl_tail = ("Однако оптимизатор упёрся в заданные границы каналов - "
                               "потенциал перераспределения ограничен настройками.")
            else:
                _compl_tail = ("Однако перераспределение в заданных границах не даёт "
                               "ощутимого прироста - разрыв объясняется насыщением каналов.")
            complication_body = (
                (". ".join(complication_parts) + ". " + _compl_tail)
                if complication_parts else _compl_tail
            )

            # L15 (math-fix v1.4 Section C, 2026-04-29): use cut_source /
            # scale_destination from action_summary вместо leader/hero.
            cut_source = f.get("cut_source_channel")
            scale_dest = f.get("scale_destination_channel")
            answer_parts = []
            if cut_source and scale_dest and realloc >= 1:
                answer_parts.append(f"Перераспределить {realloc:.0f} млн ₽ из {cut_source} в {scale_dest}")
            elif scale_dest and realloc >= 1:
                answer_parts.append(f"Нарастить {scale_dest} на ~{realloc:.0f} млн ₽")
            elif cut_source and realloc >= 1:
                answer_parts.append(f"Сократить {cut_source} ({realloc:.0f} млн ₽)")
            if underperf:
                answer_parts.append(f"остановить {underperf_str}")
            answer_body = "; ".join(answer_parts) + "." if answer_parts else f"Сохранить текущую аллокацию по {leader} с контролем насыщения."

            # B1-fix R-13: вопрос согласован со сценарием — при hold-исходе
            # «как перераспределить» обещал то, чего ответ не даёт.
            _question_body = (
                "Как перераспределить бюджет, чтобы поднять ROAS без снижения охвата знания?"
                if (_realloc_potential and _has_real_lift)
                else "Есть ли в текущем миксе резерв эффективности, который стоит задействовать?"
            )
            blocks = [
                {"label": "СИТУАЦИЯ", "height": 0.6, "body": situation_body},
                {"label": "ПРОБЛЕМА", "height": 0.8, "body": complication_body},
                {"label": "ВОПРОС", "height": 0.55,
                 "body": _question_body,
                 "accent": True},
                {"label": "ОТВЕТ", "height": 0.6, "body": answer_body},
                {"label": "РЕКОМЕНДАЦИИ", "height": 2.3, "body": None},
            ]
        else:
            action_title = "Сократить TV до 60 TRP/нед и удвоить Digital Video"
            blocks = [
                {"label":  "СИТУАЦИЯ",
                 "height": 0.6,
                 "body":   f"{self.client} размещает 286 млн ₽ в квартал через 5 активных каналов. Средневзвешенный ROI 1.5×, MQS модели {self._mstr(self.mqs_score, '{:.0f}')}/100."},
                {"label":  "ПРОБЛЕМА",
                 "height": 0.8,
                 "body":   "TV достиг насыщения: выше 80 TRP/нед маржинальный ROI падает на 22% год к году. Digital video недоинвестирован (<15% бюджета при mROAS 1.9×). Портфель не оптимизирован."},
                {"label":  "ВОПРОС",
                 "height": 0.55,
                 "body":   "Как перераспределить бюджет, чтобы поднять ROAS без снижения охвата знания?",
                 "accent": True},
                {"label":  "ОТВЕТ",
                 "height": 0.6,
                 "body":   "Сократить TV с 120 до 90 млн ₽/квартал, увеличить Digital video с 65 до 100 млн, сохранить Search/OOH, остановить Print и Radio."},
                {"label":  "РЕКОМЕНДАЦИИ",
                 "height": 2.3,  # Stage C.4: was 2.0; 3 actions need more room
                 "body":   None},
            ]

        self._action_title(
            slide,
            action_title,
            show_lime=True, y=0.80, height=0.85,
        )

        y = 2.3
        for b in blocks:
            # Gold bar
            accent = b.get("accent", False)
            bar_color = self.gold if accent else self.deep_40
            self._vbar(slide, self.safe, y + 0.05, b["height"] - 0.1, weight=3 if accent else 2, color=bar_color)
            # Label
            self._text(
                slide, self.safe + 0.2, y, 2.0, 0.28, b["label"],
                font=self.sans, size=9, bold=True,
                color=self.gold if accent else self.deep_60,
            )
            # Body (Stage C.4: font 12→11 reduces wrap → no overlap between blocks)
            if b["body"]:
                self._text(
                    slide, self.safe + 2.3, y, 7.5, b["height"] - 0.1,
                    b["body"],
                    font=self.sans, size=11,
                    italic=accent,
                    color=self.deep_100,
                    line_spacing=1.3,
                )
            else:
                # Recommendation - 3 templated actions when facts present.
                # L15 audit fix (2026-04-29): Action 01 uses cut_source/scale_dest
                # from action_summary (NOT leader/hero) - same fix как s09 SCQAR.
                if self.facts and self.channels:
                    f = self.facts
                    leader = f.get("leader_channel") or "Лидер"
                    hero = f.get("hero_channel") or leader
                    realloc = f.get("reallocation_mln") or 0
                    lift = f.get("expected_lift_pct")
                    underperf = [c.get("name") for c in self.channels if c.get("verdict") in ("Cut",)]
                    cut_source = f.get("cut_source_channel")
                    scale_dest = f.get("scale_destination_channel")

                    if cut_source and scale_dest and realloc >= 1:
                        action_01_body = (
                            f" {realloc:.0f} млн ₽ из {cut_source} в {scale_dest}. "
                            "Отложенный эффект (adstock) компенсирует краткосрочный спад охвата."
                        )
                    elif scale_dest and realloc >= 1:
                        action_01_body = (
                            f" Нарастить {scale_dest} на ~{realloc:.0f} млн ₽ – "  # П8-1
                            "за счёт roll-over бюджета или дополнительных средств."
                        )
                    elif cut_source and realloc >= 1:
                        action_01_body = (
                            f" Сократить {cut_source} ({realloc:.0f} млн ₽) – "  # П8-1
                            "текущая аллокация неэффективна."
                        )
                    elif hero != leader and realloc >= 1:
                        # Legacy fallback (cut_source/scale_dest unavailable)
                        action_01_body = (
                            f" {realloc:.0f} млн ₽ из {leader} в {hero}. "
                            "Отложенный эффект (adstock) компенсирует краткосрочный спад охвата."
                        )
                    else:
                        action_01_body = f" Сохранить аллокацию по {leader} с контролем индикаторов насыщения."

                    # B1-fix R-09 (2026-07-03): рекомендации 02/03 были выдуманы
                    # («пульсирующее размещение - экономия 15-20%», «целевой
                    # ретаргетинг по сегментам») — модель не вычисляет ни
                    # flighting, ни сегменты; числа фабриковались. Честные
                    # рекомендации — только выводимые из модели: контроль
                    # насыщения (Hill-механика), снятие неопределённости
                    # вердиктов, сверка прогноза с фактом при следующей волне.
                    _uncertain_names = [
                        c.get("name") for c in self.channels
                        if c.get("verdict") == "Uncertain" and c.get("name")
                    ]
                    if _uncertain_names:
                        _action_02 = (
                            "Снять неопределённость вердиктов.",
                            f" {', '.join(_uncertain_names[:3])}: интервалы эффективности "
                            "широки - добавить историю данных или вариацию трат, затем переобучить модель.",
                        )
                    else:
                        _action_02 = (
                            f"Контролировать насыщение {leader}.",
                            " Наращивание трат сверх текущего уровня даёт убывающую отдачу "
                            "(Hill-насыщение) - отслеживать mROAS при изменениях бюджета.",
                        )
                    # Пласт 2 (2026-07-11): KPI-aware lift fragment — для count/effectiveness «ROAS» не применим.
                    _lift_frag = (
                        f" {_lift_phrase_pptx(lift, self.kpi)} - проверить фактом."
                        if (lift is not None and lift >= 0.5)
                        else " Проверить, что фактические продажи соответствуют прогнозу модели."
                    )
                    actions = [
                        ("01", "Перераспределить бюджет.", action_01_body),
                        ("02", *_action_02),
                        ("03",
                         "Сверить прогноз с фактом после следующего периода.",
                         f" Обновить модель новыми данными и подтвердить рекомендации.{_lift_frag}"),
                    ]
                else:
                    actions = [
                        ("01",
                         "Перераспределить бюджет.",
                         " 25 млн ₽ из TV в Digital video. Экономия при сохранении инкрементальных продаж (отложенный эффект компенсирует)."),
                        ("02",
                         "Недельные пульсы вместо непрерывного размещения.",
                         " TV флайты по 60 TRP × 3 недели с паузами. Экономия 18%, охват устойчив."),
                        ("03",
                         "Целевой ретаргетинг W25-54 через CTV/OLV.",
                         " Сегмент с априорным ROI 2.1×; сейчас недоинвестирован. Потенциал +12 млн ₽ инкрементальных продаж."),
                    ]
                ay = y
                for num, lead, body in actions:
                    # Number
                    self._text(
                        slide, self.safe + 2.3, ay, 0.4, 0.3, num,
                        font=self.serif, size=13, color=self.gold,
                    )
                    # Lead (bold) on first line
                    self._text(
                        slide, self.safe + 2.8, ay, 7.0, 0.22, lead,
                        font=self.sans, size=11, bold=True, color=self.deep_100,
                    )
                    # Body (regular) on next line
                    self._text(
                        slide, self.safe + 2.8, ay + 0.25, 7.0, 0.35, body.strip(),
                        font=self.sans, size=11, color=self.deep_80, line_spacing=1.25,
                    )
                    ay += 0.65
            y += b["height"]

        # Expected impact box (bottom right)
        impact_x = 10.5
        impact_y = 5.2
        self._text(
            slide, impact_x, impact_y, 2.5, 0.22, "ОЖИДАЕМЫЙ ЭФФЕКТ",
            font=self.sans, size=8, bold=True, color=self.gold,
        )
        self._hairline(slide, impact_x, impact_y + 0.27, 1.2, weight=0.75, color=self.gold)
        # 2026-05-24 fix: do not fabricate `+12 пп` placeholder when optimize result
        # missing — это customer-facing slide, любое hardcoded numerical claim — INV-50
        # violation candidate + customer confusion. Sprint Buffer #44 (2026-05-23):
        # semantic «Не определён» вместо em-dash для consultant context — em-dash
        # читается как «забыли заполнить», текст явно signals «не посчитано / нет данных».
        # Font size scaled down (42→22pt) чтобы фраза влезла в 2.5" box без переноса.
        # B1-fix R-09/R-13: «+0 пп» — пустое обещание; незначимый lift (<0.5)
        # подписываем честно, не нулевым числом.
        _lift_box = self.facts.get("expected_lift_pct") if self.facts else None
        if _lift_box is not None and _lift_box >= 0.5:
            impact_num = f"+{_lift_box:.0f} пп"
            impact_size = 42
        elif _lift_box is not None:
            impact_num = "Незначим (<0.5 пп)"
            impact_size = 18
        else:
            impact_num = "Не определён"
            impact_size = 22
        self._text(
            slide, impact_x, impact_y + 0.35, 2.5, 0.7, impact_num,
            font=self.serif, size=impact_size, color=self.deep_100,
        )
        # v1.3.2: KPI-aware impact label.
        if self.kpi["mode"] == "effectiveness":
            impact_label_text = "Прогнозный прирост доли"
        elif self.kpi["kpi_kind"] == "count":
            impact_label_text = "Прогнозный прирост продаж"
        else:
            impact_label_text = "Прогнозный ROAS"
        self._text(
            slide, impact_x, impact_y + 1.1, 2.5, 0.22, impact_label_text,
            font=self.sans, size=10, italic=True, color=self.deep_60,
        )

        self._footer(slide, 5)

    # ----------------------------------------------------------------
    # SLIDE 10 - METHODOLOGY + LIMITATIONS
    # ----------------------------------------------------------------

    # C.6.3: s10 methodology content now at physical page 11 (divider at 10).
    def s10_methodology(self):
        slide = self._blank()
        self._header(slide, slide_num=9 + self._page_shift)

        # B4-2: плашка секции вместо отдельного слайда-дивайдера.
        self._section_intro(
            slide, 3, "Методология",
            "Байесовская MMM с отложенным эффектом и Hill-насыщением – прозрачная модель с интервалами доверия",  # П8-1
        )

        _title_text = (
            "OLS MMM с отложенным эффектом (adstock) + Hill-насыщением (closed-form + bootstrap CI)"  # П8-2
            if self.is_ols
            else "Байесовская MMM с отложенным эффектом (adstock) + Hill-насыщением"  # П8-2
        )
        self._action_title(
            slide,
            _title_text,
            show_lime=True, y=0.80, height=0.80,
        )

        # Two columns
        left_x = self.safe
        left_y = 1.85
        left_w = (self.w - 2 * self.safe) * 0.48

        # LEFT: Formula card
        self._text(
            slide, left_x, left_y, left_w, 0.25, "СПЕЦИФИКАЦИЯ",
            font=self.sans, size=9, bold=True, color=self.gold,
        )
        self._hairline(slide, left_x, left_y + 0.28, 1.0, weight=0.75, color=self.gold)

        self._rect(slide, left_x, left_y + 0.45, left_w, 2.3, fill=self.bg_quiet)
        if self.is_ols:
            formulas = [
                "y_t = baseline_t + Σ β_i · sat(adstock(x_i,t)) + ε_t",
                "",
                "adstock(x, θ) = x_t + θ·x_{t-1} + θ²·x_{t-2} + …",
                "   θ = 0.5 (fixed, geometric decay)",
                "",
                "sat(z, α, γ) = z^α / (z^α + γ^α)",
                "   Hill function (fixed): α=1.5, γ=0.5",
                "",
                "β_i estimated via closed-form OLS",
                "ε_t ~ Normal(0, σ) (frequentist)",
                # B1-fix R-07-семейство: фактический n_boot=200 (ols_bootstrap.py),
                # «n=1000» был враньём того же класса, что и conformal-глоссарий (F-18).
                "CI на ROI/mROAS - bootstrap n=200",
            ]
        else:
            formulas = [
                "y_t = baseline_t + Σ β_i · sat(adstock(x_i,t)) + ε_t",
                "",
                "adstock(x, θ) = x_t + θ·x_{t-1} + θ²·x_{t-2} + …",
                "   θ_i ∈ [0, 0.95]  (geometric decay)",
                "",
                "sat(z, α, γ) = z^α / (z^α + γ^α)",
                "   Hill function, half-max = γ_i",
                "",
                # B1-fix R-08: конкретная σ приора врала (0.5 устарел 2026-04-19;
                # факт: 0.3 базовый / 0.7-0.3 hierarchical по категории канала).
                # Формулировка без числа честна во всех конфигурациях.
                "β_i ~ HalfNormal (слабоинформативный) · CPP-normalized",
                "ε_t ~ Normal(0, σ)",
            ]
        self._paragraphs(
            slide, left_x + 0.25, left_y + 0.6, left_w - 0.5, 2.1,
            [(line, {
                "font": self.mono, "size": 10,
                "color": self.deep_100 if line.strip() else self.deep_60,
            }) for line in formulas],
            line_spacing=1.25,
        )

        # LEFT bottom: diagnostics mini
        diag_y = left_y + 3.0
        self._text(
            slide, left_x, diag_y, left_w, 0.25, "ДИАГНОСТИКА",
            font=self.sans, size=9, bold=True, color=self.gold,
        )
        self._hairline(slide, left_x, diag_y + 0.28, 1.0, weight=0.75, color=self.gold)

        # B1-fix R-01: метрики через _mstr (None → «н/д», не wireframe-числа);
        # ESS может прийти дробным (min bulk/tail из arviz) — показываем целым.
        _ess_str = self._mstr(
            int(round(self.ess_min)) if isinstance(self.ess_min, (int, float)) else None,
            "{:,}").replace(",", " ")
        if self.is_ols:
            # OLS не имеет MCMC diagnostics. Показываем frequentist метрики.
            diag = [
                ("R²",                self._mstr(self.r_squared, "{:.3f}")),
                ("MAPE",              self._mstr(self.mape_pct, "{:.1f}%")),
                ("Метод",             "closed-form OLS"),
                ("CI",                "bootstrap n=200"),  # факт ols_bootstrap.py (был n=1000)
            ]
        else:
            diag = [
                ("R²",                self._mstr(self.r_squared, "{:.3f}")),
                ("MAPE",              self._mstr(self.mape_pct, "{:.1f}%")),
                ("R-hat (max)",       self._mstr(self.r_hat_max, "{:.3f}")),
                ("ESS (min)",         _ess_str),
            ]
        dy = diag_y + 0.4
        for label, val in diag:
            self._text(
                slide, left_x, dy, left_w * 0.55, 0.25, label,
                font=self.sans, size=10, color=self.deep_60,
            )
            self._text(
                slide, left_x + left_w * 0.55, dy, left_w * 0.45, 0.25, val,
                font=self.sans, size=10, bold=True, color=self.deep_100,
                align=PP_ALIGN.RIGHT,
            )
            self._hairline(slide, left_x, dy + 0.27, left_w, weight=0.25)
            dy += 0.3

        # RIGHT: Limitations (tier-1 differentiator)
        right_x = left_x + left_w + 0.5
        right_w = self.w - self.safe - right_x

        self._text(
            slide, right_x, left_y, right_w, 0.25, "ЧТО МОДЕЛЬ НЕ УЧИТЫВАЕТ",
            font=self.sans, size=9, bold=True, color=self.gold,
        )
        self._hairline(slide, right_x, left_y + 0.28, 1.2, weight=0.75, color=self.gold)

        limits = [
            ("Долгосрочные бренд-эффекты (>26 недель).",
             "Модель учитывает краткосрочный и среднесрочный эффект, но не долгосрочное строительство бренда."),
            ("Каннибализация между категориями.",
             f"Если у клиента несколько SKU - модель считает их единым KPI."),
            ("Вариация качества креативов.",
             "Влияние качества роликов на ROI не моделируется - предполагается постоянным."),
            ("Конкурентная медиа-активность.",
             "Доля голоса (SoV) конкурентов в модели отсутствует - оценка справедлива для стабильной категории."),
            ("Макроэкономические шоки.",
             "Экстремальные события (валютные скачки, регуляторные изменения) вне области анализа."),
        ]
        ly = left_y + 0.5
        for lead, body in limits:
            # Lead (bold) on first line
            self._text(
                slide, right_x + 0.2, ly, right_w - 0.2, 0.22, lead,
                font=self.sans, size=11, bold=True, color=self.deep_100,
            )
            # Body (regular) on next line
            self._text(
                slide, right_x + 0.2, ly + 0.27, right_w - 0.2, 0.45, body,
                font=self.sans, size=10, color=self.deep_60, line_spacing=1.3,
            )
            # Bullet dot (gold square, aligned with lead)
            self._rect(slide, right_x, ly + 0.07, 0.08, 0.08, fill=self.deep_60)
            ly += 0.80

        # Trust Level 3 (v1.1.0): brand vs performance disclosure если активен.
        # Speaker note + bottom note - minimal disclosure без layout disruption.
        # HTML report имеет full block в methodology; здесь kompakt single-line note.
        hier = (self.data.get('diagnostics') or {}).get('hierarchical') or {}
        if hier.get('enabled'):
            cats = hier.get('channel_categories') or {}
            n_brand = sum(1 for v in cats.values() if v == 'brand')
            n_perf = sum(1 for v in cats.values() if v == 'performance')
            t3_text = (
                f"v1.1.0: Brand vs Performance split - {n_brand} brand, {n_perf} performance каналов. "
                f"Hierarchical priors разделяют long-decay (бренд ~12 нед) и short-decay (perf ~1-2 нед)."
            )
            self._source(slide, 6.65, text=t3_text)
            # B1-fix R-06: «12+ FMCG-проектов Aurora» — недоказуемое заявление
            # (INV-50); честная формулировка о слабоинформативных приорах.
            _bottom_note = (
                "Параметры отложенного эффекта (adstock)/насыщения: индустриальные бенчмарки OLS MMM, фиксированные."  # П8-2
                if self.is_ols
                else "Приоры: слабоинформативные, на основе индустриальных бенчмарков Bayesian MMM."
            )
            self._source(slide, 6.97, text=_bottom_note)
        else:
            # Bottom note (concise: ≤100 chars fits 1 line at 7pt in 8.3" column)
            _bottom_note2 = (
                "Параметры отложенного эффекта (adstock)/насыщения: индустриальные бенчмарки OLS MMM, фиксированные."  # П8-2
                if self.is_ols
                else "Приоры: слабоинформативные, на основе индустриальных бенчмарков Bayesian MMM."
            )
            self._source(slide, 6.87, text=_bottom_note2)

        self._footer(slide, 9 + self._page_shift)

    # ----------------------------------------------------------------
    # SLIDE 11 - SOURCES + MQS
    # ----------------------------------------------------------------

    # C.6.3: s11 sources content now at physical page 13 (divider at 12).
    def s11_sources(self):
        slide = self._blank()
        self._header(slide, slide_num=10 + self._page_shift)

        # B4-2: плашка секции вместо отдельного слайда-дивайдера
        # (takeaway — из фактов, как в бывшем s_divider_data).
        _tb_mln = self.facts.get("total_budget_mln") if self.facts else None
        _mqs_part = f"MQS модели {self._mstr(self.mqs_score, '{:.0f}')}/100"
        _window_part = (
            f"Данные охватывают {self.data_window_label}" if self.data_window_label
            else "Данные загружены из файла клиента"
        )
        _tk_data = (
            f"{_window_part}, бюджет {_tb_mln:.0f} млн руб - {_mqs_part}"
            if _tb_mln else f"{_window_part} - {_mqs_part}"
        )
        self._section_intro(slide, 4, "Данные и качество", _tk_data)

        self._text(
            slide, self.safe, 0.70, 10.0, 0.7, "Источники и качество",
            font=self.serif, size=32, color=self.deep_100,
        )
        self._hairline(slide, self.safe, 1.50, 1.2, weight=0.75, color=self.gold)

        # MQS big card (left 40%) - spacious layout
        card_x = self.safe
        card_y = 1.85
        card_w = 4.5
        card_h = 3.9
        self._rect(slide, card_x, card_y, card_w, card_h, fill=self.bg_quiet)
        # Top gold stripe
        self._rect(slide, card_x, card_y, card_w, 0.06, fill=self.gold)

        self._text(
            slide, card_x + 0.35, card_y + 0.3, card_w - 0.7, 0.25,
            "ОЦЕНКА КАЧЕСТВА МОДЕЛИ (MQS)",
            font=self.sans, size=10, bold=True, color=self.gold,
        )

        # The big MQS score + /100 pair - centered vertically between
        # title (top) and status hairline (bottom). Spacing balanced so
        # "70 → /" gap ≈ "/ → 100" gap (visually symmetric pair).
        self._text(
            slide, card_x + 0.45, card_y + 0.30, 2.0, 1.8,
            self._mstr(self.mqs_score, "{:.0f}", na="—"),
            font=self.serif, size=120, color=self.deep_100, align=PP_ALIGN.RIGHT,
        )
        # "/ 100" left-aligned, gap to "70" matches gap "/" → "100"
        self._text(
            slide, card_x + 2.55, card_y + 1.10, 1.5, 0.5, "/ 100",
            font=self.serif, size=32, color=self.deep_60,
        )

        # Status - below number with clear gap
        self._hairline(slide, card_x + 0.35, card_y + 2.55, card_w - 0.7, weight=0.5)
        self._text(
            slide, card_x + 0.35, card_y + 2.7, card_w - 0.7, 0.3,
            self.mqs_tier_label or "н/д",
            font=self.sans, size=12, italic=True, color=self.deep_100,
        )

        # 4 key metrics inside card (2x2)
        # B1-fix R-01: _mstr (None → «н/д»); ESS дробный (min bulk/tail) → целое.
        km_y = card_y + 3.1
        km_col_w = (card_w - 0.7) / 2
        _ess_card = self._mstr(
            int(round(self.ess_min)) if isinstance(self.ess_min, (int, float)) else None,
            "{:,}").replace(",", " ")
        mets = [
            ("R²",     self._mstr(self.r_squared, "{:.2f}")),
            ("MAPE",   self._mstr(self.mape_pct, "{:.1f}%")),
            ("R-hat",  self._mstr(self.r_hat_max, "{:.3f}")),
            ("ESS",    _ess_card),
        ]
        for i, (label, val) in enumerate(mets):
            row, col = divmod(i, 2)
            mx = card_x + 0.35 + col * km_col_w
            my = km_y + row * 0.22
            self._text(
                slide, mx, my, km_col_w * 0.45, 0.22, label,
                font=self.sans, size=9, color=self.deep_60,
            )
            self._text(
                slide, mx + km_col_w * 0.45, my, km_col_w * 0.55, 0.22, val,
                font=self.sans, size=9, bold=True, color=self.deep_100,
            )

        # RIGHT: Data overview
        right_x = card_x + card_w + 0.5
        right_w = self.w - self.safe - right_x

        self._text(
            slide, right_x, card_y, right_w, 0.25,
            "ОХВАТ И ПОЛНОТА",
            font=self.sans, size=9, bold=True, color=self.gold,
        )
        self._hairline(slide, right_x, card_y + 0.28, 1.2, weight=0.75, color=self.gold)

        # B1-fix R-03/R-04 (2026-07-03): охват данных — ТОЛЬКО факты. Прежде
        # «Наблюдений» = каналы×13 (выдуманная арифметика: Kagocel-зонд рисовал
        # 52 при реальных 31), «Частота: Еженедельно» — hardcode на месячных
        # данных, «Полнота 100%» и «Аномалии обработаны» — не проверялись.
        # Теперь: data_coverage от адаптера (реальные даты); нет факта — «—»
        # или строка скрыта (live не выдумывает).
        active_count = len(self.channels) if self.channels else 0
        _cov = self.data_coverage
        if self.is_live:
            _n_obs = _cov.get("n_observations")
            _gaps = _cov.get("date_gaps")
            data_info = [
                ("Период",           self.data_window_label or "—"),
                ("Наблюдений",       f"{_n_obs}" if _n_obs else "—"),
                ("Активных каналов", f"{active_count}" if active_count else "—"),
                ("Частота",          _cov.get("frequency_label") or "—"),
            ]
            if _gaps is not None:
                data_info.append(
                    ("Разрывы оси дат", "нет" if _gaps == 0 else f"{_gaps}"))
        else:
            data_info = [
                ("Период",              self.data_window_label),
                ("Наблюдений",          f"{active_count * 13}" if active_count else "-"),
                ("Активных каналов",    f"{active_count}" if active_count else "-"),
                ("Частота",             "Еженедельно (Пн-Вс)"),
                ("Полнота",             "100% (0 пропусков)"),
                ("Аномалии",            "обработаны (праздничные недели)"),
            ]
        dy = card_y + 0.55
        for label, val in data_info:
            self._text(
                slide, right_x, dy, right_w * 0.45, 0.3, label,
                font=self.sans, size=10, color=self.deep_60,
            )
            self._text(
                slide, right_x + right_w * 0.45, dy, right_w * 0.55, 0.3, val,
                font=self.sans, size=10, color=self.deep_100,
                align=PP_ALIGN.RIGHT,
            )
            self._hairline(slide, right_x, dy + 0.32, right_w, weight=0.25)
            dy += 0.34

        # B1-fix R-16 (2026-07-03): honesty-контур модели в клиентском отчёте.
        # Вердикт надёжности + preflight-провал приоров вычислялись движком и
        # показывались в программе, но PPTX молчал («вычисленная, но не
        # доставленная честность»). Kagocel-зонд: prior_predictive=fail
        # (coverage 42%) — ни слова в отчёте.
        _hy = dy + 0.05
        _verdict_ru = {
            "reliable": "модель надёжна",
            "uncertain": "требует осторожности",
            "unreliable": "ненадёжна — выводы ориентировочные",
            "unknown": "качество не измерено",
        }.get(self.honesty_verdict)
        if _verdict_ru:
            self._text(
                slide, right_x, _hy, right_w, 0.25,
                f"Вердикт надёжности: {_verdict_ru}",
                font=self.sans, size=10, bold=True,
                color=(self.deep_100 if self.honesty_verdict == "reliable" else self.gold),
            )
            _hy += 0.28
        _pp_status = (self.preflight or {}).get("prior_predictive_status")
        if _pp_status == "fail":
            _pp_cov = (self.preflight or {}).get("prior_predictive_coverage")
            _cov_sfx = (f" (покрытие {_pp_cov:.0%})" if isinstance(_pp_cov, (int, float))
                        else "")
            self._text(
                slide, right_x, _hy, right_w, 0.45,
                f"Проверка приоров: расхождение с данными{_cov_sfx} – "
                "оценки чувствительны к допущениям модели.",
                font=self.sans, size=10, italic=True, color=self.gold,
            )
            _hy += 0.45

        # E2 (2026-07-03): калибровка lift-тестами — строка «откалиброван
        # тестом от <дата>» + ЧЕСТНОЕ расхождение модель-vs-тест (within_ci
        # false не замалчивается — §E2.4).
        for _ca in (self.calibration.get("applied") or []):
            self._text(
                slide, right_x, _hy, right_w, 0.25,
                (f"Канал «{_ca.get('channel')}» откалиброван тестом "
                 f"({_ca.get('test_type')}) от {_ca.get('date_from')} – "
                 f"вошёл в модель как наблюдение."),
                font=self.sans, size=10, color=self.deep_100,
            )
            _hy += 0.26
        for _cc in (self.calibration.get("checks") or []):
            if _cc.get("within_ci"):
                continue
            _ci = _cc.get("model_contrib_ci90") or [None, None]
            self._text(
                slide, right_x, _hy, right_w, 0.45,
                (f"Модель и тест расходятся по «{_cc.get('channel')}»: вклад "
                 f"модели {self._mstr(_cc.get('model_contrib_mean'), '{:,.0f}').replace(',', ' ')} "
                 f"[{self._mstr(_ci[0], '{:,.0f}').replace(',', ' ')} – "
                 f"{self._mstr(_ci[1], '{:,.0f}').replace(',', ' ')}] против теста "
                 f"{self._mstr(_cc.get('test_lift'), '{:,.0f}').replace(',', ' ')} – "
                 f"разберите период с аналитиком."),
                font=self.sans, size=10, italic=True, color=self.gold,
            )
            _hy += 0.45

        # E4 (2026-07-03): сбывшиеся рекомендации — петля доверия в отчёте.
        # Только сверенные (kept/missed); «не сбылось» показывается так же
        # прямо, как «сбылось».
        if self.promises_summary:
            _ps = self.promises_summary
            self._text(
                slide, right_x, _hy, right_w, 0.25,
                (f"Проверка прошлых рекомендаций: сбылось {_ps.get('kept', 0)}, "
                 f"не сбылось {_ps.get('missed', 0)}."),
                font=self.sans, size=10, bold=True, color=self.deep_100,
            )
            _hy += 0.28
            for _ex in (_ps.get("examples") or []):
                self._text(
                    slide, right_x, _hy, right_w, 0.4,
                    f"«{_ex.get('action_text')}» — {_ex.get('status_ru')}.",
                    font=self.sans, size=9.5,
                    color=(self.deep_100 if _ex.get('status') == 'kept' else self.gold),
                )
                _hy += 0.30

        # INV-50 F-DELIVERABLE-1 (2026-06-07): честная оговорка о тонких данных /
        # переобучении — та же формулировка, что в вердикте программы и письме
        # (utils.diagnostics.format_thinness_caveat — SSOT). Прежде роняла на
        # report-шве: клиентский PPTX показывал «MQS 70 Хорошее» без предупреждения.
        from utils.diagnostics import format_thinness_caveat
        _caveat = format_thinness_caveat(self.ratio_eff, self.thinness_cap, leading_space=False)
        if _caveat:
            self._text(
                slide, right_x, _hy + 0.10, right_w, 0.9, _caveat,
                font=self.sans, size=10, italic=True, color=self.gold,
            )

        # Sources list bottom
        src_y = 6.0
        self._text(
            slide, self.safe, src_y, 2.0, 0.2, "PRIMARY",
            font=self.sans, size=8, bold=True, color=self.gold,
        )
        self._hairline(slide, self.safe, src_y + 0.22, 0.8, weight=0.5, color=self.gold)

        # B1-fix R-05 (2026-07-03): PRIMARY/SECONDARY источники были выдуманы
        # целиком (Mediascope/Metrica/GA/VK/бренд-трекер) — клиент таких данных
        # не передавал. В live честно: источник = загруженный файл клиента.
        if self.is_live:
            _n_obs_src = self.data_coverage.get("n_observations")
            _obs_sfx = f", наблюдений: {_n_obs_src}" if _n_obs_src else ""
            primary = [
                f"Данные {self.sources_client_label}: файл, загруженный в проект"
                f" (медиа-траты и продажи{_obs_sfx})",
            ]
        else:
            primary = [
                f"Mediascope TV {self.data_window_label} (TRP / Reach / CPP по target W25-54)",
                f"{self.sources_client_label} internal sales ledger (weekly, SKU-aggregated)",
            ]
        py = src_y + 0.32
        for i, s in enumerate(primary, start=1):
            self._rich(
                slide, self.safe, py, 6.0, 0.22,
                runs=[
                    (f"{i}.  ", {"font": self.sans, "size": 8, "bold": True, "color": self.gold}),
                    (s, {"font": self.sans, "size": 9, "color": self.deep_100}),
                ],
            )
            py += 0.23

        # B1-fix R-05: SECONDARY-блок в live не рисуем (выдуманных источников нет).
        if not self.is_live:
            sec_x = 7.0
            self._text(
                slide, sec_x, src_y, 2.0, 0.2, "SECONDARY",
                font=self.sans, size=8, bold=True, color=self.deep_60,
            )
            self._hairline(slide, sec_x, src_y + 0.22, 0.8, weight=0.5, color=self.deep_60)

            secondary = [
                "Yandex.Metrica + Google Analytics (digital clicks, conversions)",
                "VK Ads + MyTarget (social impressions / clicks)",
                f"{self.sources_client_label} бренд-трекер (quarterly W25 54)",
            ]
            sy = src_y + 0.32
            for i, s in enumerate(secondary, start=1):
                self._rich(
                    slide, sec_x, sy, 6.0, 0.22,
                    runs=[
                        (f"{i}.  ", {"font": self.sans, "size": 8, "bold": True, "color": self.deep_60}),
                        (s, {"font": self.sans, "size": 9, "color": self.deep_100}),
                    ],
                )
                sy += 0.23

        self._footer(slide, 10 + self._page_shift)

    # ----------------------------------------------------------------
    # SLIDE 12 - COLOPHON (narrative)
    # ----------------------------------------------------------------

    # ----------------------------------------------------------------
    # SLIDE 12 - GLOSSARY (тезаурус терминов)
    # ----------------------------------------------------------------

    # C.6.3: s12 glossary content now at physical page 15 (appendix divider at 14).
    def s10b_backtest(self):
        """E1 (2026-07-03): слайд «Проверка на истории» — модель против факта.

        Rolling-origin: модель переобучалась только на прошлом и предсказывала
        удержанные кварталы; прогноз сверен с фактом (канон Gelman BW §9.4 —
        доверие проверяется вне выборки). Рисуется ТОЛЬКО при живой витрине
        (self.backtest) — wireframe-дефолта нет по построению.
        """
        bt = self.backtest
        slide = self._blank()
        # П5: витрина — слайд №6, финал секции «Главное».
        self._header(slide, slide_num=6)

        verdict = bt.get("verdict")
        hit = bt.get("windows_hit_total")
        n_int = bt.get("windows_with_interval") or 0
        gran = bt.get("granularity")
        h = bt.get("horizon_periods")
        is_quarter = (gran == "M" and h == 3) or (gran == "W" and h == 13) or (gran == "D" and h == 90)
        win_word = "кварталов" if is_quarter else "окон проверки"

        if verdict == "validated" and hit is not None and n_int:
            title = (f"Проверка на истории: {hit} из {n_int} {win_word} — "
                     f"факт в прогнозном интервале")
        elif verdict == "worse_than_naive":
            title = "Проверка на истории: модель пока не точнее наивного прогноза"
        elif verdict == "coverage_low":
            title = "Проверка на истории: интервалы модели самоуверенны"
        else:
            title = "Проверка модели на истории"
        self._action_title(slide, title, show_lime=True, y=0.80, height=0.80)

        left_x = self.safe
        left_y = 1.85
        left_w = (self.w - 2 * self.safe) * 0.40

        # LEFT: герой-цифра + факты
        self._text(
            slide, left_x, left_y, left_w, 0.25, "МОДЕЛЬ ПРОТИВ ФАКТА",
            font=self.sans, size=9, bold=True, color=self.gold,
        )
        self._hairline(slide, left_x, left_y + 0.28, 1.0, weight=0.75, color=self.gold)
        if hit is not None and n_int:
            self._text(
                slide, left_x, left_y + 0.55, left_w, 1.0,
                f"{hit} из {n_int}",
                font=self.serif, size=44, bold=True, color=self.deep_100,
            )
            self._text(
                slide, left_x, left_y + 1.6, left_w, 0.6,
                f"{win_word} — фактические продажи попали в 90%-интервал прогноза",
                font=self.sans, size=11, color=self.deep_60, line_spacing=1.2,
            )
        facts_y = left_y + 2.4
        mape_model = bt.get("mape_model")
        mape_naive = bt.get("mape_naive_best")
        cov_pp = bt.get("coverage_per_period")
        n_pts = bt.get("n_holdout_points_with_interval")
        fact_lines = [
            (f"Ошибка прогноза (MAPE): {self._mstr(mape_model, '{:.1f}%')}",
             {"font": self.sans, "size": 11, "color": self.deep_100}),
        ]
        if mape_naive is not None:
            gain = (1 - float(mape_model) / float(mape_naive)) * 100 if mape_model is not None and mape_naive else None
            naive_line = f"Наивный прогноз: {self._mstr(mape_naive, '{:.1f}%')}"
            if gain is not None and gain > 0:
                naive_line += f" — модель точнее на {gain:.0f}%"
            fact_lines.append((naive_line, {"font": self.sans, "size": 11, "color": self.deep_100}))
        if cov_pp is not None:
            fact_lines.append((
                f"Покрытие по периодам: {float(cov_pp) * 100:.0f}% ({n_pts} точек, норма ≈ 90%)",
                {"font": self.sans, "size": 11, "color": self.deep_100},
            ))
        self._paragraphs(
            slide, left_x, facts_y, left_w, 1.6, fact_lines, line_spacing=1.5,
        )

        # RIGHT: таблица окон (период / факт / прогноз / попадание)
        right_x = left_x + left_w + 0.5
        right_w = self.w - self.safe - right_x
        ry = left_y
        self._text(
            slide, right_x, ry, right_w, 0.25, "ОКНА ПРОВЕРКИ",
            font=self.sans, size=9, bold=True, color=self.gold,
        )
        self._hairline(slide, right_x, ry + 0.28, 1.0, weight=0.75, color=self.gold)
        ry += 0.5
        # Заголовок таблицы
        col_period_w = right_w * 0.40
        col_num_w = right_w * 0.18
        for label, cx, cw in (
            ("Период", right_x, col_period_w),
            ("Факт", right_x + col_period_w, col_num_w),
            ("Прогноз", right_x + col_period_w + col_num_w, col_num_w),
            ("90%-интервал", right_x + col_period_w + 2 * col_num_w, right_w - col_period_w - 2 * col_num_w - 0.5),
            ("✓", right_x + right_w - 0.4, 0.4),
        ):
            self._text(
                slide, cx, ry, cw, 0.22, label,
                font=self.sans, size=8.5, bold=True, color=self.deep_60,
            )
        ry += 0.3
        windows = bt.get("windows") or []
        # До 8 окон — строки по 0.34"; при больших N таблица остаётся читаемой.
        for w_row in windows[:8]:
            hit_mark = w_row.get("hit_total")
            self._text(
                slide, right_x, ry, col_period_w, 0.24,
                str(w_row.get("window") or "—"),
                font=self.sans, size=9, color=self.deep_100,
            )
            self._text(
                slide, right_x + col_period_w, ry, col_num_w, 0.24,
                self._mstr(w_row.get("actual_total"), "{:,.0f}").replace(",", " "),
                font=self.sans, size=9, color=self.deep_100,
            )
            self._text(
                slide, right_x + col_period_w + col_num_w, ry, col_num_w, 0.24,
                self._mstr(w_row.get("predicted_total"), "{:,.0f}").replace(",", " "),
                font=self.sans, size=9, color=self.deep_100,
            )
            lo, hi = w_row.get("pi_low_total"), w_row.get("pi_high_total")
            interval = (
                f"{lo:,.0f} – {hi:,.0f}".replace(",", " ")
                if lo is not None and hi is not None else "—"
            )
            self._text(
                slide, right_x + col_period_w + 2 * col_num_w, ry,
                right_w - col_period_w - 2 * col_num_w - 0.5, 0.24,
                interval, font=self.sans, size=9, color=self.deep_60,
            )
            self._text(
                slide, right_x + right_w - 0.4, ry, 0.4, 0.24,
                "✓" if hit_mark else ("—" if hit_mark is None else "✕"),
                font=self.sans, size=9, bold=True,
                color=self.deep_100 if hit_mark else self.deep_60,
            )
            self._hairline(slide, right_x, ry + 0.27, right_w, weight=0.25)
            ry += 0.34

        # Сноска метода — по-русски, без жаргона на слайде.
        _mean_only = bt.get("pi_method") == "posterior_hdi_90_mean_only"
        self._text(
            slide, left_x, 6.35, self.w - 2 * self.safe, 0.5,
            ("Метод: скользящая проверка — модель каждый раз обучается только на прошлом "
             "и предсказывает следующий период; будущее ей не показывают. "
             + ("Интервал — 90% по средней прогноза (без шума наблюдения)."
                if _mean_only else
                "Интервал — 90% предиктивный (неопределённость модели и шум наблюдений).")),
            font=self.sans, size=9, color=self.deep_60, line_spacing=1.25,
        )
        self._footer(slide, 6)

    def s10c_generation_compare(self):
        """E3 (2026-07-03): слайд «Что изменилось с прошлого квартала».

        Сравнение поколений модели: per-channel ROI был [CI] → стал [CI]
        с вердиктом по перекрытию интервалов (Jin 2017). Рисуется ТОЛЬКО
        при живом models/generation_compare.json — wireframe-режима нет.
        """
        gc = self.gen_compare
        slide = self._blank()
        self._header(slide, slide_num=6 + int(bool(self.backtest)))

        summary = gc.get("summary") or {}
        counts = summary.get("counts") or {}
        headline = summary.get("headline")
        title = headline or "Что изменилось с прошлого квартала"
        self._action_title(slide, title, show_lime=True, y=0.80, height=0.80, size=20)

        left_x = self.safe
        left_y = 1.85
        left_w = (self.w - 2 * self.safe) * 0.34

        # LEFT: счётчики стабильности + базовое поколение
        self._text(
            slide, left_x, left_y, left_w, 0.25, "ПОКОЛЕНИЯ МОДЕЛИ",
            font=self.sans, size=9, bold=True, color=self.gold,
        )
        self._hairline(slide, left_x, left_y + 0.28, 1.0, weight=0.75, color=self.gold)
        _base_ts = (gc.get("baseline") or {}).get("timestamp") or "—"
        _rows = [
            (f"Сравнение с версией от {_base_ts[:8]}",
             {"font": self.sans, "size": 11, "color": self.deep_100}),
            (f"Стабильных каналов: {counts.get('stable', 0)}",
             {"font": self.sans, "size": 11, "color": self.deep_100}),
            (f"Сдвиг в пределах неопределённости: {counts.get('shift_within_ci', 0)}",
             {"font": self.sans, "size": 11, "color": self.deep_100}),
            (f"Резких сдвигов: {counts.get('shift_strong', 0)}",
             {"font": self.sans, "size": 11, "bold": counts.get('shift_strong', 0) > 0,
              "color": self.deep_100}),
        ]
        self._paragraphs(slide, left_x, left_y + 0.5, left_w, 2.2, _rows, line_spacing=1.5)

        _causes = summary.get("probable_causes") or []
        if counts.get('shift_strong', 0) > 0 and _causes:
            self._text(
                slide, left_x, left_y + 2.6, left_w, 0.25, "ВЕРОЯТНЫЕ ПРИЧИНЫ",
                font=self.sans, size=9, bold=True, color=self.gold,
            )
            self._paragraphs(
                slide, left_x, left_y + 2.9, left_w, 1.6,
                [(f"• {c}", {"font": self.sans, "size": 10, "color": self.deep_60})
                 for c in _causes[:3]],
                line_spacing=1.3,
            )

        # RIGHT: таблица каналов «был [CI] → стал [CI] · вердикт»
        right_x = left_x + left_w + 0.5
        right_w = self.w - self.safe - right_x
        ry = left_y
        self._text(
            slide, right_x, ry, right_w, 0.25, "ROI ПО КАНАЛАМ",
            font=self.sans, size=9, bold=True, color=self.gold,
        )
        self._hairline(slide, right_x, ry + 0.28, 1.0, weight=0.75, color=self.gold)
        ry += 0.5
        name_w = right_w * 0.34
        vals_w = right_w * 0.42
        for ch in (gc.get("channels") or [])[:8]:
            def _ci_txt(ci):
                if not ci or ci[0] is None or ci[1] is None:
                    return ""
                return f" [{float(ci[0]):.1f}–{float(ci[1]):.1f}]"
            self._text(
                slide, right_x, ry, name_w, 0.24, str(ch.get("name") or "—"),
                font=self.sans, size=9, color=self.deep_100,
            )
            self._text(
                slide, right_x + name_w, ry, vals_w, 0.24,
                (f"{float(ch.get('roi_old') or 0):.1f}{_ci_txt(ch.get('roi_ci_old'))} → "
                 f"{float(ch.get('roi_new') or 0):.1f}{_ci_txt(ch.get('roi_ci_new'))}"),
                font=self.sans, size=9, color=self.deep_100,
            )
            self._text(
                slide, right_x + name_w + vals_w, ry, right_w - name_w - vals_w, 0.24,
                str(ch.get("verdict_ru") or ""),
                font=self.sans, size=9,
                bold=ch.get("verdict") == "shift_strong",
                color=self.deep_100 if ch.get("verdict") == "shift_strong" else self.deep_60,
            )
            self._hairline(slide, right_x, ry + 0.27, right_w, weight=0.25)
            ry += 0.34

        self._text(
            slide, left_x, 6.35, self.w - 2 * self.safe, 0.5,
            ("Метод: обе версии модели пересчитаны на сегодняшних данных; вердикт — "
             "по перекрытию интервалов неопределённости (у оценок есть разброс, "
             "сравниваются интервалы, а не голые точки)."),
            font=self.sans, size=9, color=self.deep_60, line_spacing=1.25,
        )
        self._footer(slide, 6 + int(bool(self.backtest)))

    def s_forecast_plan(self):
        """E5 (2026-07-10): слайд «Прогноз на будущий период».

        Сравнительная таблица сценариев бюджетного плана: имя, бюджет,
        прогноз KPI, доверительный интервал, ROAS. Рисуется ТОЛЬКО при
        живом results/planning.json + сценариях — wireframe-режима нет.
        """
        fc = self.forecast
        slide = self._blank()
        slide_num = 6 + int(bool(self.backtest)) + int(bool(self.gen_compare))
        self._header(slide, slide_num=slide_num)

        self._action_title(
            slide, "Прогноз на будущий период",
            show_lime=True, y=0.80, height=0.80,
        )

        left_x = self.safe
        content_y = 1.85
        content_w = self.w - 2 * self.safe

        # Заголовок таблицы сценариев
        self._text(
            slide, left_x, content_y, content_w, 0.25, "ВАРИАНТЫ БЮДЖЕТНОГО ПЛАНА",
            font=self.sans, size=9, bold=True, color=self.gold,
        )
        self._hairline(slide, left_x, content_y + 0.28, content_w, weight=0.75, color=self.gold)

        ry = content_y + 0.55
        scenarios = (fc.get("scenarios") or [])[:4]

        # Ширины колонок
        col_name_w = content_w * 0.28
        col_budget_w = content_w * 0.20
        col_kpi_w = content_w * 0.20
        col_ci_w = content_w * 0.20
        col_roas_w = content_w - col_name_w - col_budget_w - col_kpi_w - col_ci_w

        # Строка заголовков колонок
        for label, cx, cw in (
            ("Сценарий", left_x, col_name_w),
            ("Бюджет, ₽", left_x + col_name_w, col_budget_w),
            ("Прогноз KPI", left_x + col_name_w + col_budget_w, col_kpi_w),
            ("90%-интервал", left_x + col_name_w + col_budget_w + col_kpi_w, col_ci_w),
            ("ROAS", left_x + col_name_w + col_budget_w + col_kpi_w + col_ci_w, col_roas_w),
        ):
            self._text(
                slide, cx, ry, cw, 0.22, label,
                font=self.sans, size=8.5, bold=True, color=self.deep_60,
            )
        ry += 0.32

        accepted = fc.get("accepted_variant")
        for sc in scenarios:
            is_accepted = sc.get("variant_id") == accepted
            row_color = self.deep_100

            # Интервал СУММЫ за горизонт (totals.predicted_kpi_ci_*) — та же
            # величина, что «Прогноз KPI» рядом и что карточка в GUI (SSOT).
            # P-2 fix 2026-07-16: прежний код брал CI ПОСЛЕДНЕГО периода —
            # число другого масштаба рядом с суммарным KPI читалось как чужое.
            ci_low_val = sc.get("total_kpi_ci_low")
            ci_high_val = sc.get("total_kpi_ci_high")
            ci_str = (
                f"{ci_low_val:,.0f} – {ci_high_val:,.0f}".replace(",", " ")
                if ci_low_val is not None and ci_high_val is not None
                else "—"
            )
            budget_str = (
                f"{sc.get('total_spend_money'):,.0f}".replace(",", " ")
                if sc.get("total_spend_money") is not None else "—"
            )
            kpi_str = (
                f"{sc.get('total_kpi'):,.0f}".replace(",", " ")
                if sc.get("total_kpi") is not None else "—"
            )
            roas_str = (
                f"{sc.get('roas_money'):.2f}"
                if sc.get("roas_money") is not None else "—"
            )
            name_str = str(sc.get("name") or sc.get("variant_id") or "—")
            if is_accepted:
                name_str = f"★ {name_str}"

            for txt, cx, cw in (
                (name_str, left_x, col_name_w),
                (budget_str, left_x + col_name_w, col_budget_w),
                (kpi_str, left_x + col_name_w + col_budget_w, col_kpi_w),
                (ci_str, left_x + col_name_w + col_budget_w + col_kpi_w, col_ci_w),
                (roas_str, left_x + col_name_w + col_budget_w + col_kpi_w + col_ci_w, col_roas_w),
            ):
                self._text(
                    slide, cx, ry, cw, 0.24, txt,
                    font=self.sans, size=9,
                    bold=is_accepted,
                    color=row_color,
                )
            self._hairline(slide, left_x, ry + 0.27, content_w, weight=0.25)
            ry += 0.34

        # Оговорки
        disclaimers = fc.get("disclaimers") or []
        if disclaimers:
            disc_text = " · ".join(str(d) for d in disclaimers[:3])
            self._text(
                slide, left_x, 6.35, content_w, 0.5,
                f"Оговорки: {disc_text}",
                font=self.sans, size=9, color=self.deep_60, line_spacing=1.25,
            )

        # Сравнительный график вариантов (при ≥2 вариантах) — вставляем в правую
        # часть слайда под таблицей. Таблица занимает верхние строки, график —
        # нижнюю половину. Используем тот же PNG base64 → BytesIO паттерн
        # что и в s13_colophon (add_picture из BytesIO объекта).
        if len(scenarios) >= 2:
            try:
                import base64
                import io
                try:
                    from econometrica.charts.generators import scenarios_comparison_chart
                except ImportError:
                    from charts.generators import scenarios_comparison_chart
                kpi_label = self.kpi.get("target_axis") or "Прогноз KPI"
                data_uri = scenarios_comparison_chart(scenarios, kpi_label=kpi_label)
                if data_uri:
                    img_bytes = base64.b64decode(data_uri)
                    img_stream = io.BytesIO(img_bytes)
                    # График под таблицей: слева half + отступ, снизу до footer
                    chart_x = left_x
                    chart_y = ry + 0.15  # чуть ниже последней строки таблицы
                    chart_w = content_w
                    chart_h = max(0.5, 6.0 - chart_y)
                    if chart_h >= 0.5:
                        slide.shapes.add_picture(
                            img_stream,
                            Inches(chart_x), Inches(chart_y),
                            width=Inches(chart_w), height=Inches(chart_h),
                        )
            except Exception:
                pass  # График опционален — ошибка не ломает PPTX

        self._footer(slide, slide_num)

    def s12_glossary(self):
        """Glossary / тезаурус: compact 3-column reference for deck terms."""
        slide = self._blank()
        self._header(slide, slide_num=11 + self._page_shift)

        # B4-2: плашка секции вместо отдельного слайда-дивайдера.
        self._section_intro(
            slide, 5, "Приложение и источники",
            "Глоссарий терминов, методологические ссылки и контакты",
        )

        self._text(
            slide, self.safe, 0.80, self.w - 2 * self.safe, 0.70,
            "Глоссарий терминов",
            font=self.serif, size=36, color=self.deep_100, line_spacing=1.0,
        )
        self._lime_under(slide, self.safe, 1.55, 2.5)

        self._text(
            slide, self.safe, 1.68, self.w - 2 * self.safe, 0.3,
            "Краткие определения ключевых терминов, используемых в отчёте",
            font=self.serif, size=14, italic=True, color=self.deep_60,
        )

        # 3-column layout
        col_gap = 0.4
        col_w = (self.w - 2 * self.safe - 2 * col_gap) / 3
        content_y = 2.30
        entry_h = 0.54

        if self.is_ols:
            _methodology_terms = [
                ("MMM", "Marketing Mix Modeling - статистическая декомпозиция вклада каналов в продажи."),
                ("OLS", "Ordinary Least Squares - метод наименьших квадратов, closed-form оценка β-коэффициентов."),
                # Мат-аудит 2026-07-02 (F-18/INV-50): глоссарий врал о методологии.
                # Было: «n=1000» (движок: n_boot=200), «Conformal PI с гарантией
                # покрытия» (conformal.py F3-caveat: на временных рядах
                # exchangeability нарушена - гарантия НЕ строгая; позиционирование
                # honest PI + caveats, не math-guaranteed).
                ("Bootstrap", "Resampling-метод для построения доверительных интервалов (200 итераций)."),
                ("Frequentist", "Классический статистический подход: точечные оценки + CI на основе sampling distribution."),
                ("Adstock", "Отложенный эффект медиа: часть воздействия переносится на следующие периоды."),
                ("Насыщение", "Убывающая отдача: кривая Хилла, S-образное насыщение эффекта."),
                ("Conformal PI", "Интервал предсказания без предположений о распределении; на временных рядах маркетинга покрытие приближённое (строгая гарантия требует взаимозаменяемости наблюдений)."),
            ]
            _quality_terms = [
                ("R²", "Доля объяснённой моделью вариации продаж (0 до 1). Целевое > 0.7."),
                ("MAPE", "Средняя абсолютная процентная ошибка прогноза."),
                # F-18 (2026-07-02): было «95%» — движок возвращает 90% HDI
                # (compute_ci_hdi, DEFAULT_HDI_PROB=0.9, M-OLS-1).
                ("Bootstrap CI", "90% HDI-интервал на ROI/mROAS из resampling-распределения."),
                ("β-SE", "Стандартная ошибка β-коэффициентов из (XᵀX)⁻¹."),
                ("MQS", "Model Quality Score - композитный индекс качества Aurora (0-100)."),
                ("VIF / Leverage", "OLS-диагностики мультиколлинеарности и влиятельности точек."),
                ("Базовый / инкрементальный", "Органические продажи без медиа и продажи, вызванные медиа-инвестициями."),
            ]
        else:
            _methodology_terms = [
                ("MMM", "Marketing Mix Modeling - статистическая декомпозиция вклада каналов в продажи."),
                ("Байесовский вывод", "Вероятностный подход: апостериорные распределения вместо точечных оценок."),
                ("NUTS", "No-U-Turn Sampler - эффективный метод MCMC для многомерных распределений."),
                ("NumPyro / JAX", "Python-стек для байесовских вычислений с компиляцией JIT."),
                ("Априорное / апостериорное", "Исходные предположения → обновлённая оценка после данных."),
                ("Adstock", "Отложенный эффект медиа: часть воздействия переносится на следующие периоды."),
                ("Насыщение", "Убывающая отдача: кривая Хилла, S-образное насыщение эффекта."),
            ]
            _quality_terms = [
                ("R²", "Доля объяснённой моделью вариации продаж (0 до 1). Целевое > 0.7."),
                ("MAPE", "Средняя абсолютная процентная ошибка прогноза."),
                ("R-hat", "Диагностика сходимости MCMC-цепей; целевое значение ≤ 1.01."),
                ("ESS", "Effective Sample Size - число независимых выборок из апостериорного распределения."),
                # B1-fix R-07: фактический уровень — 90% HDI (DEFAULT_HDI_PROB=0.9);
                # «95%» в глоссарии — то же семейство F-18, не догрепанное b501708.
                ("CI (90% HDI)", "Байесовский интервал правдоподобия - в нём истинное значение лежит с 90% вероятностью."),
                ("MQS", "Model Quality Score - композитный индекс качества Aurora (0-100)."),
                ("Базовый / инкрементальный", "Органические продажи без медиа и продажи, вызванные медиа-инвестициями."),
            ]
        columns = [
            ("МЕТОДОЛОГИЯ MMM", _methodology_terms),
            ("КАЧЕСТВО МОДЕЛИ", _quality_terms),
            ("МЕДИА-МЕТРИКИ", [
                ("mROAS", "Marginal ROAS - возврат с последнего вложенного рубля (×-коэффициент)."),
                ("ROI", "Return on Investment - общая возвратность вложений (инкрементальный вклад / расход)."),
                ("TRP / GRP", "Target / Gross Rating Points - охват целевой аудитории в рейтинг-пунктах."),
                ("CPP", "Cost per Point - стоимость одного рейтингового пункта в рублях."),
                ("Охват (Reach)", "Процент целевой аудитории, встретивших рекламу минимум один раз."),
                ("Доля голоса (SoV)", "Доля рекламного голоса бренда относительно конкурентов в категории."),
                ("Рекомендация", "Вердикт по каналу: Увеличить / Держать / Наблюдать / Сократить / Остановить."),
            ]),
        ]

        for i, (header, terms) in enumerate(columns):
            col_x = self.safe + i * (col_w + col_gap)
            # Column header
            self._text(
                slide, col_x, content_y, col_w, 0.22, header,
                font=self.sans, size=9, bold=True, color=self.gold,
            )
            self._hairline(
                slide, col_x, content_y + 0.25, 1.2, weight=0.75, color=self.gold,
            )

            ey = content_y + 0.45
            for term, defn in terms:
                # Term (bold)
                self._text(
                    slide, col_x, ey, col_w, 0.18, term,
                    font=self.sans, size=9, bold=True, color=self.deep_100,
                )
                # Definition (regular, on next line, may wrap to 2 lines)
                self._text(
                    slide, col_x, ey + 0.20, col_w, 0.32, defn,
                    font=self.sans, size=8, color=self.deep_60, line_spacing=1.15,
                )
                ey += entry_h

        self._footer(slide, 11 + self._page_shift)

    # ----------------------------------------------------------------
    # SLIDE 13 - COLOPHON (closing)
    # ----------------------------------------------------------------

    # C.6.3: s13 colophon content now at physical page 16 (final slide).
    def s13_colophon(self):
        """Closing slide - inspirational brand statement + forward-looking CTA + narrative.
        Flow: statement → CTA (with lime) → narrative → wordmark → copyright.
        No duplication of metrics from other slides."""
        slide = self._blank()
        self._header(slide, slide_num=12 + self._page_shift)

        # Big closing statement - starts higher (more top breathing), no category tag
        self._text(
            slide, self.safe, 1.30, self.w - 2 * self.safe, 1.1,
            "Решения, основанные на данных.",
            font=self.serif, size=52, italic=True, color=self.deep_100,
            line_spacing=1.0,
        )
        # Muted gold second line (softer than primary gold - less shouty)
        self._text(
            slide, self.safe, 2.10, self.w - 2 * self.safe, 0.9,
            "Не на интуиции.",
            font=self.serif, size=52, italic=True, color=self.gold_muted,
            line_spacing=1.0,
        )

        # Forward-looking CTA block with lime accent (action-first)
        cta_y = 4.00
        self._hairline(slide, self.safe, cta_y, 3.5, weight=0.75, color=self.gold)
        self._text(
            slide, self.safe, cta_y + 0.08, 4.0, 0.2, "ДАЛЬШЕ",
            font=self.sans, size=9, bold=True, color=self.gold,
        )
        # B1-fix R-10: «через 90 дней» — шаблонное обещание без основания;
        # в live нейтральная формулировка (срок волны — предмет договорённости).
        self._text(
            slide, self.safe, cta_y + 0.32, self.w - 2 * self.safe, 0.35,
            ("Следующая волна анализа - после накопления новых данных."
             if self.is_live else "Следующая волна анализа - через 90 дней."),
            font=self.serif, size=18, italic=True, color=self.deep_100,
        )
        # Sacred lime underlines CTA (emphasizes action, not statement)
        self._lime_under(slide, self.safe, cta_y + 0.75, 4.5)

        # Narrative paragraph - platform philosophy (positioned AFTER CTA as rationale footer)
        # B1-fix R-06-семейство: «результаты уровня ведущих консалтинговых групп»
        # — недоказуемое сравнительное заявление (INV-50); заменено фактичным.
        # B4-3 (стайлгайд §4): сплошной массив 4 строки разбит — лид-мысль
        # выделена, дальше два коротких предложения (Ильяхов: одна мысль —
        # одно предложение; массивы без акцентов тяжело воспринимать).
        self._rich(
            slide, self.safe, 5.00, self.w - 2 * self.safe, 1.5,
            runs=[
                ("Медиабюджет - управляемый инструмент роста, а не статья затрат. ",
                 {"font": self.serif, "size": 11, "bold": True, "italic": True, "color": self.deep_80}),
                ("Байесовский вывод измеряет эффективность каналов вместе с границами "
                 "неопределённости - это основа доверия к решениям. Методология следует "
                 "отраслевым стандартам MMM; платформа масштабируется от квартального отчёта "
                 "до ежемесячного мониторинга, от одной позиции до портфеля брендов.",
                 {"font": self.serif, "size": 11, "italic": True, "color": self.deep_60}),
            ],
            line_spacing=1.5,
        )

        # Big centered wordmark
        # 5d (audit-iter 2026-05-04): replace typographic wordmark с full PNG brand mark
        # на back cover. Customer reference image showed sigil + AURORA AI как hero
        # brand element. Centered, height 0.6" (matches old wordmark size visually).
        brand_full = self._brand_mark_full_path()
        if brand_full is not None:
            # Square 0.6×0.6 inch, centered horizontally.
            slide.shapes.add_picture(
                str(brand_full),
                Inches((self.w / 2) - 0.30), Inches(6.45),
                height=Inches(0.6),
            )
        else:
            self._wordmark(slide, (self.w / 2) - 0.9, 6.55, size=18, color=self.deep_100)
        # Copyright compact (no "Подготовлено для Client" - client name not repeated on closing)
        self._text(
            slide, self.safe, 6.90, self.w - 2 * self.safe, 0.18,
            "© 2026 Aurora AI.  Не подлежит распространению без письменного согласия",
            font=self.sans, size=7, italic=True, color=self.deep_60,
            align=PP_ALIGN.CENTER,
        )
        # Post-audit: use the shared footer helper (suppress left wordmark to
        # avoid duplicating the hero wordmark in the colophon body).
        self._footer(slide, self.total_slides, show_wordmark=False)

    # ---------- Build ----------

    def build(self):
        # B4-2 (2026-07-03, стайлгайд §1): 12-slide layout БЕЗ отдельных
        # слайдов-дивайдеров. Прежняя 16-слайдовая схема (Stage C.6.3) несла
        # 4 полупустых перебивки (25% деки); теперь секцию открывает компактная
        # плашка _section_intro на первом содержательном слайде.
        # Cover(1) → TOC(2) → Главное (3 at-a-glance, 4 key message, 5 SCQAR)
        # → Декомпозиция (6 chart+плашка, 7 table, 8 timeline) →
        # Методология (9+плашка) → Данные (10+плашка) →
        # Приложение (11 glossary+плашка, 12 colophon).
        self.s01_cover()
        self.s03_toc()
        self.s02_at_a_glance()
        self.s05_key_message()
        self.s09_scqar()
        # E1+П5+E3 (2026-07-03): вставные слайды честности завершают «Главное»
        # (после SCQAR): витрина «Проверка на истории», затем «Что изменилось
        # с прошлого квартала» — только при живых артефактах.
        if self.backtest:
            self.s10b_backtest()
        if self.gen_compare:
            self.s10c_generation_compare()
        if self.forecast:
            self.s_forecast_plan()
        self.s06_action_chart()
        self.s07_action_table()
        self.s08_action_timeline()
        self.s10_methodology()
        self.s11_sources()
        self.s12_glossary()
        self.s13_colophon()
        return self.prs


# Module-level - no main() entry. Use AuroraPPTXBuilder().build() or aurora_pptx.build_pptx().
