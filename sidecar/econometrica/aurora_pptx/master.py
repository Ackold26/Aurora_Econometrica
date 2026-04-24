"""
Sacred DNA — master elements applied to every slide.

Contract (from Standards/CLIENT_READY_ANATOMY.md §3):
  - Logo (top-left cover OR bottom-left content)
  - Breadcrumb footer (bottom-left, 8pt, Aurora Deep 60)
  - Page number (bottom-right, 9pt, Aurora Deep 60)
  - Sacred lime 2pt line (under action-title) — optional per layout
  - Confidentiality marker (bottom-left, italic 7pt) — if ctx.confidential

Entry: apply_master_elements(slide, ctx) — called by every layout render_*().
"""

from __future__ import annotations

from dataclasses import dataclass, field
from datetime import date
from pathlib import Path
from typing import Optional

from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.util import Inches, Pt

from .tokens import COLOR, FONT, SIZE, SIG_LIME


@dataclass
class SlideContext:
    """Per-slide master element configuration."""
    page_num: int
    total_pages: int
    project_id: str = "UNKNOWN"
    version: str = "1.0.11"
    generated_at: date = field(default_factory=date.today)
    client_name: str = "Клиент"
    confidential: bool = True
    show_lime: bool = True
    show_page_num: bool = True
    show_logo: bool = True
    is_cover: bool = False


def _breadcrumb_text(ctx: SlideContext) -> str:
    return f"AURORA AI · {ctx.generated_at.isoformat()} · v{ctx.version} · DEL-{ctx.project_id}"


def _add_logo_placeholder(slide, left: float, top: float, size: float):
    """Placeholder logo — gold square with "A". Replace with real SVG→PNG later."""
    box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(size), Inches(size)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR.brand.gold
    box.line.fill.background()
    tf = box.text_frame
    tf.margin_left = Inches(0); tf.margin_right = Inches(0)
    tf.margin_top = Inches(0); tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "A"
    r.font.name = FONT.family.serif
    r.font.size = Pt(max(8, size * 50))  # scale font with box
    r.font.bold = True
    r.font.color.rgb = COLOR.brand.bg_white


def _add_text(slide, left, top, width, height, text, *, font_name, size, color,
              bold=False, italic=False, align=None):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0); tf.margin_right = Inches(0)
    tf.margin_top = Inches(0); tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font_name
    r.font.size = size
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def apply_master_elements(slide, ctx: SlideContext) -> None:
    """Apply sacred DNA to a slide per context."""
    safe = 0.4  # inches, from Standards/02
    slide_w_in = 13.333
    slide_h_in = 7.5

    # --- Logo ---
    if ctx.show_logo:
        if ctx.is_cover:
            _add_logo_placeholder(slide, safe, safe, 0.6)
        else:
            _add_logo_placeholder(slide, safe, slide_h_in - safe - 0.2, 0.2)

    # --- Breadcrumb (bottom-left) ---
    if not ctx.is_cover:
        _add_text(
            slide,
            left=safe + 0.3, top=slide_h_in - safe - 0.2,
            width=slide_w_in - 2 * safe - 2.0, height=0.2,
            text=_breadcrumb_text(ctx),
            font_name=FONT.family.sans,
            size=FONT.size.footnote,
            color=COLOR.brand.deep_60,
        )

    # --- Page number (bottom-right) ---
    if ctx.show_page_num and not ctx.is_cover:
        _add_text(
            slide,
            left=slide_w_in - safe - 1.0, top=slide_h_in - safe - 0.3,
            width=1.0, height=0.25,
            text=f"{ctx.page_num} / {ctx.total_pages}",
            font_name=FONT.family.sans, size=Pt(9),
            color=COLOR.brand.deep_60,
            align=PP_ALIGN.RIGHT,
        )

    # --- Confidentiality marker (top-right, italic 7pt) ---
    if ctx.confidential and not ctx.is_cover:
        _add_text(
            slide,
            left=slide_w_in - safe - 4.0, top=safe,
            width=4.0, height=0.2,
            text=f"CONFIDENTIAL — For {ctx.client_name} only",
            font_name=FONT.family.sans, size=Pt(7),
            color=COLOR.brand.deep_60, italic=True,
            align=PP_ALIGN.RIGHT,
        )


def add_signature_lime(slide, left_in: float, top_in: float, width_in: float) -> None:
    """Sacred lime 2pt line — called by layouts that show it under action titles."""
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(left_in), Inches(top_in),
        Inches(left_in + width_in), Inches(top_in),
    )
    line.line.color.rgb = SIG_LIME
    line.line.width = SIZE.signature_lime_width
