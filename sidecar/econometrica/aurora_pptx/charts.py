"""
Chart helpers - native PPTX chart generation with tokens-driven styling.

Four MMM chart types per CLIENT_READY_ANATOMY.md §6:
  1. Waterfall (decomposition) - BAR_STACKED with baseline + per-channel contribs
  2. ROI bar - BAR_CLUSTERED horizontal, muted-except-one-gold
  3. Share of Spend vs Effect - 100% stacked COLUMN_CLUSTERED
  4. Timeline stacked area - AREA_STACKED over dates

All use theme colors auto-assigned by PowerPoint from our themed .pptx container.
Overrides applied via format.fill.solid() for specific highlights.

M3 Session 3: implement in detail. Currently stubs returning the chart frame
for downstream positioning.
"""

from __future__ import annotations

from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches, Pt

from .tokens import COLOR, FONT


def _style_chart_text(chart) -> None:
    """Apply tokens-driven axes/legend styling to a chart."""
    # Axes - 8pt for ticks (compact, leaves room for many categories)
    for axis in (chart.category_axis, chart.value_axis):
        try:
            axis.tick_labels.font.size = Pt(8)
            axis.tick_labels.font.name = FONT.family.sans
            axis.tick_labels.font.color.rgb = COLOR.brand.deep_80
        except Exception:
            pass
    # Legend - 9pt (slightly smaller than body text, but readable)
    if chart.has_legend:
        chart.legend.font.size = Pt(9)
        chart.legend.font.name = FONT.family.sans
        chart.legend.font.color.rgb = COLOR.brand.deep_80


def make_roi_bar(slide, x_in, y_in, w_in, h_in, *, channels, roi_values, highlight_idx=None):
    """Horizontal bar chart: ROI per channel. M3 stub."""
    data = CategoryChartData()
    data.categories = channels
    data.add_series("ROI", roi_values)
    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in),
        data
    )
    chart = graphic_frame.chart
    chart.has_legend = False
    _style_chart_text(chart)
    # Color series: muted by default
    for i, point in enumerate(chart.plots[0].series[0].points):
        point.format.fill.solid()
        if highlight_idx == i:
            point.format.fill.fore_color.rgb = COLOR.brand.gold
        else:
            point.format.fill.fore_color.rgb = COLOR.brand.deep_40
    return graphic_frame


def make_share_comparison(slide, x_in, y_in, w_in, h_in, *, channels, spend_shares, effect_shares):
    """Two-column stacked comparison: Spend vs Effect per channel. M3 stub."""
    data = CategoryChartData()
    data.categories = ["Доля бюджета", "Доля эффекта"]
    for i, channel in enumerate(channels):
        values = (spend_shares[i], effect_shares[i])
        data.add_series(channel, values)
    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED_100,
        Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in),
        data
    )
    chart = graphic_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    _style_chart_text(chart)
    # Color each series via channel_colors palette
    for i, series in enumerate(chart.plots[0].series):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = COLOR.data.channel_colors[i % len(COLOR.data.channel_colors)]
    return graphic_frame


def make_decomposition_stacked(slide, x_in, y_in, w_in, h_in, *, categories, series_data):
    """Stacked bar chart (waterfall fallback): decomposition of KPI.
    series_data: dict {name: [values]}. M3 stub.
    """
    data = CategoryChartData()
    data.categories = categories
    for name, values in series_data.items():
        data.add_series(name, values)
    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_STACKED,
        Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in),
        data
    )
    chart = graphic_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    _style_chart_text(chart)
    for i, series in enumerate(chart.plots[0].series):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = COLOR.data.channel_colors[i % len(COLOR.data.channel_colors)]
    return graphic_frame


# Аудит #12 (2026-06-07): цвета вынесенных факторов по типу — зеркало палитры
# программы (ChannelTimeline.svelte FACTOR_COLORS).
_FACTOR_RGB = {
    'signed_competitor': 'DC2626',
    'signed_price': 'EA580C',
    'signed_weather': 'F59E0B',
    'signed_macro': 'D97706',
    'holiday': '84CC16',
    'seasonality': '8B5CF6',  # violet-500 — сезонность (цикл), отдельно от внешних
    'category': '10B981',     # emerald-500 — спрос категории/рынка (Фаза Б)
    'positive_control': '06B6D4',
}


def make_timeline_area(slide, x_in, y_in, w_in, h_in, *, dates, baseline,
                       channel_series, factor_series=None):
    """Stacked area: KPI over time with baseline + channel contributions.
    channel_series: dict {channel_name: [values]}.
    factor_series: optional list [{name, type, side, data}] — вынесенные
        signed/holiday факторы (аудит #12); рендерятся теми же полосами, что в
        программе, чтобы отчёт показывал ТОТ ЖЕ набор факторов.

    Chart formatting:
      - Date labels compact "MM.YY" (e.g. "10.21" instead of "2021-10-01")
      - Y-axis values formatted in millions ("25 М" instead of raw 25000000)
      - X-axis sparse labels via tick_label_skip when >= 18 categories
      - Channel name truncation in legend (max 22 chars + ellipsis)
    """
    # Compact date labels: "2021-10-01" → "10.21" (saves ~50% horizontal space).
    # Falls back to original string if parsing fails (non-ISO categories).
    def _short_date(s):
        try:
            parts = str(s).split("-")
            if len(parts) >= 2 and len(parts[0]) == 4:
                return f"{parts[1]}.{parts[0][2:]}"
        except Exception:
            pass
        return str(s)
    short_dates = [_short_date(d) for d in dates]

    def _short(name):
        # Truncate long names to keep legend readable. Soft trim at first
        # newline for multi-line names like "Performance Бюджет\nДО НДС".
        s = str(name).split("\n")[0].rstrip()
        return (s[:21].rstrip() + "…") if len(s) > 22 else s

    from pptx.dml.color import RGBColor

    data = CategoryChartData()
    data.categories = short_dates
    data.add_series("Baseline", baseline)
    for name, values in channel_series.items():
        data.add_series(_short(name), values)
    # Вынесенные факторы — теми же полосами, что в программе (тип несёт цвет).
    # Т3-плюс: свёрнутый обзор 4 групп передаёт explicit 'rgb' на агрегат
    # (тип-агностично), приоритет над _FACTOR_RGB[type].
    factor_types = []
    factor_rgbs = []
    for f in (factor_series or []):
        if not f or not f.get("data"):
            continue
        data.add_series(_short(f.get("name")), f["data"])
        factor_types.append(f.get("type"))
        factor_rgbs.append(f.get("rgb"))
    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.AREA_STACKED,
        Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in),
        data
    )
    chart = graphic_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    _style_chart_text(chart)
    # Baseline = muted grey, channels = palette, факторы = цвет по типу.
    series_list = list(chart.plots[0].series)
    n_ch = len(channel_series)
    if series_list:
        series_list[0].format.fill.solid()
        series_list[0].format.fill.fore_color.rgb = COLOR.brand.deep_40
        for i, series in enumerate(series_list[1:], start=0):
            series.format.fill.solid()
            if i < n_ch:
                series.format.fill.fore_color.rgb = COLOR.data.channel_colors[i % len(COLOR.data.channel_colors)]
            else:
                fi = i - n_ch
                ftype = factor_types[fi] if fi < len(factor_types) else None
                frgb = factor_rgbs[fi] if fi < len(factor_rgbs) else None
                series.format.fill.fore_color.rgb = RGBColor.from_string(
                    frgb or _FACTOR_RGB.get(ftype, '94A3B8'))

    # Y-axis: format in millions. Excel format string "0,," collapses 25000000 → "25"
    # (each comma divides by 1000), suffix " М" appended for "25 М" display.
    try:
        chart.value_axis.tick_labels.number_format = '0,, "М"'
    except Exception:
        pass

    # X-axis: when many dates, skip every Nth label to prevent overlap.
    try:
        n_dates = len(list(dates))
        if n_dates >= 18:
            skip = max(1, n_dates // 12)  # show ~12 labels max
            chart.category_axis.tick_label_skip = skip
    except Exception:
        pass

    return graphic_frame
