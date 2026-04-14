#!/usr/bin/env python3
"""
Aurora AI — PPTX Pipeline
Preprocessing, notes injection, and DOCX generation for media-analyst cabinet.

Modes:
  preprocess   <input.pptx> <output_dir>     → slides.json + styles.json
  inject-notes <input.pptx> <notes.json> <out.pptx>  → PPTX with notes
  generate-docx <input.pptx> <notes.json> <styles.json> <out.docx> → formatted DOCX

Requirements: python-pptx, python-docx
"""

import sys
import json
import os
import re
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Pt
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip install python-pptx", file=sys.stderr)
    sys.exit(1)

try:
    from docx import Document
    from docx.shared import Pt as DocxPt, RGBColor
except ImportError:
    print("ERROR: python-docx not installed. Run: pip install python-docx", file=sys.stderr)
    sys.exit(1)


# ─── PREPROCESS ────────────────────────────────────────────────────────

def extract_text_from_shape(shape):
    """Extract text from a shape, handling text frames and tables."""
    parts = []

    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if text:
                parts.append(text)

    if shape.has_table:
        table = shape.table
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(cells)
        if rows:
            parts.append({"type": "table", "rows": rows})

    return parts


def extract_chart_data(shape):
    """Extract data from chart shapes."""
    if not shape.has_chart:
        return None

    chart = shape.chart
    result = {
        "chart_type": str(chart.chart_type) if chart.chart_type else "unknown",
        "series": [],
    }

    try:
        for series in chart.series:
            series_data = {
                "name": str(series.name) if hasattr(series, 'name') else "",
                "values": [],
            }
            try:
                series_data["values"] = [
                    float(v) if v is not None else None
                    for v in series.values
                ]
            except Exception:
                pass
            result["series"].append(series_data)

        # Try to get category labels
        try:
            plot = chart.plots[0]
            cats = plot.categories
            if cats:
                result["categories"] = [str(c) for c in cats]
        except Exception:
            pass

    except Exception as e:
        result["error"] = str(e)

    return result


def classify_slide(texts, has_chart):
    """Classify slide type: data, title, divider, empty."""
    total_text = " ".join(t for t in texts if isinstance(t, str))

    if not total_text.strip() and not has_chart:
        return "empty"

    # Title/divider slides: short text, no data
    words = total_text.split()
    if len(words) <= 8 and not has_chart:
        # Check for common divider patterns
        lower = total_text.lower()
        if any(w in lower for w in ["содержание", "оглавление", "agenda", "спасибо", "thank",
                                     "источники", "источник", "references", "methodology",
                                     "методология", "приложени"]):
            return "divider"
        if len(words) <= 3:
            return "title"

    if has_chart:
        return "data"

    # Check for data indicators (numbers, percentages, currency)
    data_pattern = re.compile(r'[\d]+[.,]?\d*\s*(%|₽|руб|млн|тыс|млрд|\$|€)')
    if data_pattern.search(total_text):
        return "data"

    # Check for tables in content
    return "data" if len(words) > 15 else "title"


def detect_source(texts):
    """Detect data source from slide text (footer, caption)."""
    source_patterns = [
        (r'(?i)источник[:\s]+(.+?)(?:\n|$)', None),
        (r'(?i)source[:\s]+(.+?)(?:\n|$)', None),
        (r'(?i)(mediascope|ADEX|brand\s*pulse)', "Mediascope"),
        (r'(?i)(DSM\s*Group|DSM)', "DSM Group"),
        (r'(?i)(digital\s*budget)', "Digital Budget"),
        (r'(?i)(wordstat|яндекс\s*wordstat)', "Wordstat"),
        (r'(?i)(nielsen)', "Nielsen"),
        (r'(?i)(similarweb)', "SimilarWeb"),
        (r'(?i)(GfK|гфк)', "GfK"),
    ]

    all_text = " ".join(t for t in texts if isinstance(t, str))

    for pattern, fixed_name in source_patterns:
        m = re.search(pattern, all_text)
        if m:
            return fixed_name or m.group(1).strip()

    return None


def preprocess(input_path, output_dir):
    """Extract all content from PPTX into structured JSON."""
    prs = Presentation(input_path)
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    slides_data = []
    style_fonts = set()
    style_colors = set()
    style_sizes = set()

    for idx, slide in enumerate(prs.slides, 1):
        texts = []
        charts = []
        tables = []
        has_chart = False

        for shape in slide.shapes:
            # Extract text
            parts = extract_text_from_shape(shape)
            for p in parts:
                if isinstance(p, dict) and p.get("type") == "table":
                    tables.append(p["rows"])
                elif isinstance(p, str):
                    texts.append(p)

            # Extract charts
            if shape.has_chart:
                has_chart = True
                chart_data = extract_chart_data(shape)
                if chart_data:
                    charts.append(chart_data)

            # Collect styles
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name:
                            style_fonts.add(run.font.name)
                        if run.font.size:
                            style_sizes.add(run.font.size.pt)
                        try:
                            if run.font.color and run.font.color.type is not None:
                                rgb = run.font.color.rgb
                                if rgb:
                                    style_colors.add(str(rgb))
                        except (AttributeError, TypeError):
                            pass

        # Get existing notes
        notes_text = ""
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            notes_text = notes_frame.text.strip() if notes_frame else ""

        slide_type = classify_slide(texts, has_chart)
        source = detect_source(texts)

        # Determine title (first short text or shape title)
        title = ""
        if slide.shapes.title and slide.shapes.title.text:
            title = slide.shapes.title.text.strip()
        elif texts:
            # Use first text block as title if short enough
            candidate = texts[0]
            if len(candidate.split()) <= 12:
                title = candidate

        slide_info = {
            "slide_num": idx,
            "title": title,
            "type": slide_type,
            "texts": texts,
            "existing_notes": notes_text,
        }

        if charts:
            slide_info["charts"] = charts
        if tables:
            slide_info["tables"] = tables
        if source:
            slide_info["source"] = source

        slides_data.append(slide_info)

    # Determine presentation dimensions
    width_emu = prs.slide_width
    height_emu = prs.slide_height

    styles = {
        "fonts": sorted(style_fonts),
        "colors": sorted(style_colors),
        "sizes": sorted(style_sizes),
        "slide_width_cm": round(width_emu / 914400 * 2.54, 1) if width_emu else 33.9,
        "slide_height_cm": round(height_emu / 914400 * 2.54, 1) if height_emu else 19.1,
        "slide_count": len(slides_data),
        "data_slide_count": sum(1 for s in slides_data if s["type"] == "data"),
    }

    # Write outputs
    slides_path = output_dir / "slides.json"
    styles_path = output_dir / "styles.json"

    with open(slides_path, "w", encoding="utf-8") as f:
        json.dump(slides_data, f, ensure_ascii=False, indent=2)

    with open(styles_path, "w", encoding="utf-8") as f:
        json.dump(styles, f, ensure_ascii=False, indent=2)

    # Summary to stdout
    summary = {
        "status": "ok",
        "slides": len(slides_data),
        "data_slides": styles["data_slide_count"],
        "charts": sum(len(s.get("charts", [])) for s in slides_data),
        "tables": sum(len(s.get("tables", [])) for s in slides_data),
        "slides_json": str(slides_path),
        "styles_json": str(styles_path),
    }
    print(json.dumps(summary, ensure_ascii=False))


# ─── INJECT NOTES ─────────────────────────────────────────────────────

def inject_notes(input_path, notes_path, output_path):
    """Write commentary into PPTX notes pane for each slide."""
    with open(notes_path, "r", encoding="utf-8") as f:
        notes_data = json.load(f)

    if not notes_data:
        print(json.dumps({
            "status": "warning",
            "warning": "notes_data is empty — no slides will be updated",
            "slides_updated": 0,
            "output": str(output_path),
        }, ensure_ascii=False))
        # Still copy original to output so file exists
        import shutil
        shutil.copy2(input_path, output_path)
        return

    prs = Presentation(input_path)

    # notes_data: list of {slide_num, text} or {slide_num, action_title, ceo, cmo, bm}
    notes_map = {}
    for entry in notes_data:
        num = entry.get("slide_num")
        if num is None:
            continue

        # Build notes text from structured or plain format
        if "text" in entry:
            notes_map[num] = entry["text"]
        else:
            parts = []
            if entry.get("action_title"):
                parts.append(f"ACTION TITLE: {entry['action_title']}")
            if entry.get("ceo"):
                parts.append(f"\n[CEO] {entry['ceo']}")
            if entry.get("cmo"):
                parts.append(f"\n[CMO] {entry['cmo']}")
            if entry.get("bm"):
                parts.append(f"\n[BM] {entry['bm']}")
            if parts:
                notes_map[num] = "\n".join(parts)

    for idx, slide in enumerate(prs.slides, 1):
        if idx not in notes_map:
            continue

        notes_text = notes_map[idx]

        # Get or create notes slide with a working text frame.
        # PowerPoint often creates empty notes slides (no body placeholder),
        # causing notes_text_frame to return None. Fix: copy the body
        # placeholder shape from the notes master into the notes slide.
        import copy
        from pptx.oxml.ns import qn

        ns = slide.notes_slide
        if ns.notes_text_frame is None:
            # Find body placeholder in notes master
            body_sp = None
            try:
                for sp in prs.notes_master.shapes:
                    if sp.is_placeholder and sp.placeholder_format.type == 2:  # BODY
                        body_sp = sp._element
                        break
            except Exception:
                pass

            if body_sp is not None:
                new_sp = copy.deepcopy(body_sp)
                sp_tree = ns._element.find(qn('p:cSld'))
                if sp_tree is not None:
                    tree = sp_tree.find(qn('p:spTree'))
                    if tree is not None:
                        tree.append(new_sp)

        tf = ns.notes_text_frame
        if tf is None:
            # Still None after recreation — skip this slide
            continue

        # Clear existing notes: clear all runs from existing paragraphs,
        # then reuse them for new content.
        existing_paras = list(tf.paragraphs)
        for p in existing_paras:
            p.clear()

        # Write new notes with basic formatting
        lines = notes_text.split("\n")
        for i, line in enumerate(lines):
            if i < len(existing_paras):
                para = existing_paras[i]
            else:
                para = tf.add_paragraph()

            # Bold for ACTION TITLE and audience markers
            if line.startswith("ACTION TITLE:"):
                run = para.add_run()
                run.text = line
                run.font.bold = True
                run.font.size = Pt(11)
            elif line.startswith("[CEO]") or line.startswith("[CMO]") or line.startswith("[BM]"):
                # Bold marker, normal text
                marker_end = line.index("]") + 1
                run_marker = para.add_run()
                run_marker.text = line[:marker_end]
                run_marker.font.bold = True
                run_marker.font.size = Pt(10)
                if line[marker_end:]:
                    run_text = para.add_run()
                    run_text.text = line[marker_end:]
                    run_text.font.size = Pt(10)
            else:
                run = para.add_run()
                run.text = line
                run.font.size = Pt(10)

    prs.save(output_path)
    print(json.dumps({
        "status": "ok",
        "slides_updated": len(notes_map),
        "output": str(output_path),
    }, ensure_ascii=False))


# ─── GENERATE DOCX ────────────────────────────────────────────────────

def generate_docx(input_path, notes_path, styles_path, output_path):
    """Create a formatted DOCX document with commentary per slide."""
    with open(notes_path, "r", encoding="utf-8") as f:
        notes_data = json.load(f)

    styles = {}
    if styles_path and os.path.exists(styles_path):
        with open(styles_path, "r", encoding="utf-8") as f:
            styles = json.load(f)

    # Get slide titles from original PPTX
    prs = Presentation(input_path)
    slide_titles = {}
    for idx, slide in enumerate(prs.slides, 1):
        if slide.shapes.title and slide.shapes.title.text:
            slide_titles[idx] = slide.shapes.title.text.strip()
        else:
            # Try to find first text
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text and len(text.split()) <= 12:
                        slide_titles[idx] = text
                        break

    # Determine primary font from styles
    primary_font = "Calibri"
    if styles.get("fonts"):
        # Prefer Russian-compatible fonts
        preferred = ["Calibri", "Arial", "Segoe UI", "Roboto", "Tahoma"]
        for pf in preferred:
            if pf in styles["fonts"]:
                primary_font = pf
                break
        if primary_font == "Calibri" and styles["fonts"]:
            primary_font = styles["fonts"][0]

    # Create document
    doc = Document()

    # Set default font
    style = doc.styles["Normal"]
    style.font.name = primary_font
    style.font.size = DocxPt(11)

    # Title
    pptx_name = Path(input_path).stem
    doc.add_heading(f"Аналитический комментарий: {pptx_name}", level=1)

    # Metadata
    meta = doc.add_paragraph()
    meta.add_run(f"Источник: {Path(input_path).name}").italic = True
    if styles.get("data_slide_count"):
        meta.add_run(f"  |  Слайдов с данными: {styles['data_slide_count']}")

    doc.add_paragraph()  # spacer

    # Process each slide's notes
    for entry in notes_data:
        num = entry.get("slide_num")
        if num is None:
            continue

        # Section header
        title = slide_titles.get(num, f"Слайд {num}")
        doc.add_heading(f"Слайд {num}: {title}", level=2)

        # Get notes text
        if "text" in entry:
            notes_text = entry["text"]
        else:
            parts = []
            if entry.get("action_title"):
                parts.append(f"ACTION TITLE: {entry['action_title']}")
            if entry.get("ceo"):
                parts.append(f"[CEO] {entry['ceo']}")
            if entry.get("cmo"):
                parts.append(f"[CMO] {entry['cmo']}")
            if entry.get("bm"):
                parts.append(f"[BM] {entry['bm']}")
            notes_text = "\n\n".join(parts) if parts else ""

        if not notes_text.strip():
            doc.add_paragraph("(без комментариев)")
            continue

        # Parse and format notes
        for line in notes_text.split("\n"):
            line = line.strip()
            if not line:
                continue

            para = doc.add_paragraph()

            if line.startswith("ACTION TITLE:"):
                run = para.add_run(line)
                run.bold = True
                run.font.size = DocxPt(12)
                run.font.color.rgb = RGBColor(0x1E, 0x40, 0xAF)  # blue
            elif line.startswith("[CEO]"):
                run_marker = para.add_run("[CEO] ")
                run_marker.bold = True
                run_marker.font.color.rgb = RGBColor(0xDC, 0x26, 0x26)  # red
                run_text = para.add_run(line[5:].strip())
                run_text.font.size = DocxPt(11)
            elif line.startswith("[CMO]"):
                run_marker = para.add_run("[CMO] ")
                run_marker.bold = True
                run_marker.font.color.rgb = RGBColor(0xEA, 0x58, 0x0C)  # orange
                run_text = para.add_run(line[5:].strip())
                run_text.font.size = DocxPt(11)
            elif line.startswith("[BM]"):
                run_marker = para.add_run("[BM] ")
                run_marker.bold = True
                run_marker.font.color.rgb = RGBColor(0x16, 0xA3, 0x4A)  # green
                run_text = para.add_run(line[4:].strip())
                run_text.font.size = DocxPt(11)
            else:
                para.add_run(line)

        # Add separator between slides
        doc.add_paragraph("─" * 60)

    doc.save(output_path)
    print(json.dumps({
        "status": "ok",
        "slides_documented": len(notes_data),
        "output": str(output_path),
    }, ensure_ascii=False))


# ─── GENERATE DOCX WITH SYNTHESIS ─────────────────────────────────────

# Synthesis markdown section headers → DOCX styling
SYNTHESIS_SECTIONS = {
    "EXECUTIVE SUMMARY": {"level": 1, "color": RGBColor(0x1E, 0x40, 0xAF)},
    "ОБЩИЙ ВЫВОД": {"level": 1, "color": RGBColor(0x0F, 0x52, 0x8A)},
    "БЛОК:": {"level": 2, "color": RGBColor(0x0D, 0x92, 0x88)},
    "МОСТЫ": {"level": 1, "color": RGBColor(0x7C, 0x3A, 0xED)},
    "РЕКОМЕНДАЦИИ": {"level": 1, "color": RGBColor(0xDC, 0x26, 0x26)},
}


def render_synthesis_section(doc, line):
    """Render a synthesis markdown line with appropriate formatting."""
    stripped = line.strip()

    # Check for section headers (## EXECUTIVE SUMMARY, ## БЛОК: Name, etc.)
    if stripped.startswith("## ") or stripped.startswith("# "):
        header_text = stripped.lstrip("#").strip()

        for key, style in SYNTHESIS_SECTIONS.items():
            if header_text.upper().startswith(key):
                heading = doc.add_heading(header_text, level=style["level"])
                for run in heading.runs:
                    run.font.color.rgb = style["color"]
                return True

        # Generic heading
        doc.add_heading(header_text, level=2)
        return True

    # Numbered items (1. 2. 3.)
    if stripped and stripped[0].isdigit() and ". " in stripped[:4]:
        para = doc.add_paragraph(style="List Number")
        # Bold the number prefix
        dot_pos = stripped.index(". ")
        run_num = para.add_run(stripped[:dot_pos + 2])
        run_num.bold = True
        para.add_run(stripped[dot_pos + 2:])
        return True

    # Bullet items (- or •)
    if stripped.startswith("- ") or stripped.startswith("• "):
        para = doc.add_paragraph(stripped[2:], style="List Bullet")
        return True

    # Bold lines (text wrapped in **)
    if stripped.startswith("**") and stripped.endswith("**"):
        para = doc.add_paragraph()
        run = para.add_run(stripped.strip("*").strip())
        run.bold = True
        return True

    # Regular text
    if stripped:
        doc.add_paragraph(stripped)
        return True

    return False


def generate_docx_with_synthesis(input_path, notes_path, styles_path, synthesis_path, output_path):
    """Create DOCX with synthesis prefix + per-slide commentary."""
    with open(notes_path, "r", encoding="utf-8") as f:
        notes_data = json.load(f)

    styles = {}
    if styles_path and os.path.exists(styles_path):
        with open(styles_path, "r", encoding="utf-8") as f:
            styles = json.load(f)

    synthesis_md = ""
    if synthesis_path and os.path.exists(synthesis_path):
        with open(synthesis_path, "r", encoding="utf-8") as f:
            synthesis_md = f.read()

    # Get slide titles from PPTX
    prs = Presentation(input_path)
    slide_titles = {}
    for idx, slide in enumerate(prs.slides, 1):
        if slide.shapes.title and slide.shapes.title.text:
            slide_titles[idx] = slide.shapes.title.text.strip()
        else:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text and len(text.split()) <= 12:
                        slide_titles[idx] = text
                        break

    # Determine font
    primary_font = "Calibri"
    if styles.get("fonts"):
        preferred = ["Calibri", "Arial", "Segoe UI", "Roboto", "Tahoma"]
        for pf in preferred:
            if pf in styles["fonts"]:
                primary_font = pf
                break
        if primary_font == "Calibri" and styles["fonts"]:
            primary_font = styles["fonts"][0]

    doc = Document()
    style = doc.styles["Normal"]
    style.font.name = primary_font
    style.font.size = DocxPt(11)

    # Title
    pptx_name = Path(input_path).stem
    doc.add_heading(f"Аналитический комментарий: {pptx_name}", level=1)

    meta = doc.add_paragraph()
    meta.add_run(f"Источник: {Path(input_path).name}").italic = True
    if styles.get("data_slide_count"):
        meta.add_run(f"  |  Слайдов с данными: {styles['data_slide_count']}")

    # ─── SYNTHESIS PREFIX ───
    if synthesis_md.strip():
        doc.add_paragraph()  # spacer
        for line in synthesis_md.split("\n"):
            render_synthesis_section(doc, line)
        # Separator between synthesis and per-slide commentary
        doc.add_paragraph()
        sep = doc.add_paragraph()
        run = sep.add_run("═" * 60)
        run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
        doc.add_heading("ПОСЛАЙДОВЫЕ КОММЕНТАРИИ", level=1)

    # ─── PER-SLIDE COMMENTARY (same logic as generate_docx) ───
    for entry in notes_data:
        num = entry.get("slide_num")
        if num is None:
            continue

        title = slide_titles.get(num, f"Слайд {num}")
        doc.add_heading(f"Слайд {num}: {title}", level=2)

        if "text" in entry:
            notes_text = entry["text"]
        else:
            parts = []
            if entry.get("action_title"):
                parts.append(f"ACTION TITLE: {entry['action_title']}")
            if entry.get("ceo"):
                parts.append(f"[CEO] {entry['ceo']}")
            if entry.get("cmo"):
                parts.append(f"[CMO] {entry['cmo']}")
            if entry.get("bm"):
                parts.append(f"[BM] {entry['bm']}")
            notes_text = "\n\n".join(parts) if parts else ""

        if not notes_text.strip():
            doc.add_paragraph("(без комментариев)")
            continue

        for line in notes_text.split("\n"):
            line = line.strip()
            if not line:
                continue

            para = doc.add_paragraph()
            if line.startswith("ACTION TITLE:"):
                run = para.add_run(line)
                run.bold = True
                run.font.size = DocxPt(12)
                run.font.color.rgb = RGBColor(0x1E, 0x40, 0xAF)
            elif line.startswith("[CEO]"):
                run_marker = para.add_run("[CEO] ")
                run_marker.bold = True
                run_marker.font.color.rgb = RGBColor(0xDC, 0x26, 0x26)
                run_text = para.add_run(line[5:].strip())
                run_text.font.size = DocxPt(11)
            elif line.startswith("[CMO]"):
                run_marker = para.add_run("[CMO] ")
                run_marker.bold = True
                run_marker.font.color.rgb = RGBColor(0xEA, 0x58, 0x0C)
                run_text = para.add_run(line[5:].strip())
                run_text.font.size = DocxPt(11)
            elif line.startswith("[BM]"):
                run_marker = para.add_run("[BM] ")
                run_marker.bold = True
                run_marker.font.color.rgb = RGBColor(0x16, 0xA3, 0x4A)
                run_text = para.add_run(line[4:].strip())
                run_text.font.size = DocxPt(11)
            else:
                para.add_run(line)

        doc.add_paragraph("─" * 60)

    doc.save(output_path)
    print(json.dumps({
        "status": "ok",
        "slides_documented": len(notes_data),
        "has_synthesis": bool(synthesis_md.strip()),
        "output": str(output_path),
    }, ensure_ascii=False))


# ─── INJECT SUMMARY SLIDES ─────────────────────────────────────────────

def parse_synthesis_sections(text):
    """Парсит synthesis markdown в список секций {title, body}."""
    sections = []
    current = None

    for line in text.split('\n'):
        stripped = line.strip()
        if stripped.startswith('## '):
            if current:
                sections.append(current)
            title = stripped[3:].strip()
            current = {'title': title, 'body': []}
        elif current is not None:
            if stripped:
                clean = stripped.lstrip('- •*').strip()
                if clean:
                    current['body'].append(clean)

    if current:
        sections.append(current)

    return sections


def inject_summary_slides(input_pptx, synthesis_md_path, styles_json_path, slides_json_path, output_pptx):
    """Добавить summary-слайды в PPTX из synthesis markdown."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    if not os.path.exists(synthesis_md_path):
        print(json.dumps({"status": "skipped", "reason": "synthesis.md not found"}))
        return

    prs = Presentation(input_pptx)

    styles = {}
    if os.path.exists(styles_json_path):
        with open(styles_json_path, 'r', encoding='utf-8') as f:
            styles = json.load(f)

    with open(synthesis_md_path, 'r', encoding='utf-8') as f:
        synthesis = f.read()

    sections = parse_synthesis_sections(synthesis)
    if not sections:
        print(json.dumps({"status": "skipped", "reason": "no sections in synthesis"}))
        return

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    font_name = 'Calibri'
    if styles.get('fonts'):
        font_name = styles['fonts'][0].get('name', 'Calibri')

    blank_layout = prs.slide_layouts[-1]

    for section in sections:
        slide = prs.slides.add_slide(blank_layout)

        left = Inches(0.5)
        top = Inches(0.3)
        width = slide_width - Inches(1)

        title_box = slide.shapes.add_textbox(left, top, width, Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = section['title']
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.name = font_name
        p.font.color.rgb = RGBColor(0x1E, 0x3A, 0x5F)

        body_box = slide.shapes.add_textbox(left, Inches(1.3), width, slide_height - Inches(1.8))
        tf = body_box.text_frame
        tf.word_wrap = True

        for i, line in enumerate(section['body']):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(14)
            p.font.name = font_name
            p.space_after = Pt(4)

    prs.save(output_pptx)
    print(json.dumps({"status": "ok", "slides_added": len(sections)}))


# ─── MAIN ──────────────────────────────────────────────────────────────

def main():
    if len(sys.argv) < 2:
        print("Usage: pptx_pipeline.py <mode> [args...]", file=sys.stderr)
        print("Modes: preprocess, inject-notes, generate-docx, inject-summary-slides", file=sys.stderr)
        sys.exit(1)

    mode = sys.argv[1]

    if mode == "preprocess":
        if len(sys.argv) != 4:
            print("Usage: pptx_pipeline.py preprocess <input.pptx> <output_dir>", file=sys.stderr)
            sys.exit(1)
        preprocess(sys.argv[2], sys.argv[3])

    elif mode == "inject-notes":
        if len(sys.argv) != 5:
            print("Usage: pptx_pipeline.py inject-notes <input.pptx> <notes.json> <output.pptx>", file=sys.stderr)
            sys.exit(1)
        inject_notes(sys.argv[2], sys.argv[3], sys.argv[4])

    elif mode == "generate-docx":
        if len(sys.argv) != 6:
            print("Usage: pptx_pipeline.py generate-docx <input.pptx> <notes.json> <styles.json> <output.docx>", file=sys.stderr)
            sys.exit(1)
        generate_docx(sys.argv[2], sys.argv[3], sys.argv[4], sys.argv[5])

    elif mode == "generate-docx-with-synthesis":
        if len(sys.argv) != 7:
            print("Usage: pptx_pipeline.py generate-docx-with-synthesis <input.pptx> <notes.json> <styles.json> <synthesis.md> <output.docx>", file=sys.stderr)
            sys.exit(1)
        generate_docx_with_synthesis(sys.argv[2], sys.argv[3], sys.argv[4], sys.argv[5], sys.argv[6])

    elif mode == "inject-summary-slides":
        if len(sys.argv) != 7:
            print("Usage: pptx_pipeline.py inject-summary-slides <input.pptx> <synthesis.md> <styles.json> <slides.json> <output.pptx>", file=sys.stderr)
            sys.exit(1)
        inject_summary_slides(sys.argv[2], sys.argv[3], sys.argv[4], sys.argv[5], sys.argv[6])

    else:
        print(f"Unknown mode: {mode}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
