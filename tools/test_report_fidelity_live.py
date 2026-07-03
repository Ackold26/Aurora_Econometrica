"""Репортинг-канарейка B1 (2026-07-03): live-отчёт не выдумывает факты.

Корень (зонд tmp/probe_b1_pptx_extract.py на реальном Kagocel): билдер PPTX
имеет легитимный wireframe/preview-режим с демо-контентом Kagocel-пилота, но
на ЖИВОМ клиентском экспорте пробелы в data тихо латались wireframe-дефолтами:
- ESS 1247 при реальном tail 382.7<400 (маппер терял ess_bulk/tail после F-11);
- период «W01 W13 2026» при данных 2023–2025 (адаптер не передавал);
- «Наблюдений» = каналы×13; «Частота: Еженедельно» на месячных данных;
- источники Mediascope/Metrica/VK — выдуманы целиком;
- подписи «CI 95%» при фактических 90% HDI; «bootstrap n=1000» при n=200;
- приор «HalfNormal(0.5)» (факт 0.3/hierarchical); «12+ FMCG-проектов»;
- honesty-вердикт и prior-predictive FAIL не доставлялись вовсе.

Канарейка: live-режим (channels в data) не содержит wireframe-маркеров;
честные значения (период из дат, n наблюдений, ESS из bulk/tail, вердикт)
доставлены. Preview-режим (без data) сохраняет демку — это легитимно.
"""
from __future__ import annotations

import sys
from pathlib import Path

import pytest

ROOT = Path(__file__).resolve().parents[1]
SIDECAR = ROOT / 'sidecar'
for _p in (str(SIDECAR), str(SIDECAR / 'econometrica')):
    if _p not in sys.path:
        sys.path.insert(0, _p)

from test_goalseek_honesty import _build_project  # noqa: E402
from engines.narrative_adapter import (  # noqa: E402
    _derive_data_coverage,
    _map_pipeline_to_builder_data,
)

# Строки, которым ЗАПРЕЩЕНО появляться в живом клиентском отчёте.
WIREFRAME_MARKERS = [
    'W01 W13 2026',
    'Q1 2026',
    'Q3-Q4 2026',
    'Mediascope',
    'Yandex.Metrica',
    'VK Ads',
    'бренд-трекер',
    '1 247',
    'Еженедельно (Пн-Вс)',
    '100% (0 пропусков)',
    'обработаны (праздничные недели)',
    '12+ FMCG',
    'через 90 дней',
    'интервал 95%',
    'CI (95%)',
    'bootstrap CI 95%',
    'n=1000',
    'HalfNormal(0.5)',
    'откалибрована под FMCG-бенчмарки',
    '~12 мин',
    '~2 800',
    # B1-fix-2 (R-09/R-12/R-13): выдуманные рекомендации и пустые обещания.
    'экономия 15-20%',
    'Пульсирующее размещение',
    'пульсирующее размещение',
    'Целевой ретаргетинг',
    '+0.0 пп',
    '+0 пп',
]


def _extract_all_text(pptx_path: str) -> str:
    from pptx import Presentation
    prs = Presentation(pptx_path)
    chunks = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                chunks.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    chunks.extend(c.text for c in row.cells)
    return '\n'.join(chunks)


# ─── Юнит-уровень адаптера ────────────────────────────────────────────────────

def test_adapter_ess_min_from_bulk_tail():
    """R-01: после F-11 движок отдаёт ess_bulk_min/ess_tail_min — адаптер обязан
    доставить min из них как ess_min (прежде терял, билдер рисовал 1247)."""
    model_data = {'diagnostics': {'metrics': {
        'ess_bulk_min': 590.7, 'ess_tail_min': 382.7, 'r_squared': 0.5}}}
    data = _map_pipeline_to_builder_data(model_data, {}, {}, None)
    assert data['diagnostics']['ess_min'] == 382
    # Прямой ess_min (legacy) имеет приоритет.
    model_data2 = {'diagnostics': {'metrics': {
        'ess_min': 1100, 'ess_bulk_min': 590.7, 'ess_tail_min': 382.7}}}
    data2 = _map_pipeline_to_builder_data(model_data2, {}, {}, None)
    assert data2['diagnostics']['ess_min'] == 1100


def test_adapter_data_coverage_from_dates():
    """R-02/R-03/R-04: период/наблюдения/частота считаются из реальных дат."""
    monthly = {'time_series': {'dates': [f'202{3 + m // 12}-{m % 12 + 1:02d}-01'
                                         for m in range(31)],
                               'baseline': [1.0] * 31}}
    cov = _derive_data_coverage(monthly)
    assert cov is not None
    assert cov['n_observations'] == 31
    assert cov['frequency_label'] == 'Ежемесячно'
    assert cov['window_label'].startswith('янв 2023')
    assert cov['date_gaps'] == 0

    weekly = {'time_series': {'dates': ['2026-01-05', '2026-01-12', '2026-01-19',
                                        '2026-01-26', '2026-02-02'],
                              'baseline': [1.0] * 5}}
    cov_w = _derive_data_coverage(weekly)
    assert cov_w['frequency_label'] == 'Еженедельно'

    # Непарсибельные метки (legacy порядковые) → None: НЕ выдумываем.
    ordinal = {'time_series': {'dates': ['1', '2', '3'], 'baseline': [1.0] * 3}}
    assert _derive_data_coverage(ordinal) is None


def test_adapter_delivers_honesty_and_preflight():
    """R-16: вердикт надёжности + preflight-сводка доезжают до builder-data."""
    model_data = {'diagnostics': {
        'metrics': {'r_squared': 0.5, 'mape_pct': 37.0, 'r_hat_max': 1.01},
        'mqs': {'score': 58.6, 'tier_label': 'Приемлемое'},
        'checks': {},
        'preflight': {
            'prior_predictive': {'status': 'fail', 'coverage': 0.42},
            'quick_proxy': {'tier': 'reliable'},
        },
    }}
    data = _map_pipeline_to_builder_data(model_data, {}, {}, None)
    diag = data['diagnostics']
    assert diag['preflight']['prior_predictive_status'] == 'fail'
    assert diag['preflight']['prior_predictive_coverage'] == pytest.approx(0.42)
    assert diag['preflight']['quick_proxy_tier'] == 'reliable'
    # Вердикт вычислен SSOT-функцией (не изобретён отчётом).
    assert 'honesty_verdict' in diag
    assert diag['honesty_verdict'] in ('reliable', 'uncertain', 'unreliable', 'unknown')


# ─── Интеграционный: живой PPTX без wireframe-вранья ─────────────────────────

@pytest.fixture(scope='module')
def live_pptx_text(tmp_path_factory):
    """Полный live-путь на synthetic-проекте: decompose → optimize →
    build_pptx → извлечённый текст всех слайдов."""
    tmp = tmp_path_factory.mktemp('b1_live')
    pdir = _build_project(tmp, 'fidelity', beta_sd=0.2, seed=7)

    from engines.decomposer import decompose
    from engines.optimizer import optimize
    dec = decompose(str(pdir))
    assert dec.get('status') != 'error', dec.get('message')
    opt = optimize({'min_pct': 0.0, 'max_pct': 100.0}, str(pdir))
    assert opt.get('status', 'ok') != 'error', opt.get('message')

    # model_data с диагностикой нового формата (после F-11: bulk/tail, не ess_min).
    model_data = {'diagnostics': {
        'metrics': {'r_squared': 0.81, 'mape_pct': 12.3, 'r_hat_max': 1.004,
                    'ess_bulk_min': 812.0, 'ess_tail_min': 640.0},
        'mqs': {'score': 71.0, 'tier_label': 'Хорошее'},
        'checks': {},
    }}
    out = str(tmp / 'report.pptx')
    from engines.pptx_export import build_pptx
    res = build_pptx(model_data, dec, opt, out, scenarios=[], project_id='fidelity_test')
    assert res.get('status') == 'ok', res.get('message')
    return _extract_all_text(out)


def test_live_pptx_has_no_wireframe_markers(live_pptx_text):
    """Ни один wireframe-маркер не смешивается с реальными числами клиента."""
    hits = [m for m in WIREFRAME_MARKERS if m in live_pptx_text]
    assert not hits, (
        f'Wireframe-демка протекла в живой клиентский отчёт: {hits}. '
        f'Регрессия B1 (см. docstring).'
    )


def test_live_pptx_delivers_honest_values(live_pptx_text):
    """Честные значения доставлены: ESS из bulk/tail (640), реальный период из
    дат проекта (2023–2025, месячная частота), реальное число наблюдений (36),
    честный источник данных."""
    assert '640' in live_pptx_text, 'ESS (min bulk/tail) не доставлен в отчёт'
    assert '2023' in live_pptx_text and '2025' in live_pptx_text, \
        'Реальный период данных не доставлен'
    assert 'Ежемесячно' in live_pptx_text, 'Частота из реальных дат не доставлена'
    assert '36' in live_pptx_text, 'Реальное число наблюдений не доставлено'
    assert 'файл, загруженный в проект' in live_pptx_text, \
        'Честная строка источника данных отсутствует'
    assert '90% HDI' in live_pptx_text, 'Честный уровень CI (90% HDI) не заявлен'


def test_preview_mode_keeps_wireframe():
    """Preview (build без data) — легитимная демка: wireframe-контент сохранён,
    live-гейт его НЕ ломает (back-compat превью/витрины)."""
    from aurora_pptx.builder import AuroraPPTXBuilder
    b = AuroraPPTXBuilder(data=None, lang='ru')
    assert b.is_live is False
    assert b.ess_min == 1247  # wireframe-дефолт жив в preview
    assert b.data_window_label == 'W01 W13 2026'


if __name__ == '__main__':
    sys.exit(pytest.main([__file__, '-q']))
