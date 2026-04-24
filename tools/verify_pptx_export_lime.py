"""
Smoke test for Aurora Hybrid signature-lime in REAL Econometrica PPTX engine
(sidecar/econometrica/engines/pptx_export.py::_add_title_text).

This is the actual generator used by Econometrica Report flow — different from
pptx_pipeline.py (which is for media-analyst cabinet).

Run:
    python tools/verify_pptx_export_lime.py
"""
from pathlib import Path
import sys

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT / "sidecar" / "econometrica"))

from engines.pptx_export import (  # noqa: E402
    _add_title_text,
    SIGNATURE_LIME,
    AURORA_DARK,
    _set_slide_bg,
)
from pptx import Presentation  # noqa: E402
from pptx.util import Inches  # noqa: E402


def count_lime(prs):
    found = []
    for idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            try:
                if shape.line.color.rgb == SIGNATURE_LIME:
                    found.append((idx, shape.left, shape.top, shape.width))
            except Exception:
                continue
    return found


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[-1]

    # Slide 1: default signature=True (should get lime)
    s1 = prs.slides.add_slide(blank)
    _set_slide_bg(s1, AURORA_DARK)
    _add_title_text(s1, "TV drives 42% of incremental sales — digital 2.1× ROAS")

    # Slide 2: signature=False (should NOT get lime — e.g. cover)
    s2 = prs.slides.add_slide(blank)
    _set_slide_bg(s2, AURORA_DARK)
    _add_title_text(s2, "Cover slide — no signature here", signature=False)

    # Slide 3: default again
    s3 = prs.slides.add_slide(blank)
    _set_slide_bg(s3, AURORA_DARK)
    _add_title_text(s3, "Recommendation: shift 15% TV → YouTube = +23% ROAS")

    out = Path(__file__).parent / "test_pptx_export_lime.pptx"
    prs.save(out)
    print(f"Wrote: {out}")

    found = count_lime(prs)
    print(f"\nLime connectors found: {len(found)}")
    for idx, left, top, width in found:
        print(
            f"  slide {idx}: left={left/914400:.3f}\" "
            f"top={top/914400:.3f}\" width={width/914400:.3f}\""
        )

    # Expected: 2 lime (slides 0 and 2), 0 on slide 1 (signature=False)
    slide_indices = sorted({idx for idx, *_ in found})
    if slide_indices == [0, 2]:
        print("\n[PASS] lime on slides 0,2; slide 1 (cover, signature=False) correctly skipped")
    else:
        print(f"\n[FAIL] expected lime on [0, 2], got {slide_indices}")
        sys.exit(1)


if __name__ == "__main__":
    main()
