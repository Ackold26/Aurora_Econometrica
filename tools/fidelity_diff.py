"""Аудит #12 «отчёты = программа»: fidelity-diff harness + регрессионный гейт.

Сверяет НАБОР факторов timeline-декомпозиции и ключевые числа каждого отчёта
(HTML, PPTX) против источника истины (decomposition.json) и против канонического
backend-поля `decomposition_series`. Цель: программа и ВСЕ отчёты показывают
ОДИН И ТОТ ЖЕ набор факторов (медиа + вынесенные signed/holiday), числа сходятся.

Контекст (2026-06-07): до фикса все 3 билдера рисовали timeline только как
baseline+media и игнорировали `signed_factor_contributions`, которые программа
(ChannelTimeline) выносит полосами → отчёты показывали МЕНЬШЕ факторов.

Запуск:
  python tools/fidelity_diff.py                 # тестовый Кагоцел 0706
  python tools/fidelity_diff.py "<results_dir>" # конкретный проект (папка results)
  python tools/fidelity_diff.py --all           # все проекты
Возврат: код 0 если все проверки PASS, иначе 1 (для CI-гейта).
"""
from __future__ import annotations
import sys, os, re, json, tempfile

HERE = os.path.dirname(os.path.abspath(__file__))
REPO = os.path.dirname(HERE)
SID = os.path.join(REPO, 'sidecar', 'econometrica')
sys.path.insert(0, SID)

PROJECTS_ROOT = os.path.join(os.environ.get('APPDATA', ''), 'aurora-econometrica-gui', 'projects')
TEST_PROJECT = os.path.join(
    PROJECTS_ROOT, 'кагоцел-рф--данные-для-эконометрики---на-ммх-0706-26', 'results')

# Типы signed-факторов, которые программа (ChannelTimeline) выносит отдельными
# полосами. positive_control (запросы/дистрибуция) остаётся внутри baseline.
BREAKOUT_TYPES = {'signed_competitor', 'signed_price', 'signed_weather', 'signed_macro', 'holiday'}
EPS = 1e-6


def norm(s: str) -> str:
    """Нормализация имени для сравнения наборов (newline/пробелы/регистр)."""
    return re.sub(r'\s+', ' ', str(s or '')).strip().lower()


def is_zero_series(per_period) -> bool:
    return not any(abs(float(v or 0)) > EPS for v in (per_period or []))


def expected_factor_set(decomp: dict) -> tuple[set, set, dict]:
    """Возвращает (media_names_norm, breakout_factor_names_norm, raw_meta)."""
    media = {norm(c.get('name')) for c in decomp.get('channels', []) if c.get('name')}
    sfc = decomp.get('signed_factor_contributions') or {}
    factors = set()
    raw = {}
    for name, fact in sfc.items():
        if not isinstance(fact, dict):
            continue
        if fact.get('type') not in BREAKOUT_TYPES:
            continue
        if is_zero_series(fact.get('per_period')):
            continue
        factors.add(norm(name))
        raw[norm(name)] = name
    return media, factors, raw


def contains_any(name_norm: str, haystack_norms: list[str]) -> bool:
    """factor присутствует если его имя — подстрока какого-то имени серии (renderer
    может префиксовать «Праздники: ...» или обрезать длинное имя в легенде «…»)."""
    for h in haystack_norms:
        if not h:
            continue
        if name_norm in h or h in name_norm:
            return True
        # Обрезанная подпись легенды (заканчивается на …/...): префиксное совпадение.
        if h.endswith('…') or h.endswith('...'):
            stem = h.rstrip('….').rstrip()
            if stem and name_norm.startswith(stem):
                return True
    return False


# ── загрузка артефактов проекта ───────────────────────────────────────────
def load_artifacts(results_dir: str) -> dict:
    out = {}
    for key, fname in (('decompose', 'decomposition.json'),
                       ('model', 'model-diagnostics.json'),
                       ('optimize', 'optimization.json')):
        p = os.path.join(results_dir, fname)
        out[key] = json.load(open(p, encoding='utf-8')) if os.path.exists(p) else {}
    return out


# ── проверка A: backend SSOT (decomposition_series) ───────────────────────
def check_backend(decomp: dict) -> dict:
    media, factors, _ = expected_factor_set(decomp)
    series = decomp.get('decomposition_series')
    if not isinstance(series, dict) or not series.get('series'):
        return {'pass': False, 'detail': 'нет поля decomposition_series',
                'missing': sorted(media | factors)}
    names = [norm(s.get('name')) for s in series['series']]
    roles = [s.get('role') for s in series['series']]
    has_baseline = 'baseline' in roles
    present = names
    missing = [m for m in sorted(media | factors) if not contains_any(m, present)]
    # тождество per-period: Σ(всех серий) == baseline_orig + media_orig (источник)
    ident_ok, ident_detail = _check_identity(decomp, series)
    ok = has_baseline and not missing and ident_ok
    detail = []
    if not has_baseline:
        detail.append('нет baseline-серии')
    if missing:
        detail.append(f'нет факторов: {missing}')
    if not ident_ok:
        detail.append(f'тождество нарушено: {ident_detail}')
    return {'pass': ok, 'detail': '; '.join(detail) or 'ok', 'missing': missing,
            'n_series': len(series['series'])}


def _check_identity(decomp: dict, series: dict) -> tuple[bool, str]:
    """Σ всех per-period серий == исходный total (baseline_full + media) по периодам."""
    ts = decomp.get('time_series') or {}
    base = ts.get('baseline') or []
    chans = ts.get('channels') or {}
    n = len(base)
    if not n:
        return True, 'нет time_series'
    orig = [float(base[t]) for t in range(n)]
    for arr in chans.values():
        for t in range(min(n, len(arr))):
            orig[t] += float(arr[t] or 0)
    got = [0.0] * n
    for s in series.get('series', []):
        data = s.get('data') or []
        for t in range(min(n, len(data))):
            got[t] += float(data[t] or 0)
    max_err = max((abs(orig[t] - got[t]) for t in range(n)), default=0.0)
    scale = max((abs(v) for v in orig), default=1.0) or 1.0
    rel = max_err / scale
    return rel < 1e-4, f'max_rel_err={rel:.2e}'


# ── проверка B: HTML отчёт ────────────────────────────────────────────────
def _extract_chart_data(html: str) -> dict:
    idx = html.index('{', html.index('var CHART_DATA'))
    obj, _ = json.JSONDecoder().raw_decode(html, idx)
    return obj


def check_html(art: dict, tmpdir: str) -> dict:
    """Т3-плюс (аудит №2 Б-4): timeline двухрежимный. detail обязан нести ВСЕ
    медиа+факторы (прежняя суть полноты, Аудит #12); overview — ровно агрегаты
    SSOT-свёртки collapse_series_to_top_groups (паритет с программой)."""
    from engines.html_export import build_html
    media, factors, _ = expected_factor_set(art['decompose'])
    out = os.path.join(tmpdir, 'fidelity.html')
    res = build_html(art['model'], art['decompose'], art['optimize'], out,
                     project_id='fidelity-test')
    if res.get('status') != 'ok':
        return {'pass': False, 'detail': f"build error: {res.get('message')}"}
    html = open(out, encoding='utf-8').read()
    try:
        cd = _extract_chart_data(html)
    except Exception as e:
        return {'pass': False, 'detail': f'CHART_DATA parse fail: {e}'}
    tl = cd.get('timeline') or {}
    det = tl.get('detail') or {}
    ovr = tl.get('overview') or {}
    # 1) detail: полнота каналов и факторов.
    series_names = [norm(k) for k in (det.get('channels') or {}).keys()]
    for f in (det.get('factors') or []):
        series_names.append(norm(f.get('name') or f.get('label') or ''))
    missing = [m for m in sorted(media | factors) if not contains_any(m, series_names)]
    # 2) overview: набор агрегатов == SSOT-свёртка decomposition_series.
    detail_msgs = []
    ds_series = ((art['decompose'].get('decomposition_series') or {}).get('series')) or []
    if ds_series:
        from engines.decomposer import collapse_series_to_top_groups
        expected_groups = {norm(c['name']) for c in collapse_series_to_top_groups(ds_series)}
        got_groups = set()
        if ovr.get('baseline'):
            got_groups.add(norm(ovr.get('baseline_label') or ''))
        got_groups |= {norm(k) for k in (ovr.get('channels') or {}).keys()}
        got_groups |= {norm(f.get('name') or '') for f in (ovr.get('factors') or [])}
        if expected_groups != got_groups:
            detail_msgs.append(
                f'overview-группы разошлись: ждали {sorted(expected_groups)}, '
                f'получили {sorted(got_groups)}')
    if missing:
        detail_msgs.append(f'нет в detail: {missing}')
    ok = not missing and not any('разошлись' in m for m in detail_msgs)
    return {'pass': ok, 'detail': '; '.join(detail_msgs) or 'ok',
            'missing': missing, 'n_series': len(series_names)}


# ── проверка C: PPTX отчёт ────────────────────────────────────────────────
def check_pptx(art: dict, tmpdir: str) -> dict:
    """Т3-плюс (аудит №2 Б-4): PPTX-timeline по канону — ОБЗОР 4 верхних групп
    (паритет с дефолтом программы), детализация вкладов живёт в waterfall и
    таблице каналов. Проверяем: каждая группа SSOT-свёртки присутствует серией
    среди чартов деки (старый контракт «каждый канал/фактор серией timeline»
    устарел вместе с переходом на свёрнутый обзор)."""
    from engines.pptx_export import build_pptx
    from pptx import Presentation
    out = os.path.join(tmpdir, 'fidelity.pptx')
    res = build_pptx(art['model'], art['decompose'], art['optimize'], out,
                     project_id='fidelity-test')
    if res.get('status') != 'ok':
        return {'pass': False, 'detail': f"build error: {res.get('message')}"}
    prs = Presentation(out)
    series_names = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_chart:
                continue
            for plot in shape.chart.plots:
                for s in plot.series:
                    nm = getattr(s, 'name', None)
                    if nm:
                        series_names.append(norm(nm))
    series_names = [s for s in series_names if s]
    ds_series = ((art['decompose'].get('decomposition_series') or {}).get('series')) or []
    if not ds_series:
        return {'pass': True, 'detail': 'нет decomposition_series (legacy) — скип',
                'missing': [], 'n_series': len(set(series_names))}
    from engines.decomposer import collapse_series_to_top_groups
    expected = [norm(c['name']) for c in collapse_series_to_top_groups(ds_series)]
    # Аудит №3 предложение 2: группы — СТРОГОЕ равенство имён (короткие «база»/
    # «медиа» подстрочным contains_any ловили бы ложный зелёный через случайное
    # вхождение в имена каналов других чартов). contains_any остаётся только для
    # детальных факторов HTML (там реальны префиксы/обрезания легенд).
    got = set(series_names)
    missing = [g for g in expected if g not in got]
    return {'pass': not missing,
            'detail': f'нет групп: {missing}' if missing else 'ok',
            'missing': missing, 'n_series': len(set(series_names))}


# ── оркестрация ───────────────────────────────────────────────────────────
def run_project(results_dir: str) -> bool:
    name = os.path.basename(os.path.dirname(results_dir))
    art = load_artifacts(results_dir)
    if not art['decompose']:
        print(f"SKIP {name}: нет decomposition.json")
        return True
    media, factors, raw = expected_factor_set(art['decompose'])
    print(f"\n=== {name} ===")
    print(f"  ожидаемый набор timeline: baseline + {len(media)} медиа + "
          f"{len(factors)} вынесенных факторов {sorted(raw.values())}")
    results = {}
    with tempfile.TemporaryDirectory() as td:
        results['backend(decomposition_series)'] = check_backend(art['decompose'])
        results['HTML'] = check_html(art, td)
        results['PPTX'] = check_pptx(art, td)
    all_ok = True
    for label, r in results.items():
        flag = 'PASS' if r['pass'] else 'FAIL'
        if not r['pass']:
            all_ok = False
        extra = f" [{r.get('n_series')} серий]" if r.get('n_series') is not None else ''
        print(f"  [{flag}] {label}{extra}: {r['detail']}")
    return all_ok


def main():
    args = [a for a in sys.argv[1:]]
    if args and args[0] == '--all':
        dirs = [os.path.join(PROJECTS_ROOT, d, 'results') for d in os.listdir(PROJECTS_ROOT)
                if os.path.isdir(os.path.join(PROJECTS_ROOT, d, 'results'))]
    elif args:
        dirs = [args[0]]
    else:
        dirs = [TEST_PROJECT]
    ok = True
    for d in dirs:
        ok = run_project(d) and ok
    print(f"\n{'='*50}\nИТОГО: {'ВСЕ PASS ✓' if ok else 'ЕСТЬ FAIL ✗'}")
    sys.exit(0 if ok else 1)


if __name__ == '__main__':
    main()
