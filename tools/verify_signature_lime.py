"""
One-off smoke test for Aurora Hybrid signature-lime under action-title.

Purpose: verify that `_add_signature_lime()` in pptx_pipeline.py actually
produces a 2pt #CCFF00 connector line below the title textbox on a generated
PPTX slide, WITHOUT needing to launch the full Tauri app.

Run from Aurora_Econometrica project root:
    python tools/verify_signature_lime.py

Output:
    tools/test_signature_lime.pptx - open in PowerPoint to inspect visually
    stdout - programmatic check that lime connector exists with correct color

This file is a temporary diagnostic, not part of the build. It is listed in
the next section of Standards/DOGFOOD_2026-04-24.md as verification artefact.
"""
from pathlib import Path
import sys

# Make pptx_pipeline importable
ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT / "src-tauri" / "sidecar"))

from pptx_pipeline import (  # noqa: E402
    _add_signature_lime,
    SIGNATURE_LIME_HEX,
)
from pptx import Presentation  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402

SAMPLE_TITLE = (
    "TV drives 42% of incremental sales, "
    "but digital delivers 2.1× higher marginal ROAS"
)


def build_fixture():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[-1]
    slide = prs.slides.add_slide(blank_layout)

    # Mimic the exact layout used by inject_summary_slides
    left = Inches(0.5)
    top = Inches(0.3)
    width = prs.slide_width - Inches(1)

    title_box = slide.shapes.add_textbox(left, top, width, Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = SAMPLE_TITLE
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.name = "Calibri"
    p.font.color.rgb = RGBColor(0x1E, 0x3A, 0x5F)

    # The function we are actually testing
    _add_signature_lime(slide, title_box)

    # Add a tiny body to contextualize position
    body_box = slide.shapes.add_textbox(
        left, Inches(1.3), width, Inches(5.0)
    )
    body_tf = body_box.text_frame
    body_tf.word_wrap = True
    body_tf.paragraphs[0].text = (
        "42% TV contribution peaks at GRP 250\n"
        "Digital mROAS 2.4:1 vs TV 1.1:1\n"
        "Reallocating 15% TV → YouTube yields +23% ROAS"
    )
    for para in body_tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(14)
            run.font.name = "Calibri"

    return prs


def verify(prs):
    expected_rgb = RGBColor(*SIGNATURE_LIME_HEX)
    found = []
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            try:
                rgb = shape.line.color.rgb
            except Exception:
                continue
            try:
                if rgb == expected_rgb:
                    found.append({
                        "slide": slide_idx,
                        "shape_id": shape.shape_id,
                        "left": shape.left,
                        "top": shape.top,
                        "width": shape.width,
                        "height": shape.height,
                        "line_width_emu": shape.line.width,
                    })
            except Exception:
                continue
    return found


def main():
    out_path = Path(__file__).resolve().parent / "test_signature_lime.pptx"

    prs = build_fixture()
    prs.save(out_path)
    print(f"Wrote: {out_path}")
    print("Open in PowerPoint and visually verify:")
    print("  • 2pt lime (#CCFF00) horizontal line directly under action-title")
    print("  • Line spans full title width (left + width)")
    print("  • Line is positioned between title (bottom ~1.1\") and body (starts 1.3\")")
    print()

    found = verify(prs)
    if not found:
        print("❌ FAIL - no lime-colored connector found on any slide")
        sys.exit(1)

    print(f"✅ PASS - found {len(found)} lime connector(s):")
    for f in found:
        # EMU units: 914400 EMU = 1 inch
        def emu_in(v):
            return f"{v / 914400:.3f}\""
        print(
            f"  slide {f['slide']}: "
            f"pos=({emu_in(f['left'])}, {emu_in(f['top'])}), "
            f"size=({emu_in(f['width'])} × {emu_in(f['height'])}), "
            f"line_width={f['line_width_emu'] / 12700:.1f}pt"
        )


if __name__ == "__main__":
    main()
