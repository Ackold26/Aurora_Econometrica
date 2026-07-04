"""
Unit tests for engines.narrative_adapter helpers.

Covers the post-audit invariants (2026-04-25):
- compute_report_id: deterministic + HTML/PPTX unification.
- derive_action_headline: zero-effect guard, negative-lift handling,
  strict-majority underperf threshold, all 4 slide hints.
- _normalize_channel_name: Excel column-header stripping.
- _sanitize_project_slug: internal-marker cleanup.
- _merge_channels: collision detection and drop-logging.
- derive_verdict: 5-way classification edge cases.

Run:
    cd sidecar && python ../tools/test_narrative_adapter.py

Exit code 0 on success, 1 on any failure. Plain stdlib - no pytest dep.
"""
from __future__ import annotations

import logging
import sys
from pathlib import Path

try:
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
except Exception:
    pass

REPO = Path(__file__).resolve().parents[1]
SIDECAR = REPO / "sidecar"
sys.path.insert(0, str(SIDECAR))
sys.path.insert(0, str(SIDECAR / "econometrica"))


PASSED = 0
FAILED = 0


def _ok(label: str) -> None:
    global PASSED
    PASSED += 1
    print(f"[OK]   {label}")


def _fail(label: str, detail: str = "") -> None:
    global FAILED
    FAILED += 1
    line = f"[FAIL] {label}"
    if detail:
        line += f" - {detail}"
    print(line)


def assert_eq(label: str, actual, expected) -> None:
    if actual == expected:
        _ok(label)
    else:
        _fail(label, f"got {actual!r}, expected {expected!r}")


def assert_true(label: str, cond: bool, detail: str = "") -> None:
    if cond:
        _ok(label)
    else:
        _fail(label, detail)


def _reset_logger():
    pass  # placeholder for finally restoration


def assert_in(label: str, substr: str, text: str) -> None:
    if substr in (text or ""):
        _ok(label)
    else:
        _fail(label, f"{substr!r} not in {text!r}")


def assert_not_in(label: str, substr: str, text: str) -> None:
    if substr not in (text or ""):
        _ok(label)
    else:
        _fail(label, f"{substr!r} leaked into {text!r}")


def test_compute_report_id():
    from engines.narrative_adapter import compute_report_id

    # Deterministic: same inputs → same ID.
    a = compute_report_id("X", "P1", [{"name": "TV", "spend": 100, "contribution": 150, "verdict": "Hold"}], {"mqs": 80})
    b = compute_report_id("X", "P1", [{"name": "TV", "spend": 100, "contribution": 150, "verdict": "Hold"}], {"mqs": 80})
    assert_eq("compute_report_id deterministic", a, b)

    # Format: aurora-mmm-{12hex}.
    assert_true("compute_report_id format", a.startswith("aurora-mmm-") and len(a) == len("aurora-mmm-xxxxxxxxxxxx"))

    # Different client → different ID.
    c = compute_report_id("Y", "P1", [{"name": "TV", "spend": 100, "contribution": 150, "verdict": "Hold"}], {"mqs": 80})
    assert_true("compute_report_id client changes ID", a != c)

    # Different diagnostic key set → different ID.
    d = compute_report_id("X", "P1", [{"name": "TV", "spend": 100, "contribution": 150, "verdict": "Hold"}], {"mqs": 80, "r2": 0.9})
    assert_true("compute_report_id diagnostics expansion changes ID", a != d)

    # Empty inputs → still deterministic, still well-formed.
    e = compute_report_id(None, None, None, None)
    assert_true("compute_report_id accepts None args", e.startswith("aurora-mmm-"))

    # Float drift tolerance - 0.872 vs 0.8724 should round to same bucket (3dp).
    f1 = compute_report_id("X", "P", [], {"mqs": 0.872})
    f2 = compute_report_id("X", "P", [], {"mqs": 0.8724})
    assert_eq("compute_report_id 3dp rounding matches", f1, f2)

    # Channel order invariance (unsorted input → same ID as sorted input).
    chs_a = [
        {"name": "TV", "spend": 100, "contribution": 150, "verdict": "Hold"},
        {"name": "Digital", "spend": 50, "contribution": 90, "verdict": "Scale"},
    ]
    chs_b = list(reversed(chs_a))
    i1 = compute_report_id("X", "P", chs_a, {"mqs": 80})
    i2 = compute_report_id("X", "P", chs_b, {"mqs": 80})
    assert_eq("compute_report_id channel-order invariant", i1, i2)

    # Non-numeric diagnostic values don't break (e.g. tier_label string).
    g = compute_report_id("X", "P", [], {"mqs": 80, "tier_label": "GOOD"})
    assert_true("compute_report_id handles mixed-type diagnostics", g.startswith("aurora-mmm-"))


def test_compute_report_id_html_pptx_parity():
    """End-to-end: same builder data → same Report ID in HTML and PPTX."""
    from engines.narrative_adapter import _map_pipeline_to_builder_data
    from aurora_pptx import build_pptx as _ppx_build
    from aurora_html import build_html as _html_build
    import re

    model = {"diagnostics": {"mqs": {"score": 82}, "metrics": {"r_squared": 0.88, "mape_pct": 9.1}}}
    decomp = {"channels": [
        {"name": "TV", "spend": 120e6, "contribution": 180e6, "roi": 1.5},
        {"name": "Digital", "spend": 65e6, "contribution": 110e6, "roi": 1.7},
    ]}
    opt = {"channels": [
        {"name": "TV", "current_spend": 120e6, "optimal_spend": 90e6, "mroi_current": 1.2},
        {"name": "Digital", "current_spend": 65e6, "optimal_spend": 100e6, "mroi_current": 2.0},
    ], "expected_lift_pct": 8}

    data = _map_pipeline_to_builder_data(model, decomp, opt, scenarios=None, project_id="Parity")

    prs = _ppx_build(data=data)
    tmp = REPO / "tools" / "_parity.pptx"
    prs.save(str(tmp))
    from pptx import Presentation
    p = Presentation(str(tmp))
    pptx_xml = "".join(s.shapes._spTree.xml for s in p.slides)
    m_ppx = re.search(r"aurora-mmm-[a-f0-9]{12}", pptx_xml)
    import os
    try:
        os.remove(tmp)
    except OSError:
        pass

    html = _html_build(data=data)
    m_html = re.search(r"aurora-mmm-[a-f0-9]{12}", html)

    assert_true("PPTX emits Report ID", m_ppx is not None)
    assert_true("HTML emits Report ID", m_html is not None)
    if m_ppx and m_html:
        assert_eq("PPTX Report ID == HTML Report ID", m_ppx.group(0), m_html.group(0))


def test_derive_action_headline_all_hints():
    from engines.narrative_adapter import derive_action_headline
    chs = [{"name": "A", "mroas": 1.8}, {"name": "B", "mroas": 1.5}]
    facts = {
        "leader_channel": "A", "hero_channel": "B",
        "expected_lift_pct": 8.0, "reallocation_mln": 20,
        "underperformer_names": [],
    }
    mroas = derive_action_headline(chs, facts, "mroas")
    portfolio = derive_action_headline(chs, facts, "portfolio")
    timeline = derive_action_headline(chs, facts, "timeline")
    scqar = derive_action_headline(chs, facts, "scqar")

    assert_in("mroas headline mentions hero", "B", mroas or "")
    assert_in("mroas headline mentions leader", "A", mroas or "")
    assert_in("scqar headline quantifies lift", "+8 пп", scqar or "")
    assert_in("timeline headline mentions leader", "A", timeline or "")
    assert_true("portfolio headline non-empty", bool(portfolio))


def test_derive_action_headline_zero_effect_guard():
    """Negative or weak lift should NOT quantify an improvement."""
    from engines.narrative_adapter import derive_action_headline
    chs = [{"name": "A", "mroas": 1.8}, {"name": "B", "mroas": 1.5}]

    for bad_lift in (None, -1.5, -0.3, 0, 0.1, 0.4):
        facts = {
            "leader_channel": "A", "hero_channel": "B",
            "expected_lift_pct": bad_lift, "reallocation_mln": 20,
            "underperformer_names": [],
        }
        scqar = derive_action_headline(chs, facts, "scqar") or ""
        mroas = derive_action_headline(chs, facts, "mroas") or ""
        # Neither string should contain a "+N пп к ROAS" promise
        assert_not_in(f"scqar no fake promise (lift={bad_lift})", " пп к ROAS", scqar)
        assert_not_in(f"mroas no fake promise (lift={bad_lift})", " пп к ROAS", mroas)
        # No literal broken "+-" formatting
        assert_not_in(f"scqar no '+-' broken formatting (lift={bad_lift})", "+-", scqar)
        assert_not_in(f"mroas no '+-' broken formatting (lift={bad_lift})", "+-", mroas)


def test_derive_action_headline_underperf_threshold():
    """Risk scenario only fires on strict majority, not 1 of N."""
    from engines.narrative_adapter import derive_action_headline

    # 3-channel portfolio, 1 underperformer (33%) → NOT risk.
    chs3 = [{"name": c, "mroas": 1.5} for c in ("A", "B", "C")]
    facts1 = {"leader_channel": "A", "hero_channel": "A", "expected_lift_pct": 3,
              "reallocation_mln": 0, "underperformer_names": ["C"]}
    h = derive_action_headline(chs3, facts1, "scqar") or ""
    assert_not_in("3-ch 1-underperf scqar is NOT risk", "Сократить C и сфокусировать", h)

    # 4-channel, 2 underperf (50%) → risk.
    chs4 = [{"name": c, "mroas": 1.5} for c in ("A", "B", "C", "D")]
    facts2 = {"leader_channel": "A", "hero_channel": "A", "expected_lift_pct": 3,
              "reallocation_mln": 0, "underperformer_names": ["C", "D"]}
    h2 = derive_action_headline(chs4, facts2, "scqar") or ""
    assert_in("4-ch 2-underperf scqar IS risk", "Сократить", h2)

    # 2-channel, 1 underperf (50%) - should NOT trigger because floor=2.
    chs2 = [{"name": c, "mroas": 1.5} for c in ("A", "B")]
    facts3 = {"leader_channel": "A", "hero_channel": "A", "expected_lift_pct": 0,
              "reallocation_mln": 0, "underperformer_names": ["B"]}
    h3 = derive_action_headline(chs2, facts3, "scqar") or ""
    assert_not_in("2-ch 1-underperf scqar is NOT risk (floor=2)", "Сократить B и", h3)


def test_normalize_channel_name():
    from engines.narrative_adapter import _normalize_channel_name

    assert_eq("strip 'Бюджет до НДС'",
              _normalize_channel_name("Performance Бюджет до НДС"), "Performance")
    assert_eq("strip 'ДО НДС до АК'",
              _normalize_channel_name("Banners Бюджет ДО НДС до АК"), "Banners")
    assert_eq("empty after strip → None",
              _normalize_channel_name("Бюджет до НДС"), None)
    assert_eq("preserve audience quantifier",
              _normalize_channel_name("TRPs бренд (W 25-50)"), "TRPs бренд (W 25-50)")
    assert_eq("no-op on clean name",
              _normalize_channel_name("TV"), "TV")
    assert_eq("None input → None",
              _normalize_channel_name(None), None)
    assert_eq("empty string → None",
              _normalize_channel_name(""), None)
    # Аудит №3 предложение 1: агрегатные слова итого/всего/сумма/total —
    # суммарные колонки → None (гейт валидации + дроп из таблицы), реальные
    # каналы с этими словами в составе имени сохраняют инструмент.
    assert_eq("'ИТОГО Бюджет' → None (агрегат)",
              _normalize_channel_name("ИТОГО Бюджет"), None)
    assert_eq("'Total Бюджет' → None (агрегат)",
              _normalize_channel_name("Total Бюджет"), None)
    assert_eq("'Всего' → None (агрегат)",
              _normalize_channel_name("Всего"), None)
    assert_eq("'Сумма, млн' → None (агрегат)",
              _normalize_channel_name("Сумма, млн"), None)
    assert_eq("'Total TV' → 'TV' (канал цел)",
              _normalize_channel_name("Total TV"), "TV")
    assert_eq("'Диджитал всего' → 'Диджитал' (под-агрегат виден инструментом)",
              _normalize_channel_name("Диджитал всего"), "Диджитал")


def test_sanitize_project_slug():
    from engines.narrative_adapter import _sanitize_project_slug

    label, code = _sanitize_project_slug("mmx-2021-2025-исходник-ммх-2404-26--4")
    assert_not_in("slug: no 'исходник' leak", "исходник", label)
    assert_not_in("slug: no 'ммх' leak", "ммх", label)
    assert_not_in("slug: no '--' leak", "--", label)

    label2, _ = _sanitize_project_slug("венарус-ммх-2404-26--2")
    assert_in("slug: preserves client name", "Венарус", label2)

    label3, code3 = _sanitize_project_slug(None)
    assert_eq("slug None → ('Client','PROJECT')", (label3, code3), ("Client", "PROJECT"))

    label4, _ = _sanitize_project_slug("Kagocel")
    assert_eq("slug: simple name preserved", label4, "Kagocel")


def test_merge_channels_collision_warning(caplog_list: list | None = None):
    """When two columns normalize to the same key, at least log the situation."""
    import engines.narrative_adapter as na

    # Capture logs
    records = []
    class _Handler(logging.Handler):
        def emit(self, record):
            records.append(record)
    h = _Handler()
    h.setLevel(logging.WARNING)
    na.logger.addHandler(h)
    prev_level = na.logger.level
    na.logger.setLevel(logging.WARNING)
    try:
        # Two columns both normalize to "TV"
        decomp = [
            {"name": "TV Бюджет до НДС", "spend": 100, "contribution": 150, "roi": 1.5},
            {"name": "TV Бюджет ДО НДС до АК", "spend": 90, "contribution": 140, "roi": 1.4},
        ]
        merged = na._merge_channels(decomp, None)
        collision_msgs = [r.getMessage() for r in records]
        warned = any(
            "collapse" in m.lower() or "collision" in m.lower() or
            "duplicate" in m.lower() or "same normalized key" in m.lower()
            for m in collision_msgs
        )
        # Acceptable behaviour: either keep both rows OR drop-with-warning.
        assert_true("merge: collision warns (first wins)",
                    warned, detail=f"warnings={collision_msgs!r}")
        # When duplicates collapse, only first survives - confirm no silent
        # data loss by checking first channel's name survived
        assert_true("merge: collision first wins",
                    len(merged) == 1 and merged[0]["contribution"] == 150)
    finally:
        na.logger.removeHandler(h)
        na.logger.setLevel(prev_level)


def test_derive_verdict():
    from engines.narrative_adapter import derive_verdict

    # Scale: high mROAS + increase spend
    assert_eq("verdict Scale", derive_verdict({
        "current_spend": 10, "optimal_spend": 20, "mroas": 2.0}), "Scale")
    # Cut: low mROAS
    assert_eq("verdict Cut (low mroas)", derive_verdict({
        "current_spend": 10, "optimal_spend": 5, "mroas": 0.3}), "Cut")
    # Cut: severe ratio
    assert_eq("verdict Cut (severe cut)", derive_verdict({
        "current_spend": 10, "optimal_spend": 2, "mroas": 1.5}), "Cut")
    # Reduce: profitable but trimmed
    assert_eq("verdict Reduce", derive_verdict({
        "current_spend": 10, "optimal_spend": 7, "mroas": 1.4}), "Reduce")
    # Hold: stable and profitable
    assert_eq("verdict Hold", derive_verdict({
        "current_spend": 10, "optimal_spend": 10, "mroas": 1.5}), "Hold")
    # Watch: borderline
    assert_eq("verdict Watch", derive_verdict({
        "current_spend": 10, "optimal_spend": 10, "mroas": 1.0}), "Watch")
    # Graceful fallback on invalid data
    assert_eq("verdict Watch fallback on bad data", derive_verdict({
        "current_spend": "oops", "optimal_spend": None, "mroas": None}), "Watch")


def main() -> int:
    print("=== test_narrative_adapter ===\n")
    test_compute_report_id()
    print()
    test_compute_report_id_html_pptx_parity()
    print()
    test_derive_action_headline_all_hints()
    print()
    test_derive_action_headline_zero_effect_guard()
    print()
    test_derive_action_headline_underperf_threshold()
    print()
    test_normalize_channel_name()
    print()
    test_sanitize_project_slug()
    print()
    test_merge_channels_collision_warning()
    print()
    test_derive_verdict()
    print()
    print(f"\n{PASSED}/{PASSED + FAILED} assertions passed.")
    return 0 if FAILED == 0 else 1


if __name__ == "__main__":
    sys.exit(main())
