"""
Pseudo-dogfood for Aurora Hybrid signature-lime under action-title.

Real Kagocel XLSX dataset не найден на диске; full live-dogfood через UI отложен
до user-attended сессии. Здесь - максимально приближенный MMM-нарратив
(Kagocel Q1 2026 context) как synthesis.md, прогнанный через реальный
inject_summary_slides() на templates/default.pptx. Получаем production-like PPTX
с множественными key-message slides и проверяем, что на каждом присутствует
lime connector.

Run:
    python tools/pseudo_dogfood_kagocel.py

Output:
    tools/pseudo_dogfood_output.pptx - артефакт для визуального audit по 07_CHECKLIST.md
    stdout - summary: сколько слайдов добавлено, сколько имеют lime connector
"""
from pathlib import Path
import json
import sys
import tempfile

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT / "src-tauri" / "sidecar"))

from pptx_pipeline import inject_summary_slides, SIGNATURE_LIME_HEX  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402


# Realistic MMM synthesis (Kagocel-inspired, Q1 2026)
# Format: ## Title\n  - bullet\n  - bullet ...
SYNTHESIS_MD = """\
## TV drives 42% of incremental sales, but digital delivers 2.1× higher marginal ROAS

- 42% TV contribution concentrates in W25-54, peaks at GRP 250
- Digital mROAS 2.4:1 vs TV 1.1:1 at current spend
- Reallocating 15% TV → YouTube yields +23% ROAS at flat budget
- Recommendation carries R² = 0.92 and MAPE 7.4%

## YouTube is under-invested: 15% share but 2.4× mROAS

- Current spend ₽38M of ₽247M envelope (15%)
- Saturation curve shows 68% headroom before diminishing returns
- Optimal share 23% (₽56M) per scenario optimizer output
- Evidence: 4 of 5 MCMC chains converge on ₽56M ±4M at HDI 95%

## Radio and OOH show negative incremental return in 3 of 8 markets

- Radio mROAS 0.65 vs blended 1.48
- OOH mROAS 0.85, contributes <5% of total sales
- Recommended reallocation: radio −33% (₽2M saved), OOH −33% (₽4M saved)
- Saved budget of ₽6M goes to YouTube tier A (+5pp share)

## Meta* platforms retain +21% allocation - creative A/B testing to protect ROAS

- Meta mROAS 1.85, above threshold 1.50
- Optimal share 14% (₽34M) per HDI 95% interval
- Risk: creative fatigue detected at week 8; rotate assets every 6 weeks
- * Meta Platforms - организация признана экстремистской и запрещена в Российской Федерации

## Recommendation holds R² 0.92 · MAPE 7.4% · HDI ±4pp

- Bayesian MMM fit: R² = 0.92, MAPE = 7.4% across 8 markets
- 95% HDI on ROAS delta: [+19%, +27%] - narrower than ±4pp
- Recommendation robust to posterior sampling (NUTS, 4 chains × 2000 draws)
- Next step: lock Q2 media plan by May 15 with owner = Media Director
"""


def count_lime_connectors(prs):
    """Count connector shapes with SIGNATURE_LIME color across all slides."""
    expected_rgb = RGBColor(*SIGNATURE_LIME_HEX)
    per_slide = {}
    for idx, slide in enumerate(prs.slides):
        per_slide[idx] = 0
        for shape in slide.shapes:
            try:
                rgb = shape.line.color.rgb
                if rgb == expected_rgb:
                    per_slide[idx] += 1
            except Exception:
                continue
    return per_slide


def main():
    tools_dir = Path(__file__).resolve().parent
    output = tools_dir / "pseudo_dogfood_output.pptx"

    template = ROOT / "sidecar" / "econometrica" / "_internal" / "pptx" / "templates" / "default.pptx"
    if not template.exists():
        print(f"❌ Template not found: {template}")
        sys.exit(1)

    with tempfile.TemporaryDirectory() as tmpdir:
        tmp = Path(tmpdir)
        synthesis_md = tmp / "synthesis.md"
        styles_json = tmp / "styles.json"
        slides_json = tmp / "slides.json"

        synthesis_md.write_text(SYNTHESIS_MD, encoding="utf-8")
        styles_json.write_text(json.dumps({
            "fonts": [{"name": "Calibri"}]  # matches inject_summary_slides default
        }), encoding="utf-8")
        slides_json.write_text(json.dumps([]), encoding="utf-8")

        inject_summary_slides(
            str(template),
            str(synthesis_md),
            str(styles_json),
            str(slides_json),
            str(output),
        )

    # Reload output and audit
    prs = Presentation(str(output))
    total_slides = len(prs.slides)
    per_slide = count_lime_connectors(prs)
    slides_with_lime = sum(1 for n in per_slide.values() if n > 0)
    template_baseline = Presentation(str(template))
    template_slides = len(template_baseline.slides)
    added_slides = total_slides - template_slides

    print(f"Template slides: {template_slides}")
    print(f"Output slides:   {total_slides}")
    print(f"Added summary slides: {added_slides}")
    print()
    print(f"Lime connector count per slide:")
    for idx, count in per_slide.items():
        tag = "ok" if count == 1 else ("extra" if count > 1 else "none")
        print(f"  slide {idx:2d}: {count} lime line(s) - {tag}")
    print()

    if slides_with_lime == added_slides:
        print(f"✅ PASS - all {added_slides} summary slides have signature-lime line")
    else:
        print(f"❌ FAIL - {added_slides} summary slides added, but only "
              f"{slides_with_lime} have lime line")
        sys.exit(1)

    print(f"\nArtifact: {output}")
    print("Open in PowerPoint for visual audit against Standards/07_CHECKLIST.md")


if __name__ == "__main__":
    main()
