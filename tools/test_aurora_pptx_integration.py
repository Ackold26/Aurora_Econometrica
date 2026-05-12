"""
PPTX golden snapshot integration tests (v1.3.2 audit followup I4).

Build full PPTX deck для 3 KPI modes (monetary roi / count / effectiveness),
extract concatenated text from all slides, assert KPI-aware labels appear
correctly. Поймало бы B1 (effectiveness не multiplied) + B4 (count показывал
units/₽ как 0 ₽/ед.) автоматически.

Requires aurora_tokens module (shim или production). Test environment shim в
sidecar/econometrica/aurora_tokens.py provides minimal hex values.

NOT byte-exact comparison — too brittle (python-pptx version drift / random
report_id / timestamps). Concentrate на user-visible TEXT and LABELS что
constitute KPI correctness contract.
"""
from __future__ import annotations

import io
import sys
from pathlib import Path

REPO = Path(__file__).resolve().parents[1]
SIDECAR = REPO / 'sidecar'
sys.path.insert(0, str(SIDECAR))
sys.path.insert(0, str(SIDECAR / 'econometrica'))


def _make_data(kpi_kind: str = 'monetary', mode: str = 'roi') -> dict:
    """Synthetic 3-channel data для PPTX build. mroas values intentional —
    raw mathematical KPI/spend. Для count это units/₽ (small fraction); render
    layer inverts к CPU per B4 audit fix.
    """
    base = {
        'meta': {
            'client': 'TestCo', 'project_id': 'TEST-Q1-2026',
            'version': '1.3.2', 'report_date': '2026-Q2',
            'period_label': 'Q1 2026', 'data_window_label': 'W01-W13 2026',
        },
        'diagnostics': {
            'mqs_score': 75, 'mqs_tier_label': 'Хорошее',
            'r_squared': 0.85, 'mape_pct': 7.5,
            'r_hat_max': 1.01, 'ess_min': 1200,
        },
        'channels': [
            # mroas values для count = units/₽ (~0.01 range), не CPU
            # = mroas values для effectiveness = share fraction (0..0.6)
            {
                'name': 'TV',
                'spend': 50_000_000, 'contribution': 30_000_000,
                'mroas': 0.012 if kpi_kind == 'count' else (0.45 if mode == 'effectiveness' else 1.5),
                'verdict': 'Hold', 'action': 'Hold',
                'action_label': 'Удержать', 'action_priority': 4,
                'action_reasoning': 'Метрика в стабильном диапазоне.',
                'mroas_ci_low': 0.010 if kpi_kind == 'count' else (0.35 if mode == 'effectiveness' else 1.3),
                'mroas_ci_high': 0.015 if kpi_kind == 'count' else (0.55 if mode == 'effectiveness' else 1.7),
                'current_spend': 50_000_000, 'optimal_spend': 50_000_000,
            },
            {
                'name': 'Digital',
                'spend': 30_000_000, 'contribution': 25_000_000,
                'mroas': 0.008 if kpi_kind == 'count' else (0.30 if mode == 'effectiveness' else 0.8),
                'verdict': 'Cut', 'action': 'Cut',
                'action_label': 'Сократить', 'action_priority': 0,
                'action_reasoning': 'Метрика ниже breakeven.',
                'mroas_ci_low': 0.006 if kpi_kind == 'count' else (0.20 if mode == 'effectiveness' else 0.6),
                'mroas_ci_high': 0.010 if kpi_kind == 'count' else (0.40 if mode == 'effectiveness' else 1.0),
                'current_spend': 30_000_000, 'optimal_spend': 20_000_000,
            },
            {
                'name': 'Social',
                'spend': 20_000_000, 'contribution': 15_000_000,
                'mroas': 0.011 if kpi_kind == 'count' else (0.25 if mode == 'effectiveness' else 1.0),
                'verdict': 'Watch', 'action': 'Watch',
                'action_label': 'Наблюдать', 'action_priority': 2,
                'action_reasoning': 'Стабильный канал.',
                'mroas_ci_low': 0.009 if kpi_kind == 'count' else (0.15 if mode == 'effectiveness' else 0.8),
                'mroas_ci_high': 0.013 if kpi_kind == 'count' else (0.35 if mode == 'effectiveness' else 1.2),
                'current_spend': 20_000_000, 'optimal_spend': 20_000_000,
            },
        ],
        'narrative_facts': {
            'total_budget_mln': 100, 'total_contrib_mln': 70,
            # weighted_roi backend convention = total_contrib / total_spend
            # для count это units/₽ (small), для monetary это ROI (>1), для
            # effectiveness не используется напрямую.
            'weighted_roi': 0.7 if kpi_kind == 'monetary' else 0.0125,
            'leader_channel': 'TV', 'hero_channel': 'TV',
            'leader_share_spend_pct': 50, 'leader_share_contrib_pct': 43,
            'reallocation_mln': 10, 'expected_lift_pct': 5,
            'optimization_converged': True, 'binding_constraints': False,
            'cut_source_channel': 'Digital', 'scale_destination_channel': 'TV',
            'underperformer_names': [],
            'n_active_channels': 3,
            'action_counts': {'Scale': 0, 'Hold': 1, 'Cut': 1, 'Watch': 1,
                              'Reduce': 0, 'Uncertain': 0},
            'media_contribution_pct': 50, 'baseline_pct': 30,
            'honest_narrative': False,
            'budget_dominator_channel': 'TV',
            'budget_dominator_spend_pct': 50, 'budget_dominator_contrib_pct': 43,
        },
        'kpi': {
            'kpi_kind': kpi_kind, 'derived_mode': mode,
            'value_per_count_unit': 80.0 if kpi_kind == 'count' else None,
            'value_per_count_unit_label': '80 ₽/упак' if kpi_kind == 'count' else '',
            'labels': {
                'metric_label': (
                    'CPU, ₽/ед.' if kpi_kind == 'count'
                    else ('Доля %' if mode == 'effectiveness' else 'ROI')
                ),
                'metric_short_label': (
                    'CPU' if kpi_kind == 'count'
                    else ('Доля' if mode == 'effectiveness' else 'ROI')
                ),
                'target_unit_label': 'упак / ед.' if kpi_kind == 'count' else '₽',
                'target_axis_label': 'Продажи, упак' if kpi_kind == 'count' else 'Продажи, ₽',
                'methodology_label': '',
            },
        },
    }
    return base


def _build_text_dump(data: dict) -> str:
    """Build PPTX in-memory + extract all text runs."""
    from aurora_pptx.builder import AuroraPPTXBuilder
    b = AuroraPPTXBuilder(data=data)
    b.build()
    buf = io.BytesIO()
    b.prs.save(buf)
    buf.seek(0)
    from pptx import Presentation
    prs = Presentation(buf)
    chunks = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text:
                            chunks.append(run.text)
    return '\n'.join(chunks)


# ─── Backward compat (monetary roi) ────────────────────────────────────


def test_monetary_roi_renders_mroas_labels():
    """Legacy: ROI/mROAS labels intact, value formats ×."""
    data = _make_data('monetary', 'roi')
    text = _build_text_dump(data)
    # ROI / mROAS labels present
    assert 'ROI' in text or 'mROAS' in text.upper()
    # No count-specific labels leak
    assert 'CPU' not in text
    assert '₽/ед.' not in text
    # Client identity preserved
    assert 'TestCo' in text


def test_monetary_roi_does_not_show_share_labels():
    data = _make_data('monetary', 'roi')
    text = _build_text_dump(data)
    # Effectiveness-mode labels NOT in monetary roi deck
    assert 'Средняя доля каналов' not in text
    assert 'Прогнозный прирост доли' not in text


# ─── Count KPI (CPU display) ────────────────────────────────────────────


def test_count_kpi_shows_cpu_column_header():
    """B4 audit fix: count KPI action_table column header = CPU."""
    data = _make_data('count', 'roi')
    text = _build_text_dump(data)
    assert 'CPU' in text
    assert '₽/ед.' in text


def test_count_kpi_inverts_to_cpu_values():
    """B4 audit: c.mroas=0.012 (units/₽) → CPU=1/0.012≈83 ₽/ед.

    Pre-fix: render layer showed raw mroas (e.g. «0.012») or rounded
    floor («0 ₽/ед.»). Post-fix: inverted к integer CPU before display.
    """
    data = _make_data('count', 'roi')
    text = _build_text_dump(data)
    # Raw fractions из fixture (0.012, 0.008, 0.011) НЕ должны утечь as-is.
    assert '0.012' not in text
    assert '0.008' not in text
    assert '0.011' not in text
    # CPU values от 1/raw_mroas — TV: 1/0.012≈83, Digital: 1/0.008=125, Social: 1/0.011≈91.
    # Хотя бы один из CPU values должен appear в action_table cell или commentary.
    assert ('83' in text or '91' in text or '125' in text), \
        'Expected at least one of CPU values (83/91/125) от inverted mroas'


def test_count_kpi_has_cpu_portfolio_phrase():
    """B4 audit: weighted_roi=0.0125 units/₽ → CPU портфеля 80 ₽/ед."""
    data = _make_data('count', 'roi')
    text = _build_text_dump(data)
    assert 'CPU портфеля' in text
    # 1 / 0.0125 = 80
    assert '80 ₽/ед.' in text


def test_count_kpi_impact_label_is_sales_not_roas():
    """B3+M2 audit: для count «Прогнозный прирост продаж», не ROAS."""
    data = _make_data('count', 'roi')
    text = _build_text_dump(data)
    assert 'Прогнозный прирост продаж' in text
    assert 'Прогнозный ROAS' not in text


# ─── Effectiveness mode (Доля display) ──────────────────────────────────


def test_effectiveness_mode_shows_share_column():
    """Effectiveness column header = Доля эффекта."""
    data = _make_data('monetary', 'effectiveness')
    text = _build_text_dump(data)
    assert 'Доля эффекта' in text or 'ДОЛЯ ЭФФЕКТА' in text


def test_effectiveness_mode_uses_share_phrase():
    """SCQAR situation override использует share phrasing."""
    data = _make_data('monetary', 'effectiveness')
    text = _build_text_dump(data)
    assert 'Средняя доля каналов в портфеле' in text
    assert 'ROI портфеля' not in text


def test_effectiveness_impact_label_is_share():
    """Effectiveness «Прогнозный прирост доли»."""
    data = _make_data('monetary', 'effectiveness')
    text = _build_text_dump(data)
    assert 'Прогнозный прирост доли' in text
    assert 'Прогнозный ROAS' not in text


def test_effectiveness_mode_no_breakeven_metaphor():
    """M1 audit: для effectiveness «breakeven» metaphor unnatural; используется
    «низкая доля» / «вклад в долю эффекта».
    """
    data = _make_data('monetary', 'effectiveness')
    text = _build_text_dump(data)
    # Common breakeven phrasing для legacy/count — отсутствует
    assert 'mROAS < 1' not in text
    # Effectiveness-specific phrasings present (verdict reasons / action_02)
    # Note: action_02 fires only when n_saturated > 0. Synthetic data has
    # 1 'Cut' channel — qualifies. Check «доля» appears (verdict reason).
    assert 'доля' in text.lower() or 'эффект' in text.lower()


# ─── Cross-mode SCQAR recommendation phrase (M2) ───────────────────────


def test_monetary_scqar_recommendation_uses_roas():
    """M2 audit: legacy keeps «Прогнозный ROAS / прирост ROAS»."""
    data = _make_data('monetary', 'roi')
    text = _build_text_dump(data)
    # Either «прирост ROAS» (SCQAR template) или «Прогнозный ROAS» (impact card).
    assert 'ROAS' in text


def test_count_scqar_recommendation_drops_roas():
    """M2 audit: count override → «прирост продаж», ROAS removed."""
    data = _make_data('count', 'roi')
    text = _build_text_dump(data)
    # SCQAR situation phrase + impact card both adapted.
    assert 'прирост продаж' in text or 'Прогнозный прирост продаж' in text


# ─── Smoke: builder doesn't crash for edge cases ───────────────────────


def test_builder_handles_empty_channels_wireframe_mode():
    """Builder must build wireframe deck когда channels = []."""
    data = {
        'meta': {'client': 'Empty'},
        'diagnostics': {'mqs_score': 0},
        'channels': [],
        'narrative_facts': None,
        'time_series': None,
    }
    text = _build_text_dump(data)
    # Should still build basic deck с Kagocel pilot wireframe text
    assert 'Empty' in text
