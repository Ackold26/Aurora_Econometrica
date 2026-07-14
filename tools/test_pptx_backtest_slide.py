"""E1 (2026-07-03): слайд «Проверка на истории» в клиентской PPTX-деке.

Контракт: витрина (models/backtest.json, status=ok) рождает 13-й слайд после
методологии со сдвигом хвостовой нумерации; без витрины дека остаётся
12-слайдовой и НЕ содержит следов витрины (wireframe-режима у слайда нет
по построению — урок B1 «замаскированная дефолтом честность»).
"""
from __future__ import annotations

import json
import sys
from pathlib import Path

import pytest

ROOT = Path(__file__).resolve().parents[1]
SIDECAR = ROOT / 'sidecar'
for _p in (str(SIDECAR), str(SIDECAR / 'econometrica')):
    if _p not in sys.path:
        sys.path.insert(0, _p)

from test_goalseek_honesty import _build_project  # noqa: E402
from test_report_fidelity_live import _extract_all_text  # noqa: E402
from engines.narrative_adapter import _map_pipeline_to_builder_data  # noqa: E402


def _backtest_fixture(**overrides):
    bt = {
        'status': 'ok',
        'verdict': 'validated',
        'verdict_text': 'Модель подтверждена на удержанной истории: MAPE 6.1%.',
        'granularity': 'M',
        'horizon_periods': 3,
        'mode': 'ols',
        'n_windows': 4,
        'windows_hit_total': 4,
        'windows_with_interval': 4,
        'coverage_per_window': 1.0,
        'coverage_per_period': 0.9167,
        'n_holdout_points': 12,
        'n_holdout_points_with_interval': 12,
        'mape_model': 6.1,
        'naive_mape': {'naive_last': 9.4, 'seasonal_naive': 8.2},
        'mape_naive_best': 8.2,
        'naive_best_name': 'seasonal_naive',
        'pi_method': 'conformal_90',
        'pi_level': 0.9,
        'generated_at': '2026-07-03T10:00:00+00:00',
        'model_trained_at': '2026-07-03T09:00:00+00:00',
        'windows': [
            {
                'window': '2025-01-31 — 2025-03-31', 'train_periods': 24,
                'test_periods': 3, 'actual_total': 3100.0,
                'predicted_total': 2950.0, 'pi_low_total': 2700.0,
                'pi_high_total': 3250.0, 'hit_total': True, 'mape': 5.2,
                'per_period': [],
            },
            {
                'window': '2025-04-30 — 2025-06-30', 'train_periods': 27,
                'test_periods': 3, 'actual_total': 3400.0,
                'predicted_total': 3300.0, 'pi_low_total': 3050.0,
                'pi_high_total': 3550.0, 'hit_total': True, 'mape': 4.9,
                'per_period': [],
            },
        ],
    }
    bt.update(overrides)
    return bt


# ─── Адаптер: только завершённая проверка достигает билдера ──────────────────


def test_mapper_passes_only_ok_backtest():
    ok = _backtest_fixture()
    data = _map_pipeline_to_builder_data({}, {}, {}, None, backtest=ok)
    assert data['backtest']['windows_hit_total'] == 4

    for bad in (
        None,
        {'status': 'insufficient', 'message': 'мало истории'},
        {'status': 'error', 'error_code': 'NO_MODEL'},
        {'status': 'ok', 'windows': []},  # ok без окон — не витрина
    ):
        d = _map_pipeline_to_builder_data({}, {}, {}, None, backtest=bad)
        assert 'backtest' not in d, f'{bad} не должен рождать слайд'


# ─── Дека: 13 слайдов с витриной / 12 без, содержимое слайда ─────────────────


@pytest.fixture(scope='module')
def live_decks(tmp_path_factory):
    """Один synthetic live-проект → две деки: с витриной и без."""
    tmp = tmp_path_factory.mktemp('e1_pptx')
    pdir = _build_project(tmp, 'e1deck', beta_sd=0.2, seed=7)

    from engines.decomposer import decompose
    from engines.optimizer import optimize
    dec = decompose(str(pdir))
    assert dec.get('status') != 'error', dec.get('message')
    opt = optimize({'min_pct': 0.0, 'max_pct': 100.0}, str(pdir))
    assert opt.get('status', 'ok') != 'error', opt.get('message')

    model_data = {'diagnostics': {
        'metrics': {'r_squared': 0.81, 'mape_pct': 12.3, 'r_hat_max': 1.004,
                    'ess_bulk_min': 812.0, 'ess_tail_min': 640.0},
        'mqs': {'score': 71.0, 'tier_label': 'Хорошее'},
        'checks': {},
    }}
    from engines.pptx_export import build_pptx

    out_with = str(tmp / 'with_bt.pptx')
    res_with = build_pptx(model_data, dec, opt, out_with, scenarios=[],
                          project_id='e1_test', backtest=_backtest_fixture())
    assert res_with.get('status') == 'ok', res_with.get('message')

    out_without = str(tmp / 'without_bt.pptx')
    res_without = build_pptx(model_data, dec, opt, out_without, scenarios=[],
                             project_id='e1_test', backtest=None)
    assert res_without.get('status') == 'ok', res_without.get('message')

    return {
        'with': {'result': res_with, 'text': _extract_all_text(out_with),
                 'result_path': out_with},
        'without': {'result': res_without, 'text': _extract_all_text(out_without)},
        'pipeline': {'model_data': model_data, 'dec': dec, 'opt': opt, 'tmp': str(tmp)},
    }


def test_deck_with_backtest_has_13_slides(live_decks):
    assert live_decks['with']['result']['slides'] == 13


def test_backtest_slide_position_in_main_section(live_decks):
    """П5 (одобрено Антоном): витрина — слайд №6, финал секции «Главное»
    (сразу после SCQAR), а не в глубине методологии."""
    from pptx import Presentation
    prs = Presentation(live_decks['with']['result_path'])
    slide6_text = '\n'.join(
        sh.text_frame.text for sh in prs.slides[5].shapes if sh.has_text_frame
    )
    assert 'Проверка на истории' in slide6_text, (
        'Слайд №6 обязан быть витриной «Проверка на истории» (П5)'
    )


def test_deck_without_backtest_stays_12_and_clean(live_decks):
    assert live_decks['without']['result']['slides'] == 12
    text = live_decks['without']['text']
    # Никаких следов витрины без живой проверки — у слайда нет wireframe-режима.
    assert 'Проверка на истории' not in text
    assert 'МОДЕЛЬ ПРОТИВ ФАКТА' not in text


def test_backtest_slide_content_honest(live_decks):
    text = live_decks['with']['text']
    # Заголовок-вердикт и герой-факт
    assert 'Проверка на истории: 4 из 4 кварталов' in text
    # Числа модели против наивного
    assert 'Ошибка прогноза (MAPE): 6.1%' in text
    assert 'Наивный прогноз: 8.2%' in text
    assert 'модель точнее на 26%' in text
    # Покрытие по периодам с нормой
    assert 'Покрытие по периодам: 92% (12 точек, норма ≈ 90%)' in text
    # Окна с датами из данных и метод по-русски
    assert '2025-01-31 — 2025-03-31' in text
    assert 'скользящая проверка' in text
    assert 'будущее ей не показывают' in text


def _gen_compare_fixture(**overrides):
    gc = {
        'status': 'ok',
        'baseline': {'timestamp': '20260401_120000',
                     'trained_at': '2026-04-01T12:00:00+00:00'},
        'current': {'trained_at': '2026-07-03T10:00:00+00:00'},
        'channels': [
            {'name': 'TV', 'roi_old': 3.2, 'roi_new': 3.4,
             'roi_ci_old': [2.6, 3.9], 'roi_ci_new': [2.8, 4.0],
             'delta_pct': 6.3, 'verdict': 'stable', 'verdict_ru': 'стабильно',
             'method': 'ci_overlap', 'decay_shift': False, 'contribution_new': 900},
            {'name': 'Digital', 'roi_old': 2.0, 'roi_new': 4.8,
             'roi_ci_old': [1.6, 2.4], 'roi_ci_new': [4.2, 5.4],
             'delta_pct': 140.0, 'verdict': 'shift_strong',
             'verdict_ru': 'резкий сдвиг', 'method': 'ci_overlap',
             'decay_shift': False, 'contribution_new': 500},
        ],
        'added_channels': [], 'removed_channels': [],
        'summary': {
            'counts': {'stable': 1, 'shift_within_ci': 0, 'shift_strong': 1},
            'headline': 'ROI TV: был 3.2 [2.6–3.9], стал 3.4 [2.8–4.0] — стабильно.',
            'strong_shifts': ['Digital'],
            'probable_causes': ['Новые наблюдения изменили оценку вклада.'],
        },
        'generated_at': '2026-07-03T10:05:00+00:00',
    }
    gc.update(overrides)
    return gc


def test_mapper_passes_only_ok_generation_compare():
    ok = _gen_compare_fixture()
    data = _map_pipeline_to_builder_data({}, {}, {}, None, generation_compare=ok)
    assert data['generation_compare']['summary']['counts']['stable'] == 1
    for bad in (None, {'status': 'insufficient'}, {'status': 'ok', 'channels': []}):
        d = _map_pipeline_to_builder_data({}, {}, {}, None, generation_compare=bad)
        assert 'generation_compare' not in d


def test_deck_with_both_inserts_14_slides_ordered(live_decks, tmp_path):
    """E3: оба вставных артефакта → дека 14; №6 витрина, №7 «что изменилось»."""
    from pptx import Presentation
    from engines.pptx_export import build_pptx
    p = live_decks['pipeline']
    out = str(tmp_path / 'both.pptx')
    res = build_pptx(
        p['model_data'], p['dec'], p['opt'], out, scenarios=[],
        project_id='e3_test', backtest=_backtest_fixture(),
        generation_compare=_gen_compare_fixture(),
    )
    assert res.get('status') == 'ok', res.get('message')
    assert res['slides'] == 14
    prs = Presentation(out)

    def slide_text(i):
        return '\n'.join(
            sh.text_frame.text for sh in prs.slides[i].shapes if sh.has_text_frame
        )
    assert 'Проверка на истории' in slide_text(5)          # слайд №6
    s7 = slide_text(6)                                      # слайд №7
    assert 'был 3.2 [2.6–3.9], стал 3.4' in s7
    assert 'резкий сдвиг' in s7
    assert 'перекрытию интервалов' in s7
    assert 'Резких сдвигов: 1' in s7


def test_deck_with_only_gen_compare_13_slides(live_decks, tmp_path):
    from pptx import Presentation
    from engines.pptx_export import build_pptx
    p = live_decks['pipeline']
    out = str(tmp_path / 'only_gc.pptx')
    res = build_pptx(
        p['model_data'], p['dec'], p['opt'], out, scenarios=[],
        project_id='e3_test', backtest=None,
        generation_compare=_gen_compare_fixture(),
    )
    assert res.get('status') == 'ok'
    assert res['slides'] == 13
    prs = Presentation(out)
    s6 = '\n'.join(
        sh.text_frame.text for sh in prs.slides[5].shapes if sh.has_text_frame
    )
    assert 'был 3.2' in s6, 'без витрины «что изменилось» занимает слайд №6'


def test_calibration_marks_delivered(live_decks, tmp_path):
    """E2-3: [CALIBRATED] у канала + строка калибровки + честное расхождение
    (within_ci=false) доезжают до клиентской деки."""
    from engines.pptx_export import build_pptx
    p = live_decks['pipeline']
    model_data = json.loads(json.dumps(p['model_data']))  # deep copy
    ch_name = (p['dec'].get('channels') or [{}])[0].get('name')
    model_data['diagnostics']['calibration_applied'] = [{
        'channel': ch_name, 'test_type': 'geo_lift',
        'date_from': '2026-01-01', 'date_to': '2026-03-01', 'lift_abs': 500.0,
    }]
    model_data['diagnostics']['calibration_check'] = [{
        'channel': ch_name, 'test_type': 'geo_lift',
        'date_from': '2026-01-01', 'date_to': '2026-03-01',
        'test_lift': 500.0, 'test_sigma': 60.0,
        'model_contrib_mean': 320.0, 'model_contrib_ci90': [280.0, 360.0],
        'within_ci': False,
    }]
    out = str(tmp_path / 'calib.pptx')
    res = build_pptx(model_data, p['dec'], p['opt'], out, scenarios=[],
                     project_id='e2_test')
    assert res.get('status') == 'ok', res.get('message')
    text = _extract_all_text(out)
    assert '[CALIBRATED]' in text
    assert 'откалиброван тестом' in text
    assert 'Модель и тест расходятся' in text
    assert 'разберите период с аналитиком' in text


def test_promises_lines_delivered(live_decks, tmp_path):
    """E4-3: сверенные обещания (kept/missed) доезжают до отчёта; pending —
    нет (нечего показывать), wireframe-суррогатов нет."""
    from engines.pptx_export import build_pptx
    p = live_decks['pipeline']
    promises = [
        {'status': 'kept', 'status_ru': 'сбылось',
         'action_text': 'Бюджет 12 000 000 ₽ на Q3 по плану оптимизации',
         'verdict_note': 'Факт попал в интервал.'},
        {'status': 'missed', 'status_ru': 'не сбылось',
         'action_text': 'Сдвиг 10% в TV',
         'verdict_note': 'Факт вне интервала.'},
        {'status': 'pending', 'status_ru': 'ожидает данных',
         'action_text': 'Не должно попасть в отчёт'},
    ]
    out = str(tmp_path / 'promises.pptx')
    res = build_pptx(p['model_data'], p['dec'], p['opt'], out, scenarios=[],
                     project_id='e4_test', promises=promises)
    assert res.get('status') == 'ok', res.get('message')
    text = _extract_all_text(out)
    assert 'сбылось 1, не сбылось 1' in text
    assert 'Бюджет 12 000 000' in text
    assert 'Не должно попасть в отчёт' not in text

    # Без сверенных обещаний — ни следа блока
    out2 = str(tmp_path / 'no_promises.pptx')
    res2 = build_pptx(p['model_data'], p['dec'], p['opt'], out2, scenarios=[],
                      project_id='e4_test',
                      promises=[{'status': 'pending', 'action_text': 'x'}])
    assert res2.get('status') == 'ok'
    assert 'Проверка прошлых рекомендаций' not in _extract_all_text(out2)


def test_backtest_slide_worse_than_naive_title(live_decks, tmp_path):
    """Нелестный вердикт выносится в заголовок слайда — честность на витрине."""
    from engines.pptx_export import build_pptx
    p = live_decks['pipeline']
    out = str(tmp_path / 'wtn.pptx')
    res = build_pptx(
        p['model_data'], p['dec'], p['opt'], out, scenarios=[],
        project_id='e1_test',
        backtest=_backtest_fixture(verdict='worse_than_naive'),
    )
    assert res.get('status') == 'ok', res.get('message')
    text = _extract_all_text(out)
    assert 'модель пока не точнее наивного прогноза' in text


if __name__ == '__main__':
    sys.exit(pytest.main([__file__, '-q']))
